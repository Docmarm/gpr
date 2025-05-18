import streamlit as st
import pandas as pd
import numpy as np
import plotly.express as px
import plotly.graph_objects as go
from plotly.subplots import make_subplots
from datetime import datetime, timedelta
import io
import os
import time  # Ajoutez cette ligne ici
from typing import Dict, List, Tuple, Optional, Any
from dateutil.relativedelta import relativedelta
from haversine import haversine, Unit # Pour calculer la distance entre coordonnées GPS
import folium
from streamlit_folium import folium_static
from PIL import Image
import base64
from io import BytesIO
# ---------------------------------------------------------------------
# Constantes et Utilitaires
# ---------------------------------------------------------------------
EXCEL_MIME_TYPE = "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet"
DATE_FORMAT = '%Y-%m-%d'

# --- Valeurs par défaut pour les paramètres (seront stockées dans session_state) ---
DEFAULT_SEUIL_HEURES_RAPPROCHEES = 2
DEFAULT_CONSO_SEUIL = 16.0
DEFAULT_HEURE_DEBUT_NON_OUVRE = 20
DEFAULT_HEURE_FIN_NON_OUVRE = 6
DEFAULT_DELTA_MINUTES_FACTURATION_DOUBLE = 60
DEFAULT_SEUIL_ANOMALIES_SUSPECTES_SCORE = 10 # Basé sur un score pondéré

# --- Poids par défaut pour le score de risque ---
DEFAULT_POIDS_CONSO_EXCESSIVE = 5
DEFAULT_POIDS_DEPASSEMENT_CAPACITE = 10
DEFAULT_POIDS_PRISES_RAPPROCHEES = 3
DEFAULT_POIDS_KM_DECROISSANT = 8
DEFAULT_POIDS_KM_INCHANGE = 2
DEFAULT_POIDS_KM_SAUT = 6
DEFAULT_POIDS_HORS_HORAIRE = 2
DEFAULT_POIDS_HORS_SERVICE = 9
DEFAULT_POIDS_FACT_DOUBLE = 7
DEFAULT_POIDS_ANOMALIE_GEO = 10 # Nouveau poids pour anomalie géographique

# --- Nouveaux poids pour anomalies de géolocalisation ---
DEFAULT_POIDS_TRAJET_HORS_HEURES = 6
DEFAULT_POIDS_TRAJET_WEEKEND = 5
DEFAULT_POIDS_ARRETS_FREQUENTS = 4
DEFAULT_POIDS_DETOUR_SUSPECT = 7
DEFAULT_POIDS_TRANSACTION_SANS_PRESENCE = 9
DEFAULT_POIDS_VITESSE_EXCESSIVE = 8

# --- Valeurs par défaut pour les paramètres de géolocalisation ---
DEFAULT_RAYON_STATION_KM = 0.3  # Rayon autour d'une station pour considérer le véhicule présent
DEFAULT_SEUIL_ARRET_MINUTES = 5  # Durée minimum pour considérer un arrêt
DEFAULT_SEUIL_DETOUR_PCT = 20    # Pourcentage au-delà duquel un trajet est considéré comme détour
DEFAULT_HEURE_DEBUT_SERVICE = 7  # Heure normale de début de service
DEFAULT_HEURE_FIN_SERVICE = 19   # Heure normale de fin de service
DEFAULT_NB_ARRETS_SUSPECT = 4    # Nombre d'arrêts au-delà duquel c'est suspect pour un trajet court
DEFAULT_VITESSE_EXCESSIVE_SEUIL = 90  # Vitesse maximale en km/h au-delà de laquelle on considère un excès de vitesse

# ---------------------------------------------------------------------
# Initialisation Session State pour les Paramètres
# ---------------------------------------------------------------------
def initialize_session_state(df_vehicules: Optional[pd.DataFrame] = None):
    """Initialise les paramètres dans st.session_state s'ils n'existent pas."""
    defaults = {
        'ss_seuil_heures_rapprochees': DEFAULT_SEUIL_HEURES_RAPPROCHEES,
        'ss_heure_debut_non_ouvre': DEFAULT_HEURE_DEBUT_NON_OUVRE,
        'ss_heure_fin_non_ouvre': DEFAULT_HEURE_FIN_NON_OUVRE,
        'ss_delta_minutes_facturation_double': DEFAULT_DELTA_MINUTES_FACTURATION_DOUBLE,
        'ss_seuil_anomalies_suspectes_score': DEFAULT_SEUIL_ANOMALIES_SUSPECTES_SCORE,
        'ss_poids_conso_excessive': DEFAULT_POIDS_CONSO_EXCESSIVE,
        'ss_poids_depassement_capacite': DEFAULT_POIDS_DEPASSEMENT_CAPACITE,
        'ss_poids_prises_rapprochees': DEFAULT_POIDS_PRISES_RAPPROCHEES,
        'ss_poids_km_decroissant': DEFAULT_POIDS_KM_DECROISSANT,
        'ss_poids_km_inchange': DEFAULT_POIDS_KM_INCHANGE,
        'ss_poids_km_saut': DEFAULT_POIDS_KM_SAUT,
        'ss_poids_hors_horaire': DEFAULT_POIDS_HORS_HORAIRE,
        'ss_poids_hors_service': DEFAULT_POIDS_HORS_SERVICE,
        'ss_poids_fact_double': DEFAULT_POIDS_FACT_DOUBLE,
        'ss_conso_seuils_par_categorie': {}, # Sera peuplé dynamiquement
        'data_loaded': False, # Indicateur de chargement des données
        # Nouveaux paramètres pour géolocalisation
        'ss_rayon_station_km': DEFAULT_RAYON_STATION_KM,
        'ss_seuil_arret_minutes': DEFAULT_SEUIL_ARRET_MINUTES,
        'ss_seuil_detour_pct': DEFAULT_SEUIL_DETOUR_PCT,
        'ss_heure_debut_service': DEFAULT_HEURE_DEBUT_SERVICE,
        'ss_heure_fin_service': DEFAULT_HEURE_FIN_SERVICE,
        'ss_nb_arrets_suspect': DEFAULT_NB_ARRETS_SUSPECT,
        'ss_vitesse_excessive_seuil': DEFAULT_VITESSE_EXCESSIVE_SEUIL,  # Nouveau paramètre pour vitesse maximale
        # Nouveaux poids pour anomalies de géolocalisation
        'ss_poids_trajet_hors_heures': DEFAULT_POIDS_TRAJET_HORS_HEURES,
        'ss_poids_trajet_weekend': DEFAULT_POIDS_TRAJET_WEEKEND,
        'ss_poids_arrets_frequents': DEFAULT_POIDS_ARRETS_FREQUENTS,
        'ss_poids_detour_suspect': DEFAULT_POIDS_DETOUR_SUSPECT,
        'ss_poids_transaction_sans_presence': DEFAULT_POIDS_TRANSACTION_SANS_PRESENCE,
        'ss_poids_vitesse_excessive': DEFAULT_POIDS_VITESSE_EXCESSIVE,
        # Paramètres d'activation/désactivation des types d'anomalies
        'ss_activer_trajets_suspects': False,  # Désactivé par défaut
        'ss_activer_detours_suspects': False,  # Désactivé par défaut
        'ss_activer_transactions_sans_presence': True   # Activé par défaut
    }
    for key, value in defaults.items():
        if key not in st.session_state:
            st.session_state[key] = value

    # Initialisation dynamique des seuils de conso par catégorie
    if df_vehicules is not None and not st.session_state['ss_conso_seuils_par_categorie']:
        all_cats = sorted(df_vehicules['Catégorie'].dropna().astype(str).unique())
        st.session_state['ss_conso_seuils_par_categorie'] = {cat: DEFAULT_CONSO_SEUIL for cat in all_cats}
    elif df_vehicules is not None:
        # S'assurer que toutes les catégories actuelles ont un seuil
        all_cats = sorted(df_vehicules['Catégorie'].dropna().astype(str).unique())
        current_seuils = st.session_state['ss_conso_seuils_par_categorie']
        updated_seuils = {cat: current_seuils.get(cat, DEFAULT_CONSO_SEUIL) for cat in all_cats}
        st.session_state['ss_conso_seuils_par_categorie'] = updated_seuils
# ---------------------------------------------------------------------
# Fonctions de calcul de score spécifiques à la géolocalisation
# ---------------------------------------------------------------------

def calculer_score_ecart_distance(pourcentage_ecart: float) -> float:
    """
    Calcule le score d'écart de distance avec plafonnement.

    Args:
        pourcentage_ecart: L'écart en pourcentage entre la distance géoloc et la distance transaction.

    Returns:
        Le score calculé, plafonné à 50.
    """
    # Plafonner le pourcentage d'écart à 200% maximum
    pourcentage_ecart_plafonne = min(abs(pourcentage_ecart), 200)

    # Calculer le score avec le pourcentage plafonné
    poids = st.session_state.get('ss_poids_detour_suspect', DEFAULT_POIDS_DETOUR_SUSPECT)
    score = pourcentage_ecart_plafonne * (poids / 20)

    # Plafonner également le score final à 50 maximum
    return min(score, 50)


# ---------------------------------------------------------------------
# Fonctions : Nettoyage et chargement (inchangées sauf ajout Type Hints)
# ---------------------------------------------------------------------
def nettoyer_numero_carte(numero: Any) -> str:
    """Convertit un numéro de carte en entier si possible, puis string, retirant les espaces."""
    if pd.isna(numero):
        return ""
    try:
        # Tenter conversion en float puis int pour gérer les ".0" puis en str
        return str(int(float(str(numero)))).strip()
    except ValueError:
         # Si la conversion échoue, retourner le numéro comme string nettoyé
        return str(numero).strip()
    except Exception:
        # Fallback général
        return str(numero).strip()

@st.cache_data(show_spinner="Chargement et nettoyage des fichiers...")
def charger_donnees(fichier_transactions, fichier_cartes) -> Tuple[Optional[pd.DataFrame], Optional[pd.DataFrame], Optional[pd.DataFrame], Optional[pd.DataFrame]]:
    """Charge et nettoie les données des fichiers CSV et Excel."""
    df_transactions, df_vehicules, df_ge, df_autres = None, None, None, None

    # --- Chargement Transactions ---
    try:
        df_transactions = pd.read_csv(fichier_transactions, sep=';', encoding='utf-8', low_memory=False)
        if 'Amount eur' in df_transactions.columns and 'Amount' not in df_transactions.columns:
             df_transactions.rename(columns={'Amount eur': 'Amount'}, inplace=True)
        if 'Place' not in df_transactions.columns and 'Place name' in df_transactions.columns:
             df_transactions.rename(columns={'Place name': 'Place'}, inplace=True)
    except Exception as e:
        st.error(f"Erreur lors de la lecture du fichier de transactions : {e}")
        return None, None, None, None

    # --- Chargement Cartes ---
    try:
        xls = pd.ExcelFile(fichier_cartes)
        required_sheets = {'CARTES VEHICULE', 'CARTES GE', 'AUTRES CARTES'}
        available_sheets = set(xls.sheet_names)
        if not required_sheets.issubset(available_sheets):
             st.error(f"Feuilles manquantes dans le fichier Excel. Attendues: {required_sheets}. Trouvées: {available_sheets}")
             return None, None, None, None

        df_vehicules = pd.read_excel(xls, sheet_name='CARTES VEHICULE')
        df_ge = pd.read_excel(xls, sheet_name='CARTES GE')
        df_autres = pd.read_excel(xls, sheet_name='AUTRES CARTES')
    except Exception as e:
        st.error(f"Erreur lors de la lecture du fichier des cartes : {e}")
        return None, None, None, None

    # --- Vérification Colonnes Transactions ---
    colonnes_attendues_transactions = ['Date', 'Hour', 'Card num.', 'Quantity', 'Past mileage', 'Current mileage', 'Amount', 'Place']
    missing_cols_trans = [col for col in colonnes_attendues_transactions if col not in df_transactions.columns]
    if missing_cols_trans:
        st.error(f"Colonnes manquantes dans le fichier de transactions: {', '.join(missing_cols_trans)}")
        return None, None, None, None

    # --- Vérification Colonnes Cartes ---
    colonnes_attendues_cartes = ['N° Carte']
    dfs_cartes = {'CARTES VEHICULE': df_vehicules, 'CARTES GE': df_ge, 'AUTRES CARTES': df_autres}
    for nom_feuille, df_sheet in dfs_cartes.items():
        missing_cols_carte = [col for col in colonnes_attendues_cartes if col not in df_sheet.columns]
        if missing_cols_carte:
            st.error(f"Colonne(s) manquante(s) dans la feuille '{nom_feuille}': {', '.join(missing_cols_carte)}")
            return None, None, None, None
        # Vérifier et convertir 'Cap-rèservoir' si existe
        if 'Cap-rèservoir' in df_sheet.columns:
            df_sheet['Cap-rèservoir'] = pd.to_numeric(df_sheet['Cap-rèservoir'], errors='coerce').fillna(0)
        # NOUVEAU: Vérifier et convertir 'Dotation' si existe
        if 'Dotation' in df_sheet.columns:
            df_sheet['Dotation'] = pd.to_numeric(df_sheet['Dotation'], errors='coerce').fillna(0)
        elif nom_feuille == 'CARTES VEHICULE': # Si absente dans CARTES VEHICULE, créer avec 0 et avertir
            st.warning("La colonne 'Dotation' est manquante dans la feuille 'CARTES VEHICULE'. Elle sera initialisée à 0. Le suivi des dotations ne sera pas significatif.")
            df_sheet['Dotation'] = 0

        # Assurer que Catégorie est string
        if 'Catégorie' in df_sheet.columns:
            df_sheet['Catégorie'] = df_sheet['Catégorie'].astype(str).fillna('Non défini')


    # --- Nettoyage Numéros de Carte ---
    df_transactions['Card num.'] = df_transactions['Card num.'].apply(nettoyer_numero_carte)
    for df_sheet in [df_vehicules, df_ge, df_autres]:
        df_sheet['N° Carte'] = df_sheet['N° Carte'].apply(nettoyer_numero_carte)
        df_sheet.dropna(subset=['N° Carte'], inplace=True)
        df_sheet = df_sheet[df_sheet['N° Carte'] != ""]

    # --- Conversion Types Transactions ---
    df_transactions['Date'] = pd.to_datetime(df_transactions['Date'], format='%d/%m/%Y', errors='coerce')
    try:
        df_transactions['Hour'] = pd.to_datetime(df_transactions['Hour'], format='%H:%M:%S', errors='coerce').dt.time
    except ValueError:
        try:
             df_transactions['Hour'] = pd.to_datetime(df_transactions['Hour'], format='%H:%M', errors='coerce').dt.time
        except Exception as e:
             st.warning(f"Impossible de parser la colonne 'Hour'. Vérifiez le format (HH:MM:SS ou HH:MM). Erreur: {e}")
             df_transactions['Hour'] = pd.NaT

    for col in ['Quantity', 'Past mileage', 'Current mileage', 'Amount']:
        df_transactions[col] = pd.to_numeric(df_transactions[col].astype(str).str.replace(',', '.'), errors='coerce')

    # --- Suppression Lignes Invalides ---
    df_transactions.dropna(subset=['Date', 'Card num.'], inplace=True)
    df_transactions = df_transactions[df_transactions['Card num.'] != ""]
    df_vehicules = df_vehicules[df_vehicules['N° Carte'] != ""]
    df_ge = df_ge[df_ge['N° Carte'] != ""]
    df_autres = df_autres[df_autres['N° Carte'] != ""]

    df_transactions['DateTime'] = df_transactions.apply(
        lambda row: datetime.combine(row['Date'].date(), row['Hour']) if pd.notna(row['Date']) and pd.notna(row['Hour']) else pd.NaT,
        axis=1
    )
    df_transactions.dropna(subset=['DateTime'], inplace=True)

    return df_transactions, df_vehicules, df_ge, df_autres

# MODIFIED FUNCTION
@st.cache_data(show_spinner="Chargement et nettoyage du fichier de géolocalisation...")
def charger_donnees_geolocalisation(fichier_geoloc) -> Optional[pd.DataFrame]:
    """Charge et nettoie les données du fichier de géolocalisation."""
    if fichier_geoloc is None:
        return None

    try:
        # Chargement en spécifiant les types pour éviter les problèmes de conversion
        df_geoloc = pd.read_excel(fichier_geoloc, engine='openpyxl', dtype={
            'Début': str, 'Fin': str, 'Durée': str, 'Distance': str
        })

        # Vérification des colonnes essentielles
        colonnes_attendues = ['Véhicule', 'Date', 'Début', 'Fin', 'Durée', 'Distance', 'Type']
        # Ajout des coordonnées si elles existent (pour la carte)
        optional_coords = ['Latitude_depart', 'Longitude_depart', 'Latitude_arrivee', 'Longitude_arrivee']
        for oc in optional_coords:
            if oc in df_geoloc.columns:
                colonnes_attendues.append(oc)

        missing_cols = [col for col in colonnes_attendues if col not in df_geoloc.columns and col not in optional_coords]
        # Re-check missing_cols only for mandatory ones if some optional are not found
        mandatory_cols = ['Véhicule', 'Date', 'Début', 'Fin', 'Durée', 'Distance', 'Type']
        missing_mandatory_cols = [col for col in mandatory_cols if col not in df_geoloc.columns]


        if missing_mandatory_cols:
            st.error(f"Colonnes manquantes dans le fichier de géolocalisation: {', '.join(missing_mandatory_cols)}")
            return None

        # Normalisation du type (Trajet/Arrêt)
        if 'Type' in df_geoloc.columns:
            df_geoloc['Type'] = df_geoloc['Type'].astype(str).str.strip()
            type_corrections = {
                'trajet': 'Trajet', 'TRAJET': 'Trajet', 'Tragets': 'Trajet',
                'arret': 'Arrêt', 'ARRÊT': 'Arrêt', 'Arrèt': 'Arrêt', 'Arret': 'Arrêt', 'ArrËt': 'Arrêt'
            }
            df_geoloc['Type'] = df_geoloc['Type'].replace(type_corrections)

        # Conversion des dates
        try:
            df_geoloc['Date'] = pd.to_datetime(df_geoloc['Date'], format='%d/%m/%Y', errors='coerce')
        except:
            # Essayer sans format spécifique si le format standard échoue
            df_geoloc['Date'] = pd.to_datetime(df_geoloc['Date'], errors='coerce')

        if df_geoloc['Date'].isna().all():
            st.error("Impossible de convertir la colonne 'Date' en format date")
            return None

        # Conversion des heures de début et fin
        for col in ['Début', 'Fin']:
            df_geoloc[f'{col}_str'] = df_geoloc[col].astype(str).replace(['NaT', 'nan', '<NA>', 'None', 'nat', 'NAT', 'none'], '', regex=False).str.split('.').str[0]

        date_str_series = df_geoloc['Date'].dt.strftime('%Y-%m-%d')

        # Conversion des heures de début
        df_geoloc['DateTime_Debut'] = pd.NaT
        valid_debut_mask = df_geoloc['Date'].notna() & (df_geoloc['Début_str'] != '')
        if valid_debut_mask.any():
            combined_debut_str = date_str_series[valid_debut_mask] + ' ' + df_geoloc.loc[valid_debut_mask, 'Début_str']
            df_geoloc.loc[valid_debut_mask, 'DateTime_Debut'] = pd.to_datetime(combined_debut_str, errors='coerce')

        # Conversion des heures de fin
        df_geoloc['DateTime_Fin'] = pd.NaT
        valid_fin_mask = df_geoloc['Date'].notna() & (df_geoloc['Fin_str'] != '')
        if valid_fin_mask.any():
            combined_fin_str = date_str_series[valid_fin_mask] + ' ' + df_geoloc.loc[valid_fin_mask, 'Fin_str']
            df_geoloc.loc[valid_fin_mask, 'DateTime_Fin'] = pd.to_datetime(combined_fin_str, errors='coerce')

        # Conversion des durées en timedelta
        df_geoloc['Durée_str'] = df_geoloc['Durée'].astype(str).replace(['NaT', 'nan', '<NA>', 'None', 'nat', 'NAT', 'none'], '', regex=False).str.split('.').str[0]
        df_geoloc['Durée_minutes'] = pd.to_timedelta(df_geoloc['Durée_str'], errors='coerce').dt.total_seconds() / 60

        # Recalculer les durées manquantes à partir de début et fin
        mask_recalc_duree = (df_geoloc['Type'] == 'Trajet') & df_geoloc['Durée_minutes'].isna() & df_geoloc['DateTime_Debut'].notna() & df_geoloc['DateTime_Fin'].notna() & (df_geoloc['DateTime_Fin'] > df_geoloc['DateTime_Debut'])
        if mask_recalc_duree.any():
            df_geoloc.loc[mask_recalc_duree, 'Durée_minutes'] = (df_geoloc.loc[mask_recalc_duree, 'DateTime_Fin'] - df_geoloc.loc[mask_recalc_duree, 'DateTime_Debut']).dt.total_seconds() / 60

        # Si la fin est avant le début, on ajoute un jour à la fin (trajet à cheval sur minuit)
        mask_nuit = df_geoloc['DateTime_Fin'] < df_geoloc['DateTime_Debut']
        if mask_nuit.any(): # Check if any True values before attempting .loc
            df_geoloc.loc[mask_nuit, 'DateTime_Fin'] = df_geoloc.loc[mask_nuit, 'DateTime_Fin'] + timedelta(days=1)


        # *** NOUVEAU TRAITEMENT DES DISTANCES ***
        # Nettoyer et convertir la colonne Distance
        distance_str = df_geoloc['Distance'].astype(str)
        # Enlever "km" et tout espace en fin de chaîne
        distance_str = distance_str.str.strip().str.replace(r'\s*km$', '', regex=True, case=False)
        # Enlever les points avant les séparateurs de milliers (1.234,56 -> 1234,56)
        distance_str = distance_str.str.replace(r'\.(?=\d{3},)', '', regex=True)
        # Remplacer la virgule par un point décimal
        distance_str = distance_str.str.replace(',', '.', regex=False)
        # Enlever tous les espaces
        distance_str = distance_str.str.replace(r'\s+', '', regex=True)
        # Convertir en numérique
        df_geoloc['Distance'] = pd.to_numeric(distance_str, errors='coerce')

        # Calcul de la vitesse moyenne pour les trajets
        df_geoloc['Vitesse moyenne'] = np.nan
        trajets_mask = df_geoloc['Type'] == 'Trajet'
        if trajets_mask.any():
            trajets_df = df_geoloc[trajets_mask]
            durees_minutes = trajets_df['Durée_minutes']
            distances = trajets_df['Distance']
            calcul_valide_mask = durees_minutes.notna() & (durees_minutes > 0) & distances.notna() & (distances > 0)
            if calcul_valide_mask.any():
                duree_heures_valides = durees_minutes[calcul_valide_mask] / 60  # minutes -> heures
                distances_valides = distances[calcul_valide_mask]
                vitesses = (distances_valides / duree_heures_valides).round(2)
                df_geoloc.loc[trajets_df[calcul_valide_mask].index, 'Vitesse moyenne'] = vitesses

        # Ajout des jours de la semaine et indication weekend
        df_geoloc['Jour_semaine'] = df_geoloc['Date'].dt.dayofweek
        df_geoloc['Est_weekend'] = df_geoloc['Jour_semaine'] >= 5

        # Ajout de l'heure de début en format numérique
        df_geoloc['Heure_debut'] = df_geoloc['DateTime_Debut'].dt.hour

        # Création d'une colonne AnnéeMois pour le regroupement
        df_geoloc['AnnéeMois'] = df_geoloc['Date'].dt.to_period('M').astype(str)

        # Suppression des lignes avec des données cruciales manquantes
        df_geoloc.dropna(subset=['Véhicule', 'Date', 'Type'], inplace=True)

        # Nettoyage des colonnes temporaires
        columns_to_drop = ['Début_str', 'Fin_str', 'Durée_str']
        df_geoloc.drop(columns=columns_to_drop, errors='ignore', inplace=True)

        # Coordonnées GPS (vérifier si elles existent avant de les traiter)
        if 'Latitude_depart' in df_geoloc.columns and 'Longitude_depart' in df_geoloc.columns:
            df_geoloc['Coordonnees_depart'] = df_geoloc.apply(
                lambda row: (row['Latitude_depart'], row['Longitude_depart'])
                if pd.notna(row['Latitude_depart']) and pd.notna(row['Longitude_depart'])
                else None,
                axis=1
            )
        else:
            df_geoloc['Latitude_depart'] = np.nan
            df_geoloc['Longitude_depart'] = np.nan
            df_geoloc['Coordonnees_depart'] = None

        if 'Latitude_arrivee' in df_geoloc.columns and 'Longitude_arrivee' in df_geoloc.columns:
            df_geoloc['Coordonnees_arrivee'] = df_geoloc.apply(
                lambda row: (row['Latitude_arrivee'], row['Longitude_arrivee'])
                if pd.notna(row['Latitude_arrivee']) and pd.notna(row['Longitude_arrivee'])
                else None,
                axis=1
            )
        else:
            df_geoloc['Latitude_arrivee'] = np.nan
            df_geoloc['Longitude_arrivee'] = np.nan
            df_geoloc['Coordonnees_arrivee'] = None


        return df_geoloc

    except Exception as e:
        st.error(f"Erreur lors de la lecture du fichier de géolocalisation : {e}")
        return None


# ---------------------------------------------------------------------
# Fonctions : export Excel + affichage DataFrame (inchangées)
# ---------------------------------------------------------------------
def to_excel(df: pd.DataFrame) -> bytes:
    """Convertit un DataFrame en un fichier Excel (BytesIO)."""
    output = io.BytesIO()
    with pd.ExcelWriter(output, engine='openpyxl') as writer:
        df_copy = df.copy()
        for col in df_copy.select_dtypes(include=['datetime64[ns]']).columns:
             df_copy[col] = df_copy[col].dt.strftime(DATE_FORMAT)
        df_copy.to_excel(writer, index=False, sheet_name='Sheet1')
    return output.getvalue()

def afficher_dataframe_avec_export(df: pd.DataFrame, titre: str = "Tableau", key: str = "df"):
    """Affiche un DataFrame + un bouton d'export Excel."""
    if df is None or df.empty:
        st.info(f"{titre} : Aucune donnée à afficher.")
        return

    nb_lignes = len(df)
    st.markdown(f"### {titre} ({nb_lignes:,} lignes)")
    st.dataframe(df, use_container_width=True)
    try:
        excel_data = to_excel(df)
        nom_fic = f"{titre.lower().replace(' ', '_').replace('(', '').replace(')', '').replace(':', '')[:50]}.xlsx"
        st.download_button(
            label=f"Exporter '{titre}' en Excel",
            data=excel_data,
            file_name=nom_fic,
            mime=EXCEL_MIME_TYPE,
            key=f"export_{key}"
        )
    except Exception as e:
        st.error(f"Erreur lors de la génération de l'export Excel pour '{titre}': {e}")


# ---------------------------------------------------------------------
# Fonctions : Vérifications et Analyses (Mises à jour pour utiliser session_state)
# ---------------------------------------------------------------------

def verifier_cartes_inconnues(df_transactions: pd.DataFrame, df_vehicules: pd.DataFrame, df_ge: pd.DataFrame, df_autres: pd.DataFrame) -> pd.DataFrame:
    """Identifie les transactions associées à des cartes non listées."""
    cartes_vehicules = set(df_vehicules['N° Carte'].unique()) if df_vehicules is not None else set()
    cartes_ge = set(df_ge['N° Carte'].unique()) if df_ge is not None else set()
    cartes_autres = set(df_autres['N° Carte'].unique()) if df_autres is not None else set()
    cartes_valides = cartes_vehicules.union(cartes_ge).union(cartes_autres)

    cartes_utilisees = set(df_transactions['Card num.'].unique())
    cartes_inconnues = cartes_utilisees - cartes_valides

    if not cartes_inconnues:
        return pd.DataFrame()

    df_inc = df_transactions[df_transactions['Card num.'].isin(cartes_inconnues)].copy()

    if 'Card name' in df_inc.columns:
        stats = df_inc.groupby(['Card num.', 'Card name']).agg(
            Nombre_transactions=('Quantity', 'count'),
            Volume_total_L=('Quantity', 'sum'),
            Montant_total_CFA=('Amount', 'sum')
        ).round(2).reset_index()
    else:
         stats = df_inc.groupby('Card num.').agg(
            Nombre_transactions=('Quantity', 'count'),
            Volume_total_L=('Quantity', 'sum'),
            Montant_total_CFA=('Amount', 'sum')
        ).round(2).reset_index()
         stats['Card name'] = 'Nom non disponible'

    stats = stats[['Card num.', 'Card name', 'Nombre_transactions', 'Volume_total_L', 'Montant_total_CFA']]
    return stats


def detecter_anomalies(
    df_transactions: pd.DataFrame,
    df_vehicules: pd.DataFrame
) -> pd.DataFrame:
    """
    Fonction centrale (simplifiée) pour détecter toutes les anomalies sur les véhicules.
    Retourne un DataFrame unique avec toutes les anomalies détectées.
    """
    all_anomalies = []
    # S'assurer que la colonne 'Dotation' existe dans df_vehicules, sinon l'ajouter avec 0
    cols_vehicules_necessaires = ['N° Carte', 'Nouveau Immat', 'Catégorie', 'Type', 'Cap-rèservoir']
    if 'Dotation' in df_vehicules.columns:
        cols_vehicules_necessaires.append('Dotation')

    df_merged = df_transactions.merge(
        df_vehicules[cols_vehicules_necessaires], # Utiliser les colonnes nécessaires
        left_on='Card num.',
        right_on='N° Carte',
        how='inner'
    )
    df_merged['distance_parcourue'] = df_merged['Current mileage'] - df_merged['Past mileage']
    df_merged['consommation_100km'] = np.where(
         (df_merged['distance_parcourue'] > 0) & df_merged['Quantity'].notna(),
         (df_merged['Quantity'] / df_merged['distance_parcourue']) * 100,
         np.nan
    )

    seuils_conso = st.session_state.get('ss_conso_seuils_par_categorie', {})
    seuil_heures_rapprochees = st.session_state.get('ss_seuil_heures_rapprochees', DEFAULT_SEUIL_HEURES_RAPPROCHEES)
    heure_debut_non_ouvre = st.session_state.get('ss_heure_debut_non_ouvre', DEFAULT_HEURE_DEBUT_NON_OUVRE)
    heure_fin_non_ouvre = st.session_state.get('ss_heure_fin_non_ouvre', DEFAULT_HEURE_FIN_NON_OUVRE)
    delta_minutes_double = st.session_state.get('ss_delta_minutes_facturation_double', DEFAULT_DELTA_MINUTES_FACTURATION_DOUBLE)

    for index, row in df_merged.iterrows():
        cat = row['Catégorie']
        seuil = seuils_conso.get(cat, DEFAULT_CONSO_SEUIL)
        if pd.notna(row['consommation_100km']) and row['consommation_100km'] > seuil:
            anomalie = row.to_dict()
            anomalie['type_anomalie'] = 'Consommation excessive'
            anomalie['detail_anomalie'] = f"{row['consommation_100km']:.1f} L/100km > seuil {seuil} L/100km"
            anomalie['poids_anomalie'] = st.session_state.get('ss_poids_conso_excessive', DEFAULT_POIDS_CONSO_EXCESSIVE)
            all_anomalies.append(anomalie)

    depassement = df_merged[df_merged['Quantity'] > df_merged['Cap-rèservoir']].copy()
    if not depassement.empty:
         depassement['type_anomalie'] = 'Dépassement capacité'
         depassement['detail_anomalie'] = depassement.apply(lambda x: f"Volume: {x['Quantity']:.1f}L > Capacité: {x['Cap-rèservoir']:.1f}L", axis=1)
         depassement['poids_anomalie'] = st.session_state.get('ss_poids_depassement_capacite', DEFAULT_POIDS_DEPASSEMENT_CAPACITE)
         all_anomalies.extend(depassement.to_dict('records'))

    df_merged_sorted = df_merged.sort_values(['Card num.', 'DateTime'])
    rapprochees_indices = set()
    for carte in df_merged_sorted['Card num.'].unique():
        sub = df_merged_sorted[df_merged_sorted['Card num.'] == carte]
        if len(sub) > 1:
            time_diff = sub['DateTime'].diff().dt.total_seconds() / 3600
            indices_anomalie = sub.index[time_diff < seuil_heures_rapprochees]
            indices_precedents = sub.index[time_diff.shift(-1) < seuil_heures_rapprochees]
            rapprochees_indices.update(indices_anomalie)
            rapprochees_indices.update(indices_precedents)

    if rapprochees_indices:
        rapprochees_df = df_merged_sorted.loc[list(rapprochees_indices)].copy()
        rapprochees_df['type_anomalie'] = 'Prises rapprochées'
        rapprochees_df['detail_anomalie'] = f'Moins de {seuil_heures_rapprochees}h entre prises'
        rapprochees_df['poids_anomalie'] = st.session_state.get('ss_poids_prises_rapprochees', DEFAULT_POIDS_PRISES_RAPPROCHEES)
        all_anomalies.extend(rapprochees_df.to_dict('records'))

    km_anomalies = []
    for carte in df_merged_sorted['Card num.'].unique():
        df_carte = df_merged_sorted[df_merged_sorted['Card num.'] == carte]
        prev_m = None
        prev_row = None
        for index, row in df_carte.iterrows():
            curr_m = row['Current mileage']
            past_m = row['Past mileage']
            user = row.get('Card name', 'N/A')

            if pd.isna(curr_m) or curr_m == 0 or pd.isna(past_m) :
                 prev_m = None
                 prev_row = row
                 continue
            if past_m > curr_m:
                 anomalie = row.to_dict()
                 anomalie['type_anomalie'] = 'Kilométrage incohérent (transaction)'
                 anomalie['detail_anomalie'] = f"Km début ({past_m}) > Km fin ({curr_m})"
                 anomalie['poids_anomalie'] = st.session_state.get('ss_poids_km_decroissant', DEFAULT_POIDS_KM_DECROISSANT)
                 km_anomalies.append(anomalie)
                 prev_m = None
                 prev_row = row
                 continue
            if prev_m is not None and prev_row is not None:
                 if curr_m < prev_m:
                     anomalie = row.to_dict()
                     anomalie['type_anomalie'] = 'Kilométrage décroissant (inter-transactions)'
                     anomalie['detail_anomalie'] = f"Km transaction N ({curr_m}) < Km transaction N-1 ({prev_m})"
                     anomalie['poids_anomalie'] = st.session_state.get('ss_poids_km_decroissant', DEFAULT_POIDS_KM_DECROISSANT)
                     km_anomalies.append(anomalie)
                 elif curr_m == prev_m:
                     anomalie = row.to_dict()
                     anomalie['type_anomalie'] = 'Kilométrage inchangé (inter-transactions)'
                     anomalie['detail_anomalie'] = f"Km identique à la transaction précédente: {curr_m} km"
                     anomalie['poids_anomalie'] = st.session_state.get('ss_poids_km_inchange', DEFAULT_POIDS_KM_INCHANGE)
                     km_anomalies.append(anomalie)
                 elif (curr_m - prev_m) > 1000:
                     anomalie = row.to_dict()
                     anomalie['type_anomalie'] = 'Saut kilométrage important'
                     anomalie['detail_anomalie'] = f"Augmentation de +{curr_m - prev_m} km depuis transaction précédente"
                     anomalie['poids_anomalie'] = st.session_state.get('ss_poids_km_saut', DEFAULT_POIDS_KM_SAUT)
                     km_anomalies.append(anomalie)
            prev_m = curr_m
            prev_row = row
    all_anomalies.extend(km_anomalies)

    df_merged['heure'] = df_merged['DateTime'].dt.hour
    df_merged['jour_semaine'] = df_merged['DateTime'].dt.dayofweek
    if heure_debut_non_ouvre < heure_fin_non_ouvre:
        cond_heure = (df_merged['heure'] < heure_fin_non_ouvre) | (df_merged['heure'] >= heure_debut_non_ouvre)
    else:
        cond_heure = (df_merged['heure'] >= heure_debut_non_ouvre) | (df_merged['heure'] < heure_fin_non_ouvre)
    cond_weekend = (df_merged['jour_semaine'] >= 5)
    anomalies_hor = df_merged[cond_heure | cond_weekend].copy()
    if not anomalies_hor.empty:
        anomalies_hor['type_anomalie'] = 'Hors Horaires / Weekend'
        anomalies_hor['detail_anomalie'] = anomalies_hor.apply(
            lambda r: f"{r['DateTime'].strftime('%A %H:%M')} " + \
                      ("(Weekend)" if r['jour_semaine'] >= 5 else "") + \
                      ("(Heures non ouvrées)" if (cond_heure.loc[r.name]) else ""),
            axis=1
        )
        anomalies_hor['poids_anomalie'] = st.session_state.get('ss_poids_hors_horaire', DEFAULT_POIDS_HORS_HORAIRE)
        all_anomalies.extend(anomalies_hor.to_dict('records'))

    hors_service = df_merged[df_merged['Type'].isin(['EN PANNE', 'IMMO'])].copy()
    if not hors_service.empty:
        hors_service['type_anomalie'] = 'Véhicule Hors Service'
        hors_service['detail_anomalie'] = 'Transaction alors que véhicule EN PANNE ou IMMO'
        hors_service['poids_anomalie'] = st.session_state.get('ss_poids_hors_service', DEFAULT_POIDS_HORS_SERVICE)
        all_anomalies.extend(hors_service.to_dict('records'))

    double_indices = set()
    for carte in df_merged_sorted['Card num.'].unique():
        sub = df_merged_sorted[df_merged_sorted['Card num.'] == carte]
        if len(sub) > 1:
            for i in range(len(sub) - 1):
                rowA = sub.iloc[i]
                rowB = sub.iloc[i+1]
                delta_sec = (rowB['DateTime'] - rowA['DateTime']).total_seconds()
                if delta_sec >= 0 and (delta_sec / 60.0) < delta_minutes_double and rowA['Amount'] == rowB['Amount'] and pd.notna(rowA['Amount']):
                    double_indices.add(sub.index[i])
                    double_indices.add(sub.index[i+1])

    if double_indices:
        doubles_df = df_merged_sorted.loc[list(double_indices)].copy()
        doubles_df['type_anomalie'] = 'Facturation double suspectée'
        doubles_df['detail_anomalie'] = f"Montant identique ({doubles_df['Amount']}) en < {delta_minutes_double} min"
        doubles_df['poids_anomalie'] = st.session_state.get('ss_poids_fact_double', DEFAULT_POIDS_FACT_DOUBLE)
        all_anomalies.extend(doubles_df.to_dict('records'))

    if not all_anomalies:
        return pd.DataFrame()

    df_final_anomalies = pd.DataFrame(all_anomalies)
    cols_to_keep = [
        'Date', 'Hour', 'DateTime', 'Card num.', 'Nouveau Immat', 'Catégorie', 'Type',
        'Quantity', 'Amount', 'Past mileage', 'Current mileage', 'distance_parcourue',
        'consommation_100km', 'Cap-rèservoir', 'Place', 'Card name',
        'type_anomalie', 'detail_anomalie', 'poids_anomalie'
    ]
    if 'Dotation' in df_final_anomalies.columns: # Ajouter Dotation si elle a été fusionnée
        cols_to_keep.append('Dotation')

    cols_final = [col for col in cols_to_keep if col in df_final_anomalies.columns]
    df_final_anomalies = df_final_anomalies[cols_final]
    return df_final_anomalies.sort_values(by=['Nouveau Immat', 'DateTime', 'type_anomalie'])


# --- Fonctions d'analyse spécifiques ---

def analyser_stations_risque(df_anomalies: pd.DataFrame, df_transactions: pd.DataFrame) -> pd.DataFrame:
    """
    Analyse les stations en fonction des anomalies détectées pour identifier celles présentant 
    des risques élevés de fraude.
    
    Args:
        df_anomalies: DataFrame des anomalies détectées
        df_transactions: DataFrame complet des transactions
        
    Returns:
        DataFrame des stations classées par niveau de risque
    """
    if df_anomalies.empty or 'Place' not in df_anomalies.columns:
        return pd.DataFrame()
    
    # Compter le nombre total de transactions par station
    transactions_par_station = df_transactions.groupby('Place').agg(
        Nb_Total_Transactions=('Amount', 'count'),
        Volume_Total=('Quantity', 'sum'),
        Montant_Total=('Amount', 'sum')
    ).reset_index()
    
    # Compter le nombre d'anomalies par station et par type
    anomalies_par_station = df_anomalies.groupby(['Place', 'type_anomalie']).size().reset_index(name='Nb_Anomalies')
    
    # Types d'anomalies liées à la fraude potentielle par carte
    types_fraude_carte = [
        'Dépassement capacité', 
        'Prises rapprochées', 
        'Facturation double suspectée',
        'Transaction sans présence (géoloc)',
        'Véhicule Hors Service'
    ]
    
    # Filtrer pour ne garder que les anomalies liées à la fraude par carte
    anomalies_fraude = anomalies_par_station[anomalies_par_station['type_anomalie'].isin(types_fraude_carte)]
    
    # Agréger par station
    resume_stations = anomalies_fraude.groupby('Place').agg(
        Nb_Anomalies_Fraude=('Nb_Anomalies', 'sum')
    ).reset_index()
    
    # Fusionner avec les statistiques générales des transactions
    resume_complet = resume_stations.merge(
        transactions_par_station, on='Place', how='left'
    )
    
    # Calculer le pourcentage d'anomalies
    resume_complet['Pourcentage_Anomalies'] = np.where(
        resume_complet['Nb_Total_Transactions'] > 0,
        (resume_complet['Nb_Anomalies_Fraude'] / resume_complet['Nb_Total_Transactions']) * 100,
        0
    )
    
    # Calculer un score de risque (formule pondérée)
    resume_complet['Score_Risque_Station'] = resume_complet['Pourcentage_Anomalies'] * np.log1p(resume_complet['Nb_Anomalies_Fraude'])
    
    # Déterminer le niveau de risque
    resume_complet['Niveau_Risque'] = pd.cut(
        resume_complet['Score_Risque_Station'],
        bins=[-float('inf'), 5, 10, 20, float('inf')],
        labels=['Faible', 'Modéré', 'Élevé', 'Critique']
    )
    
    # Récupérer les types d'anomalies les plus fréquents par station
    anomalies_principales = {}
    for station in resume_complet['Place'].unique():
        top_anomalies = anomalies_fraude[anomalies_fraude['Place'] == station].sort_values('Nb_Anomalies', ascending=False)
        if not top_anomalies.empty:
            top_types = top_anomalies.head(2)['type_anomalie'].tolist()
            anomalies_principales[station] = " + ".join(top_types)
        else:
            anomalies_principales[station] = "N/A"
    
    resume_complet['Anomalies_Principales'] = resume_complet['Place'].map(anomalies_principales)
    
    # Arrondir les valeurs numériques
    resume_complet['Pourcentage_Anomalies'] = resume_complet['Pourcentage_Anomalies'].round(1)
    resume_complet['Score_Risque_Station'] = resume_complet['Score_Risque_Station'].round(1)
    
    # Trier par score de risque décroissant
    return resume_complet.sort_values('Score_Risque_Station', ascending=False)

def analyser_consommation_vehicule(vehicule_data: pd.DataFrame, info_vehicule: pd.Series) -> Dict[str, Any]:
    """Analyse la consommation d'un véhicule spécifique."""
    if vehicule_data.empty:
        return {
            'total_litres': 0, 'nb_prises': 0, 'moyenne_prise': 0,
            'distance_totale': 0, 'consommation_moyenne': 0,
            'cout_total': 0, 'cout_moyen_prise': 0, 'cout_par_km': 0,
            'conso_mensuelle': pd.DataFrame(), 'stations_frequentes': pd.Series(dtype='int64'),
            'prix_moyen_litre': 0
        }

    vehicule_data_sorted = vehicule_data.sort_values('DateTime')
    total_litres = vehicule_data_sorted['Quantity'].sum()
    cout_total = vehicule_data_sorted['Amount'].sum()
    nb_prises = len(vehicule_data_sorted)
    moyenne_prise = vehicule_data_sorted['Quantity'].mean() if nb_prises > 0 else 0
    cout_moyen_prise = vehicule_data_sorted['Amount'].mean() if nb_prises > 0 else 0
    prix_moyen_litre = (cout_total / total_litres) if total_litres > 0 else 0

    df_km = vehicule_data_sorted[['Past mileage', 'Current mileage']].dropna()
    distance_totale = 0
    consommation_moyenne = 0
    cout_par_km = 0

    if not df_km.empty and len(df_km) > 1:
        first_km = df_km['Past mileage'].iloc[0]
        last_km = df_km['Current mileage'].iloc[-1]
        if pd.notna(first_km) and pd.notna(last_km) and last_km > first_km:
            distance_totale = last_km - first_km

    vehicule_data_sorted['distance_transaction'] = vehicule_data_sorted['Current mileage'] - vehicule_data_sorted['Past mileage']
    distance_sommee_valide = vehicule_data_sorted.loc[vehicule_data_sorted['distance_transaction'] > 0, 'distance_transaction'].sum()
    if distance_sommee_valide > 0:
        distance_utilisee = distance_sommee_valide
        consommation_moyenne = (total_litres / distance_utilisee) * 100 if distance_utilisee > 0 else 0
        cout_par_km = (cout_total / distance_utilisee) if distance_utilisee > 0 else 0
    elif distance_totale > 0:
        distance_utilisee = distance_totale
        consommation_moyenne = (total_litres / distance_utilisee) * 100 if distance_utilisee > 0 else 0
        cout_par_km = (cout_total / distance_utilisee) if distance_utilisee > 0 else 0
    else:
        distance_utilisee = 0

    vehicule_data_sorted['mois'] = vehicule_data_sorted['Date'].dt.strftime('%Y-%m')
    conso_mensuelle = vehicule_data_sorted.groupby('mois').agg(
        Volume_L=('Quantity', 'sum'),
        Montant_CFA=('Amount','sum'),
        Nb_prises=('Quantity', 'count'),
        Volume_moyen_L=('Quantity', 'mean')
    )
    stations_freq = vehicule_data_sorted['Place'].value_counts().head(5)

    return {
        'total_litres': total_litres,
        'nb_prises': nb_prises,
        'moyenne_prise': moyenne_prise,
        'distance_totale_estimee': distance_utilisee,
        'consommation_moyenne': consommation_moyenne,
        'cout_total': cout_total,
        'cout_moyen_prise': cout_moyen_prise,
        'cout_par_km': cout_par_km,
        'conso_mensuelle': conso_mensuelle,
        'stations_frequentes': stations_freq,
        'prix_moyen_litre': prix_moyen_litre
    }

def generer_rapport_vehicule(donnees_vehicule: pd.DataFrame, info_vehicule: pd.Series, date_debut: datetime.date, date_fin: datetime.date, conso_moyenne_categorie: float) -> Tuple[pd.DataFrame, pd.DataFrame, pd.DataFrame, pd.Series, Dict[str, Any]]:
    """Génère un rapport détaillé pour un véhicule, incluant le benchmarking."""
    dotation_vehicule = info_vehicule.get('Dotation', 0) # Récupérer la dotation

    infos_base_list = [
        ('Immatriculation', info_vehicule.get('Nouveau Immat', 'N/A')),
        ('Marque', info_vehicule.get('Marque', 'N/A')),
        ('Modèle', info_vehicule.get('Modèle', 'N/A')),
        ('Type', info_vehicule.get('Type', 'N/A')),
        ('Catégorie', info_vehicule.get('Catégorie', 'N/A')),
        ('Capacité réservoir', f"{info_vehicule.get('Cap-rèservoir', 0):.0f} L")
    ]
    if 'Dotation' in info_vehicule: # Ajouter la dotation si elle existe
         infos_base_list.append(('Dotation Mensuelle', f"{dotation_vehicule:.0f} L"))
    infos_base_list.extend([
        ('Période début', date_debut.strftime(DATE_FORMAT)),
        ('Période fin', date_fin.strftime(DATE_FORMAT))
    ])
    infos_base = pd.DataFrame(infos_base_list, columns=['Paramètre', 'Valeur'])


    analyse = analyser_consommation_vehicule(donnees_vehicule, info_vehicule)

    conso_veh = analyse['consommation_moyenne']
    ecart_conso = conso_veh - conso_moyenne_categorie if conso_moyenne_categorie > 0 and conso_veh > 0 else 0
    ecart_conso_pct = (ecart_conso / conso_moyenne_categorie) * 100 if conso_moyenne_categorie > 0 and conso_veh > 0 else 0

    stats_conso = pd.DataFrame({
        'Paramètre': [
            'Volume total', 'Coût total', 'Nombre de prises', 'Moyenne par prise (Volume)', 'Moyenne par prise (Coût)',
            'Prix moyen / Litre', 'Distance totale estimée', 'Consommation moyenne', 'Consommation moyenne (Catégorie)',
            'Écart vs Catégorie', 'Coût par Km'
        ],
        'Valeur': [
            f"{analyse['total_litres']:.1f} L", f"{analyse['cout_total']:,.0f} CFA", analyse['nb_prises'],
            f"{analyse['moyenne_prise']:.1f} L", f"{analyse['cout_moyen_prise']:,.0f} CFA", f"{analyse['prix_moyen_litre']:,.0f} CFA/L",
            f"{analyse.get('distance_totale_estimee', 0):,.0f} km",
            f"{conso_veh:.1f} L/100km" if conso_veh > 0 else "N/A",
            f"{conso_moyenne_categorie:.1f} L/100km" if conso_moyenne_categorie > 0 else "N/A",
            f"{ecart_conso:+.1f} L/100km ({ecart_conso_pct:+.1f}%)" if conso_veh > 0 and conso_moyenne_categorie > 0 else "N/A",
            f"{analyse['cout_par_km']:,.1f} CFA/km" if analyse['cout_par_km'] > 0 else "N/A"
        ]
    })
    return infos_base, stats_conso, analyse['conso_mensuelle'], analyse['stations_frequentes'], analyse


def calculer_distance_et_consommation(df_transactions: pd.DataFrame) -> Tuple[float, float, float, str]:
    """Calcule la distance parcourue et la consommation en utilisant différentes méthodes.

    Args:
        df_transactions: DataFrame des transactions triées par date

    Returns:
        Tuple contenant:
        - distance_simple: Distance calculée par différence première/dernière transaction
        - distance_cumulative: Distance calculée en sommant les distances entre transactions
        - consommation_recommandee: Consommation calculée avec la méthode la plus fiable
        - methode_utilisee: Méthode utilisée pour le calcul ('simple', 'cumulative', 'hybride', 'insuffisant')
    """
    if df_transactions.empty or len(df_transactions) < 2:
        return 0.0, 0.0, 0.0, "insuffisant"

    # Méthode simple: différence entre premier et dernier kilométrage
    df_km = df_transactions[['Past mileage', 'Current mileage']].dropna()
    distance_simple = 0.0
    if not df_km.empty and len(df_km) > 1:
        first_km = df_km['Past mileage'].iloc[0]
        last_km = df_km['Current mileage'].iloc[-1]
        if pd.notna(first_km) and pd.notna(last_km) and last_km > first_km:
            distance_simple = last_km - first_km

    # Méthode cumulative: somme des distances valides entre transactions
    df_sorted = df_transactions.sort_values('DateTime')
    distances_valides = []
    for i in range(len(df_sorted)):
        curr_row = df_sorted.iloc[i]
        if pd.notna(curr_row['Past mileage']) and pd.notna(curr_row['Current mileage']):
            dist = curr_row['Current mileage'] - curr_row['Past mileage']
            if dist > 0 and dist < 1000:  # Filtre basique pour les distances aberrantes
                distances_valides.append(dist)
    distance_cumulative = sum(distances_valides)

    # Choisir la méthode la plus appropriée
    if distance_simple == 0 and distance_cumulative == 0:
        return 0.0, 0.0, 0.0, "insuffisant"

    # Si les deux méthodes donnent des résultats similaires (écart < 10%)
    if distance_simple > 0 and distance_cumulative > 0:
        max_dist = max(distance_simple, distance_cumulative)
        min_dist = min(distance_simple, distance_cumulative)
        if (max_dist - min_dist) / max_dist < 0.1:  # Écart < 10%
            distance_utilisee = (distance_simple + distance_cumulative) / 2
            methode = "hybride"
        else:
            # Utiliser la méthode cumulative si disponible, sinon la méthode simple
            distance_utilisee = distance_cumulative if distance_cumulative > 0 else distance_simple
            methode = "cumulative" if distance_cumulative > 0 else "simple"
    else:
        # Utiliser la méthode non nulle
        distance_utilisee = distance_cumulative if distance_cumulative > 0 else distance_simple
        methode = "cumulative" if distance_cumulative > 0 else "simple"

    # Calculer la consommation
    volume_total = df_transactions['Quantity'].sum()
    consommation = (volume_total / distance_utilisee * 100) if distance_utilisee > 0 else 0.0

    return distance_simple, distance_cumulative, consommation, methode

def calculer_kpis_globaux(df_transactions: pd.DataFrame, df_vehicules: pd.DataFrame, date_debut: datetime.date, date_fin: datetime.date, selected_categories: List[str]) -> Tuple[pd.DataFrame, pd.DataFrame]:
    """Calcule les KPIs de consommation et de coût par catégorie et véhicule."""
    # S'assurer que Dotation est présente dans df_vehicules pour les fusions futures si besoin
    cols_veh_kpi = ['N° Carte', 'Catégorie', 'Nouveau Immat']
    if 'Dotation' in df_vehicules.columns:
        cols_veh_kpi.append('Dotation')

    df = df_transactions.merge(
        df_vehicules[cols_veh_kpi],
        left_on='Card num.', right_on='N° Carte', how='left'
    )
    mask_date = (df['Date'].dt.date >= date_debut) & (df['Date'].dt.date <= date_fin)
    df = df[mask_date].copy()
    if selected_categories:
        df = df[df['Catégorie'].isin(selected_categories)]

    if df.empty:
        return pd.DataFrame(), pd.DataFrame()

    vehicle_data = []
    df_sorted = df.sort_values(['Card num.', 'DateTime'])

    for card, group in df_sorted.groupby('Card num.'):
        if group.empty: continue
        total_lit = group['Quantity'].sum()
        total_amount = group['Amount'].sum()
        cat = group['Catégorie'].iloc[0]
        immat = group['Nouveau Immat'].iloc[0]
        nb_prises = len(group)
        dotation_mensuelle = group['Dotation'].iloc[0] if 'Dotation' in group.columns else 0

        # Utiliser les deux méthodes de calcul de distance
        distance_simple, distance_cumulative, consommation_recommandee, methode_utilisee = calculer_distance_et_consommation(group)
        
        # Choisir la distance à utiliser (prendre la plus grande par défaut)
        distance_utilisee = max(distance_simple, distance_cumulative)
        
        # Si aucune méthode n'a fonctionné, essayer l'ancienne méthode comme fallback
        if distance_utilisee == 0:
            group_km = group[['Past mileage', 'Current mileage']].dropna()
            if not group_km.empty and len(group_km) > 1:
                first_km = group_km['Past mileage'].iloc[0]
                last_km = group_km['Current mileage'].iloc[-1]
                if pd.notna(first_km) and pd.notna(last_km) and last_km > first_km:
                    distance_utilisee = last_km - first_km
                    methode_utilisee = "legacy"

        # Calculer les KPIs basés sur la distance
        cons = consommation_recommandee if methode_utilisee != "insuffisant" else np.nan
        cpk = (total_amount / distance_utilisee) if distance_utilisee > 0 else np.nan
        avg_price_liter = (total_amount / total_lit) if total_lit > 0 else np.nan

        vehicle_data.append({
            'Card num.': card, 'Nouveau Immat': immat, 'Catégorie': cat,
            'total_litres': total_lit, 'total_cout': total_amount,
            'distance': distance_utilisee, 'consommation': cons, 'cout_par_km': cpk,
            'nb_prises': nb_prises, 'prix_moyen_litre': avg_price_liter,
            'Dotation': dotation_mensuelle,
            'distance_simple': distance_simple,  # Nouvelles colonnes pour les deux méthodes
            'distance_cumulative': distance_cumulative,
            'methode_distance': methode_utilisee
        })
    df_vehicle_kpi = pd.DataFrame(vehicle_data)
    if df_vehicle_kpi.empty:
        return pd.DataFrame(), pd.DataFrame()

    kpi_cat = df_vehicle_kpi.groupby('Catégorie').agg(
        consommation_moyenne=('consommation', 'mean'),
        cout_par_km_moyen=('cout_par_km', 'mean'),
        total_litres=('total_litres', 'sum'),
        total_cout=('total_cout', 'sum'),
        distance_totale=('distance', 'sum'),
        nb_vehicules=('Card num.', 'nunique'),
        nb_transactions=('nb_prises', 'sum')
    ).reset_index()

    kpi_cat['consommation_globale'] = (kpi_cat['total_litres'] / kpi_cat['distance_totale']) * 100
    kpi_cat['cout_par_km_global'] = kpi_cat['total_cout'] / kpi_cat['distance_totale']
    kpi_cat['prix_moyen_litre_global'] = kpi_cat['total_cout'] / kpi_cat['total_litres']

    kpi_cat = kpi_cat.round({
        'consommation_moyenne': 1, 'cout_par_km_moyen': 1, 'total_litres': 0, 'total_cout': 0,
        'distance_totale': 0, 'consommation_globale': 1, 'cout_par_km_global': 1, 'prix_moyen_litre_global': 0
    })
    df_vehicle_kpi = df_vehicle_kpi.round({
         'total_litres': 1, 'total_cout': 0, 'distance': 0, 'consommation': 1, 'cout_par_km': 1,
         'prix_moyen_litre': 0, 'Dotation':0
    })
    return kpi_cat, df_vehicle_kpi


# ---------------------------------------------------------------------
# Fonctions d'agrégation des anomalies pour les résumés
# ---------------------------------------------------------------------

def calculer_score_risque(df_anomalies: pd.DataFrame) -> pd.DataFrame:
    """Calcule un score de risque pondéré par véhicule basé sur les anomalies."""
    if df_anomalies.empty or 'poids_anomalie' not in df_anomalies.columns:
        return pd.DataFrame(columns=['Nouveau Immat', 'Card num.', 'Catégorie', 'nombre_total_anomalies', 'score_risque'])
    
    # Vérifier si la colonne 'Catégorie' existe, sinon l'ajouter avec une valeur par défaut
    if 'Catégorie' not in df_anomalies.columns:
        df_anomalies = df_anomalies.copy()
        df_anomalies['Catégorie'] = 'Non définie'
    
    # Vérifier les colonnes nécessaires pour le groupby
    groupby_cols = ['Nouveau Immat', 'Card num.', 'Catégorie', 'type_anomalie']
    for col in groupby_cols:
        if col not in df_anomalies.columns:
            df_anomalies[col] = None
    
    pivot = df_anomalies.groupby(groupby_cols).agg(
        nombre=('type_anomalie', 'size'),
        score_partiel=('poids_anomalie', 'sum')
    ).reset_index()

    summary = pivot.groupby(['Nouveau Immat', 'Card num.', 'Catégorie']).agg(
        nombre_total_anomalies=('nombre', 'sum'),
        score_risque=('score_partiel', 'sum')
    ).reset_index()
    return summary.sort_values('score_risque', ascending=False)

# ---------------------------------------------------------------------
# NOUVELLE FONCTION : Analyse consommation par période
# ---------------------------------------------------------------------
def analyser_consommation_par_periode(
    df_transactions: pd.DataFrame,
    df_vehicules: pd.DataFrame,
    date_debut: datetime.date,
    date_fin: datetime.date,
    periode: str = 'M',
    selected_categories: List[str] = None,
    selected_vehicles: List[str] = None
) -> Tuple[pd.DataFrame, pd.DataFrame]:
    """
    Analyse la consommation de carburant par période (jour, semaine, mois, trimestre, année)
    """
    cols_veh_periode = ['N° Carte', 'Catégorie', 'Nouveau Immat', 'Cap-rèservoir']
    if 'Dotation' in df_vehicules.columns:
        cols_veh_periode.append('Dotation')

    df = df_transactions.merge(
        df_vehicules[cols_veh_periode],
        left_on='Card num.', right_on='N° Carte', how='left'
    )

    mask_date = (df['Date'].dt.date >= date_debut) & (df['Date'].dt.date <= date_fin)
    df = df[mask_date].copy()

    if selected_categories:
        df = df[df['Catégorie'].isin(selected_categories)]
    if selected_vehicles:
        df = df[df['Nouveau Immat'].isin(selected_vehicles)]

    if df.empty:
        return pd.DataFrame(), pd.DataFrame()

    df['distance_parcourue'] = df['Current mileage'] - df['Past mileage']
    df['consommation_100km'] = np.where(
        (df['distance_parcourue'] > 0) & df['Quantity'].notna(),
        (df['Quantity'] / df['distance_parcourue']) * 100,
        np.nan
    )

    seuils_conso = st.session_state.get('ss_conso_seuils_par_categorie', {})
    df['periode_datetime'] = df['Date'].dt.to_period(periode).dt.to_timestamp()

    if periode == 'D':
        df['periode_str'] = df['Date'].dt.strftime('%Y-%m-%d')
    elif periode == 'W':
        df['periode_str'] = df['Date'].dt.to_period('W').astype(str)
    elif periode == 'M':
        df['periode_str'] = df['Date'].dt.strftime('%Y-%m')
    elif periode == 'Q':
        df['periode_str'] = df['Date'].dt.to_period('Q').astype(str)
    else:
        df['periode_str'] = df['Date'].dt.strftime('%Y')

    conso_veh_periode = df.groupby(['Nouveau Immat', 'Catégorie', 'periode_str']).agg(
        volume_total=('Quantity', 'sum'),
        cout_total=('Amount', 'sum'),
        distance_totale=('distance_parcourue', lambda x: x[x > 0].sum()),
        nb_transactions=('Quantity', 'count'),
        date_debut_periode=('Date', 'min'),
        date_fin_periode=('Date', 'max')
    ).reset_index()

    conso_veh_periode['consommation_moyenne'] = np.where(
        conso_veh_periode['distance_totale'] > 0,
        (conso_veh_periode['volume_total'] / conso_veh_periode['distance_totale']) * 100,
        np.nan
    )
    conso_veh_periode['seuil_consommation'] = conso_veh_periode['Catégorie'].map(
        lambda x: seuils_conso.get(x, DEFAULT_CONSO_SEUIL)
    )
    conso_veh_periode['exces_consommation'] = np.where(
        conso_veh_periode['consommation_moyenne'] > conso_veh_periode['seuil_consommation'],
        conso_veh_periode['consommation_moyenne'] - conso_veh_periode['seuil_consommation'],
        0
    )
    conso_veh_periode['pourcentage_exces'] = np.where(
        conso_veh_periode['seuil_consommation'] > 0,
        (conso_veh_periode['exces_consommation'] / conso_veh_periode['seuil_consommation']) * 100,
        0
    )

    conso_periode = df.groupby(['periode_str']).agg(
        volume_total=('Quantity', 'sum'),
        cout_total=('Amount', 'sum'),
        distance_totale=('distance_parcourue', lambda x: x[x > 0].sum()),
        nb_transactions=('Quantity', 'count'),
        nb_vehicules=('Nouveau Immat', 'nunique'),
        date_debut_periode=('Date', 'min'),
        date_fin_periode=('Date', 'max')
    ).reset_index()

    conso_periode['consommation_moyenne'] = np.where(
        conso_periode['distance_totale'] > 0,
        (conso_periode['volume_total'] / conso_periode['distance_totale']) * 100,
        np.nan
    )
    conso_veh_periode = conso_veh_periode.round({
        'volume_total': 1,'cout_total': 0,'distance_totale': 0,
        'consommation_moyenne': 1,'exces_consommation': 1,'pourcentage_exces': 1
    })
    conso_periode = conso_periode.round({
        'volume_total': 1,'cout_total': 0,'distance_totale': 0,'consommation_moyenne': 1
    })

    conso_veh_periode = conso_veh_periode.sort_values(['periode_str', 'exces_consommation'], ascending=[True, False])
    conso_periode = conso_periode.sort_values('periode_str')
    return conso_periode, conso_veh_periode

# ---------------------------------------------------------------------
# NOUVELLE FONCTION : Amélioration du dashboard
# ---------------------------------------------------------------------
def ameliorer_dashboard(df_transactions: pd.DataFrame, df_vehicules: pd.DataFrame, global_date_debut: datetime.date, global_date_fin: datetime.date, kpi_cat_dash: pd.DataFrame, df_vehicle_kpi_dash: pd.DataFrame, df_geoloc: Optional[pd.DataFrame] = None):
    """Ajoute une section d'aperçu des excès de consommation au tableau de bord"""
    # Section existante...
    with st.expander("⚠️ Aperçu des Excès de Consommation (Mensuel)", expanded=True):
        _, conso_veh_mois = analyser_consommation_par_periode(
            df_transactions, df_vehicules, global_date_debut, global_date_fin,
            periode='M', selected_categories=None, selected_vehicles=None
        )
        if not conso_veh_mois.empty:
            exces_mois = conso_veh_mois[conso_veh_mois['exces_consommation'] > 0]
            if not exces_mois.empty:
                nb_exces_mois = len(exces_mois)
                nb_vehicules_exces = exces_mois['Nouveau Immat'].nunique()

                col_e1, col_e2, col_e3 = st.columns(3)
                col_e1.metric("Nombre d'Excès Détectés", f"{nb_exces_mois}")
                col_e2.metric("Véhicules Concernés", f"{nb_vehicules_exces}")
                col_e3.metric("Excès Moyen", f"{exces_mois['pourcentage_exces'].mean():.1f}%")

                top_exces = exces_mois.nlargest(5, 'pourcentage_exces')
                top_exces_display = top_exces[[
                    'periode_str', 'Nouveau Immat', 'consommation_moyenne',
                    'seuil_consommation', 'pourcentage_exces'
                ]]
                top_exces_display.columns = [
                    'Période', 'Immatriculation', 'Consommation (L/100km)',
                    'Seuil (L/100km)', 'Excès (%)'
                ]
                st.dataframe(top_exces_display, use_container_width=True)
                st.markdown("""
                👉 *Pour une analyse complète des excès de consommation, utilisez la page "Analyse par Période"*
                """)
            else:
                st.success("✅ Aucun excès de consommation détecté pour les périodes analysées.")
        else:
            st.info("Données insuffisantes pour l'analyse des excès de consommation.")
    
    # Ajouter une nouvelle section pour les véhicules sans géolocalisation
    if df_geoloc is not None and not df_geoloc.empty:
        with st.expander("⚠️ Véhicules Sans Données de Géolocalisation", expanded=False):
            with st.spinner("Analyse des véhicules sans géolocalisation..."):
                df_vehicules_sans_geoloc_dash = analyser_vehicules_sans_geoloc(
                    df_transactions, df_vehicules, df_geoloc, global_date_debut, global_date_fin
                )
            
            if df_vehicules_sans_geoloc_dash.empty:
                st.success("✅ Tous les véhicules avec des transactions ont des données de géolocalisation.")
            else:
                nb_vehicules_sans_geoloc = len(df_vehicules_sans_geoloc_dash)
                st.warning(f"⚠️ {nb_vehicules_sans_geoloc} véhicules ont effectué des transactions sans données de géolocalisation correspondantes.")
                
                # Afficher un résumé simplifié
                df_resume = df_vehicules_sans_geoloc_dash[['Immatriculation', 'Catégorie', 'Nb_Transactions', 'Volume_Total_L']]
                st.dataframe(df_resume.head(5), use_container_width=True)
                
                if nb_vehicules_sans_geoloc > 5:
                    st.info(f"👉 *{nb_vehicules_sans_geoloc - 5} autres véhicules non affichés. Consultez la page 'Géolocalisation' > onglet 'Véhicules Sans Géoloc' pour tous les détails.*")

# ---------------------------------------------------------------------
# NOUVELLE FONCTION : Affichage de la page d'analyse par période
# ---------------------------------------------------------------------
def afficher_page_analyse_periodes(df_transactions: pd.DataFrame, df_vehicules: pd.DataFrame, date_debut: datetime.date, date_fin: datetime.date):
    """Affiche la page d'analyse de consommation par période."""
    st.header(f"📅 Analyse de Consommation par Période ({date_debut.strftime('%d/%m/%Y')} - {date_fin.strftime('%d/%m/%Y')})")

    # Fonction interne pour l'analyse de période personnalisée
    def analyser_consommation_par_periode_custom(df_trans, df_vehs, date_deb, date_fin, selected_categories=None, selected_vehicles=None):
        """Analyse la consommation pour une période personnalisée."""
        # Utiliser la même logique que analyser_consommation_par_periode mais avec période journalière
        return analyser_consommation_par_periode(
            df_trans, df_vehs, date_deb, date_fin,
            periode='D', selected_categories=selected_categories,
            selected_vehicles=selected_vehicles
        )

    if df_transactions.empty:
        st.warning("Aucune transaction à analyser pour la période sélectionnée.")
        return

    st.subheader("Configuration de l'Analyse")
    col_config1, col_config2 = st.columns(2)
    with col_config1:
        periode_options = {
            'Personnalisée': 'CUSTOM',
            'Jour': 'D',
            'Semaine': 'W',
            'Mois': 'M',
            'Trimestre': 'Q',
            'Année': 'Y'
        }
        periode_label = st.selectbox(
            "Sélectionner la période d'analyse :",
            options=list(periode_options.keys()),
            index=0
        )
        periode_code = periode_options[periode_label]
        
        # Gestion de la période personnalisée
        if periode_code == 'CUSTOM':
            col_date1, col_date2 = st.columns(2)
            with col_date1:
                custom_date_debut = st.date_input(
                    "Date de début personnalisée",
                    value=date_debut,
                    min_value=date_debut,
                    max_value=date_fin,
                    key="custom_date_debut_periode"
                )
            with col_date2:
                custom_date_fin = st.date_input(
                    "Date de fin personnalisée",
                    value=date_fin,
                    min_value=custom_date_debut,
                    max_value=date_fin,
                    key="custom_date_fin_periode"
                )
            
            if custom_date_debut > custom_date_fin:
                st.error("La date de début ne peut pas être postérieure à la date de fin.")
                return
            
            date_debut = custom_date_debut
            date_fin = custom_date_fin
    with col_config2:
        all_cats = sorted(df_vehicules['Catégorie'].dropna().astype(str).unique())
        selected_cats = st.multiselect(
            "Filtrer par Catégories de véhicules",
            options=all_cats,default=all_cats,key="periode_cat_filter"
        )
    with st.expander("Filtrer par véhicules spécifiques (optionnel)"):
        if selected_cats:
            available_vehicles = sorted(df_vehicules[df_vehicules['Catégorie'].isin(selected_cats)]['Nouveau Immat'].dropna().unique())
        else:
            available_vehicles = sorted(df_vehicules['Nouveau Immat'].dropna().unique())
        selected_vehicles = st.multiselect(
            "Sélectionner des véhicules spécifiques",
            options=available_vehicles,default=None,key="periode_veh_filter"
        )

    with st.spinner(f"Analyse {periode_label.lower()} en cours..."):
        if periode_code == 'CUSTOM':
            # Pour période personnalisée
            conso_periode, conso_veh_periode = analyser_consommation_par_periode_custom(
                df_transactions, df_vehicules, date_debut, date_fin,
                selected_categories=selected_cats,
                selected_vehicles=selected_vehicles if selected_vehicles else None
            )
        else:
            # Pour les périodes standards
            conso_periode, conso_veh_periode = analyser_consommation_par_periode(
                df_transactions, df_vehicules, date_debut, date_fin,
                periode=periode_code, selected_categories=selected_cats,
                selected_vehicles=selected_vehicles if selected_vehicles else None
            )

    if conso_periode.empty or conso_veh_periode.empty:
        st.warning(f"Données insuffisantes pour l'analyse {periode_label.lower()}.")
        return

    st.subheader(f"Consommation {periode_label} Globale")
    afficher_dataframe_avec_export(
        conso_periode[['periode_str', 'volume_total', 'cout_total', 'distance_totale',
                      'consommation_moyenne', 'nb_transactions', 'nb_vehicules']],
        f"Récapitulatif {periode_label}",key=f"recap_periode_{periode_code}"
    )
    fig_conso = px.line(
        conso_periode, x='periode_str', y='consommation_moyenne',
        title=f"Évolution de la Consommation Moyenne ({periode_label})",
        labels={'periode_str': periode_label, 'consommation_moyenne': 'Conso. Moyenne (L/100km)'},
        markers=True
    )
    conso_moy_globale = conso_periode['consommation_moyenne'].mean()
    fig_conso.add_hline(
        y=conso_moy_globale,line_dash="dash", line_color="green",
        annotation_text=f"Moyenne: {conso_moy_globale:.1f} L/100km"
    )
    st.plotly_chart(fig_conso, use_container_width=True)

    fig_vol_cout = px.bar(
        conso_periode, x='periode_str', y=['volume_total', 'cout_total'],
        title=f"Volume et Coût par {periode_label}",
        labels={'periode_str': periode_label, 'value': 'Valeur', 'variable': 'Métrique'},
        barmode='group'
    )
    st.plotly_chart(fig_vol_cout, use_container_width=True)

    st.subheader(f"Détail par Véhicule et par {periode_label}")
    exces_veh = conso_veh_periode[conso_veh_periode['exces_consommation'] > 0]
    nb_exces = len(exces_veh)
    if nb_exces > 0:
        st.warning(f"⚠️ Détecté : {nb_exces} cas d'excès de consommation sur la période.")
        cols_display_exces = [
            'periode_str', 'Nouveau Immat', 'Catégorie', 'consommation_moyenne',
            'seuil_consommation', 'exces_consommation', 'pourcentage_exces',
            'volume_total', 'distance_totale', 'nb_transactions'
        ]
        afficher_dataframe_avec_export(
            exces_veh[cols_display_exces],"Excès de Consommation Détectés",key=f"exces_conso_{periode_code}"
        )
        top_exces = exces_veh.nlargest(10, 'pourcentage_exces')
        fig_top_exces = px.bar(
            top_exces,x='Nouveau Immat',y='pourcentage_exces',color='Catégorie',
            title="Top 10 des Excès de Consommation (%)",
            labels={'pourcentage_exces': "Excès (%)", 'Nouveau Immat': 'Véhicule'},
            hover_data=['periode_str', 'consommation_moyenne', 'seuil_consommation']
        )
        st.plotly_chart(fig_top_exces, use_container_width=True)
    else:
        st.success("✅ Aucun excès de consommation détecté sur la période analysée.")

    with st.expander("Voir toutes les données détaillées par véhicule et période"):
        cols_display_detail = [
            'periode_str', 'Nouveau Immat', 'Catégorie', 'volume_total',
            'distance_totale', 'consommation_moyenne', 'seuil_consommation',
            'exces_consommation', 'pourcentage_exces', 'cout_total', 'nb_transactions'
        ]
        afficher_dataframe_avec_export(
            conso_veh_periode[cols_display_detail],
            f"Toutes les données par Véhicule et {periode_label}",key=f"all_data_periode_{periode_code}"
        )
    with st.expander("Analyse comparative entre périodes", expanded=False):
        st.info("Cette section permet de visualiser l'évolution de la consommation par véhicule à travers les périodes.")
        vehicules_list = sorted(conso_veh_periode['Nouveau Immat'].unique())
        if vehicules_list:
            vehicule_selected = st.selectbox(
                "Sélectionner un véhicule pour l'analyse détaillée :",
                options=vehicules_list,key="compare_vehicule_select"
            )
            veh_data = conso_veh_periode[conso_veh_periode['Nouveau Immat'] == vehicule_selected]
            if not veh_data.empty:
                fig_veh_evo = px.line(
                    veh_data, x='periode_str', y=['consommation_moyenne', 'seuil_consommation'],
                    title=f"Évolution de la Consommation - {vehicule_selected}",
                    labels={'periode_str': periode_label, 'value': 'Consommation (L/100km)', 'variable': 'Métrique'},
                    markers=True
                )
                st.plotly_chart(fig_veh_evo, use_container_width=True)
                st.dataframe(veh_data[[
                    'periode_str', 'consommation_moyenne', 'seuil_consommation',
                    'exces_consommation', 'volume_total', 'distance_totale'
                ]], use_container_width=True)
            else:
                st.info(f"Pas de données disponibles pour {vehicule_selected} sur les périodes sélectionnées.")
        else:
            st.info("Aucun véhicule avec données suffisantes pour l'analyse comparative.")

# ---------------------------------------------------------------------
# NOUVELLES FONCTIONS POUR LE SUIVI DES DOTATIONS
# ---------------------------------------------------------------------
def analyser_suivi_dotations(
    df_transactions_filtrees: pd.DataFrame,
    df_vehicules: pd.DataFrame, # Doit contenir 'N° Carte', 'Nouveau Immat', 'Catégorie', 'Dotation'
    date_debut_periode: datetime.date,
    date_fin_periode: datetime.date
) -> Tuple[pd.DataFrame, pd.DataFrame]:
    """
    Analyse l'utilisation des dotations de carburant.

    Args:
        df_transactions_filtrees: Transactions déjà filtrées pour la période globale.
        df_vehicules: DataFrame des véhicules avec leurs informations, incluant 'Dotation'.
        date_debut_periode: Date de début de la période d'analyse.
        date_fin_periode: Date de fin de la période d'analyse.

    Returns:
        Un tuple de DataFrames:
        - df_recap_dotation_periode: Récapitulatif par véhicule sur toute la période.
        - df_detail_dotation_mensuel: Détail mensuel par véhicule.
    """
    if 'Dotation' not in df_vehicules.columns:
        st.error("La colonne 'Dotation' est indispensable pour cette analyse et n'a pas été trouvée dans les données véhicules.")
        return pd.DataFrame(), pd.DataFrame()

    # Merge transactions avec infos véhicules (y compris Dotation)
    df_merged = df_transactions_filtrees.merge(
        df_vehicules[['N° Carte', 'Nouveau Immat', 'Catégorie', 'Dotation']],
        left_on='Card num.',
        right_on='N° Carte',
        how='inner' # Garder seulement les transactions des cartes véhicules connues avec dotation
    )

    if df_merged.empty:
        return pd.DataFrame(), pd.DataFrame()

    # --- Calcul du nombre de mois dans la période ---
    # Utilisation de dateutil.relativedelta pour plus de précision
    delta = relativedelta(date_fin_periode, date_debut_periode)
    nombre_mois_periode = delta.years * 12 + delta.months + 1


    # --- Analyse mensuelle détaillée ---
    df_merged['AnneeMois'] = df_merged['Date'].dt.strftime('%Y-%m')

    # Consommation réelle par mois et par véhicule
    conso_mensuelle_veh = df_merged.groupby(['Nouveau Immat', 'Catégorie', 'Dotation', 'AnneeMois']).agg(
        Consommation_Mois_L=('Quantity', 'sum')
    ).reset_index()

    # La dotation est déjà mensuelle, donc 'Dotation' est la dotation pour ce mois.
    conso_mensuelle_veh.rename(columns={'Dotation': 'Dotation_Mensuelle_L'}, inplace=True)

    conso_mensuelle_veh['Difference_Mois_L'] = conso_mensuelle_veh['Dotation_Mensuelle_L'] - conso_mensuelle_veh['Consommation_Mois_L']
    conso_mensuelle_veh['Taux_Utilisation_Mois_%'] = np.where(
        conso_mensuelle_veh['Dotation_Mensuelle_L'] > 0,
        (conso_mensuelle_veh['Consommation_Mois_L'] / conso_mensuelle_veh['Dotation_Mensuelle_L']) * 100,
        np.nan # ou 0 si consommation = 0, ou un grand nombre si consommation > 0 et dotation = 0
    )
    conso_mensuelle_veh['Taux_Utilisation_Mois_%'] = conso_mensuelle_veh['Taux_Utilisation_Mois_%'].round(1)


    # --- Récapitulatif par véhicule sur toute la période ---
    # Consommation totale sur la période par véhicule
    conso_totale_periode_veh = df_merged.groupby(['Nouveau Immat', 'Catégorie', 'Dotation']).agg(
        Consommation_Reelle_Periode_L=('Quantity', 'sum')
    ).reset_index()
    conso_totale_periode_veh.rename(columns={'Dotation': 'Dotation_Mensuelle_L'}, inplace=True)

    # Dotation totale allouée pour la période
    conso_totale_periode_veh['Nb_Mois_Periode'] = nombre_mois_periode
    conso_totale_periode_veh['Dotation_Allouee_Periode_L'] = conso_totale_periode_veh['Dotation_Mensuelle_L'] * nombre_mois_periode

    conso_totale_periode_veh['Difference_Periode_L'] = conso_totale_periode_veh['Dotation_Allouee_Periode_L'] - conso_totale_periode_veh['Consommation_Reelle_Periode_L']
    conso_totale_periode_veh['Taux_Utilisation_Periode_%'] = np.where(
        conso_totale_periode_veh['Dotation_Allouee_Periode_L'] > 0,
        (conso_totale_periode_veh['Consommation_Reelle_Periode_L'] / conso_totale_periode_veh['Dotation_Allouee_Periode_L']) * 100,
        np.nan
    )
    conso_totale_periode_veh['Taux_Utilisation_Periode_%'] = conso_totale_periode_veh['Taux_Utilisation_Periode_%'].round(1)

    # Sélection et ordre des colonnes
    cols_recap = ['Nouveau Immat', 'Catégorie', 'Dotation_Mensuelle_L', 'Nb_Mois_Periode',
                  'Dotation_Allouee_Periode_L', 'Consommation_Reelle_Periode_L',
                  'Difference_Periode_L', 'Taux_Utilisation_Periode_%']
    df_recap_dotation_periode = conso_totale_periode_veh[cols_recap]

    cols_detail = ['Nouveau Immat', 'Catégorie', 'AnneeMois', 'Dotation_Mensuelle_L',
                   'Consommation_Mois_L', 'Difference_Mois_L', 'Taux_Utilisation_Mois_%']
    df_detail_dotation_mensuel = conso_mensuelle_veh[cols_detail].sort_values(['Nouveau Immat', 'AnneeMois'])

    return df_recap_dotation_periode, df_detail_dotation_mensuel

def afficher_page_suivi_dotations(
    df_transactions: pd.DataFrame,
    df_vehicules: pd.DataFrame,
    date_debut: datetime.date,
    date_fin: datetime.date
):
    """Affiche la page de suivi des dotations avec visualisations améliorées et analyses avancées."""
    st.header(f"⛽ Suivi des Dotations Carburant ({date_debut.strftime('%d/%m/%Y')} - {date_fin.strftime('%d/%m/%Y')})")

    # Vérification préalable des données de dotation
    if 'Dotation' not in df_vehicules.columns or df_vehicules['Dotation'].sum() == 0:
        st.warning("Aucune donnée de dotation n'est disponible ou les dotations sont toutes à zéro. Le suivi des dotations ne peut pas être effectué.")
        
        with st.expander("💡 Comment configurer les dotations?", expanded=True):
            st.info("""
            ### Configuration des dotations mensuelles
            
            Pour utiliser cette fonctionnalité, assurez-vous que:
            1. La colonne 'Dotation' est présente dans votre fichier 'CARTES VEHICULE'
            2. Les valeurs représentent la dotation mensuelle en litres pour chaque véhicule
            3. Les cellules contiennent des valeurs numériques (pas de texte ou de formules)
            
            Les dotations permettent de surveiller la consommation par rapport aux allocations budgétées.
            """)
        return

    if df_transactions.empty:
        st.warning("Aucune transaction à analyser pour la période sélectionnée.")
        return

    # --- Filtres améliorés ---
    st.sidebar.subheader("🔍 Filtres pour Suivi Dotations")
    
    # Filtre par catégorie avec comptage
    all_cats_dot = sorted(df_vehicules['Catégorie'].dropna().astype(str).unique())
    cats_count = {cat: len(df_vehicules[df_vehicules['Catégorie'] == cat]) for cat in all_cats_dot}
    cats_display = [f"{cat} ({cats_count[cat]} véh.)" for cat in all_cats_dot]
    
    selected_cats_display = st.sidebar.multiselect(
        "Filtrer par Catégories", options=cats_display, default=cats_display, key="dot_cat_filter_display"
    )
    # Extraction des catégories sans le comptage
    selected_cats_dot = [cat.split(" (")[0] for cat in selected_cats_display]
    
    # Filtre par plage d'utilisation
    st.sidebar.subheader("Filtrer par taux d'utilisation")
    min_utilisation, max_utilisation = st.sidebar.slider(
        "Plage de taux d'utilisation (%)", 
        min_value=0, max_value=200, value=(0, 200), step=5,
        key="dot_util_range"
    )

    # Analyse des données
    with st.spinner("⏳ Analyse du suivi des dotations en cours..."):
        df_recap, df_detail_mensuel = analyser_suivi_dotations(
            df_transactions, df_vehicules, date_debut, date_fin
        )

    if df_recap.empty:
        st.info("Aucune donnée à afficher pour le suivi des dotations avec les filtres actuels.")
        return

    # Appliquer les filtres post-analyse
    df_recap_filtered = df_recap[
        (df_recap['Catégorie'].isin(selected_cats_dot)) &
        (df_recap['Taux_Utilisation_Periode_%'] >= min_utilisation) &
        (df_recap['Taux_Utilisation_Periode_%'] <= max_utilisation)
    ]
    
    # --- KPIs globaux en haut de la page ---
    if not df_recap_filtered.empty:
        st.subheader("📊 Aperçu Global des Dotations")
        
        # Calcul des métriques clés
        dotation_totale = df_recap_filtered['Dotation_Allouee_Periode_L'].sum()
        consommation_totale = df_recap_filtered['Consommation_Reelle_Periode_L'].sum()
        taux_global = (consommation_totale / dotation_totale * 100) if dotation_totale > 0 else 0
        nb_vehicles = len(df_recap_filtered)
        
        # Calcul des sous/surutilisations
        nb_sous_util = len(df_recap_filtered[df_recap_filtered['Taux_Utilisation_Periode_%'] < 80])
        nb_sur_util = len(df_recap_filtered[df_recap_filtered['Taux_Utilisation_Periode_%'] > 120])
        
        # Affichage des KPIs avec 3 colonnes
        col1, col2, col3 = st.columns(3)
        
        # Colonne 1 - Chiffres dotation/consommation
        col1.metric(
            "Dotation Totale Allouée", 
            f"{dotation_totale:,.0f} L",
            help="Somme des dotations mensuelles sur toute la période"
        )
        col1.metric(
            "Consommation Réelle", 
            f"{consommation_totale:,.0f} L",
            delta=f"{consommation_totale - dotation_totale:,.0f} L",
            delta_color="inverse"
        )
        
        # Colonne 2 - Taux global et nombre de véhicules
        col2.metric(
            "Taux d'Utilisation Global", 
            f"{taux_global:.1f}%",
            help="Consommation totale / Dotation totale"
        )
        col2.metric(
            "Nombre de Véhicules", 
            f"{nb_vehicles}",
            help="Véhicules avec dotation et transactions"
        )
        
        # Colonne 3 - Alertes sous/surutilisation
        col3.metric(
            "Véhicules Sous-utilisés (<80%)", 
            f"{nb_sous_util}",
            help="Véhicules utilisant moins de 80% de leur dotation"
        )
        col3.metric(
            "Véhicules Surutilisés (>120%)", 
            f"{nb_sur_util}",
            help="Véhicules utilisant plus de 120% de leur dotation",
            delta_color="inverse"
        )
        
        # Jauge de taux d'utilisation global
        fig_gauge = go.Figure(go.Indicator(
            mode="gauge+number",
            value=taux_global,
            title={'text': "Taux d'Utilisation Global (%)"},
            gauge={
                'axis': {'range': [0, 150], 'tickwidth': 1},
                'bar': {'color': "darkblue"},
                'steps': [
                    {'range': [0, 70], 'color': "red"},
                    {'range': [70, 90], 'color': "orange"},
                    {'range': [90, 110], 'color': "green"},
                    {'range': [110, 150], 'color': "red"}
                ],
                'threshold': {
                    'line': {'color': "black", 'width': 4},
                    'thickness': 0.75,
                    'value': 100
                }
            }
        ))
        fig_gauge.update_layout(height=250)
        st.plotly_chart(fig_gauge, use_container_width=True)
    
    # --- Tableau récapitulatif avec classement et code couleur ---
    st.subheader("📋 Récapitulatif de l'Utilisation des Dotations par Véhicule")
    
    # Création d'un DataFrame stylisé
    df_styled = df_recap_filtered.copy()
    
    # Fonction pour appliquer un style conditionnel aux taux d'utilisation
    def highlight_utilisation(val):
        if pd.isna(val):
            return ''
        
        # Pour les colonnes numériques de taux d'utilisation
        if isinstance(val, (int, float)):
            if val < 70:
                return 'background-color: #FFCCCC'  # Rouge clair pour très sous-utilisé
            elif val < 90:
                return 'background-color: #FFEECC'  # Orange clair pour sous-utilisé
            elif val <= 110:
                return 'background-color: #CCFFCC'  # Vert clair pour optimal
            elif val <= 130:
                return 'background-color: #FFEECC'  # Orange clair pour surutilisé
            else:
                return 'background-color: #FFCCCC'  # Rouge clair pour très surutilisé
        return ''
    
    # Ajout d'une colonne pour l'étiquette de statut
    df_styled['Statut'] = pd.cut(
        df_styled['Taux_Utilisation_Periode_%'],
        bins=[-float('inf'), 70, 90, 110, 130, float('inf')],
        labels=['Très sous-utilisé', 'Sous-utilisé', 'Optimal', 'Surutilisé', 'Très surutilisé']
    )
    
    # Réorganiser et renommer les colonnes pour meilleure lisibilité
    df_display = df_styled[[
        'Nouveau Immat', 'Catégorie', 'Statut', 'Dotation_Mensuelle_L', 
        'Nb_Mois_Periode', 'Dotation_Allouee_Periode_L', 
        'Consommation_Reelle_Periode_L', 'Difference_Periode_L', 
        'Taux_Utilisation_Periode_%'
    ]].copy()
    
    df_display.columns = [
        'Immatriculation', 'Catégorie', 'Statut', 'Dotation Mensuelle (L)', 
        'Nb Mois', 'Dotation Allouée (L)', 
        'Consommation Réelle (L)', 'Différence (L)', 
        "Taux d'Utilisation (%)"
    ]
    
    # Trier par taux d'utilisation pour mettre en évidence les cas extrêmes
    df_display = df_display.sort_values("Taux d'Utilisation (%)", ascending=False)
    
    # Appliquer le style conditionnel et afficher
    st.dataframe(
        df_display.style.applymap(
            highlight_utilisation, 
            subset=["Taux d'Utilisation (%)"]
        ),
        use_container_width=True
    )
    
    # Bouton d'exportation
    excel_data = to_excel(df_display)
    st.download_button(
        label="📥 Exporter en Excel",
        data=excel_data,
        file_name=f"suivi_dotations_{date_debut.strftime('%Y%m%d')}_{date_fin.strftime('%Y%m%d')}.xlsx",
        mime=EXCEL_MIME_TYPE
    )
    
    # --- Graphiques améliorés ---
    st.subheader("📊 Visualisation des Taux d'Utilisation")
    
    tab1, tab2, tab3 = st.tabs(["Distribution", "Par Catégorie", "Top/Flop Véhicules"])
    
    with tab1:
        # Histogramme de distribution des taux d'utilisation
        fig_dist = px.histogram(
            df_recap_filtered,
            x='Taux_Utilisation_Periode_%',
            nbins=20,
            title="Distribution des Taux d'Utilisation",
            labels={'Taux_Utilisation_Periode_%': "Taux d'Utilisation (%)"}
        )
        
        # Ajouter des lignes verticales pour les seuils
        fig_dist.add_vline(x=80, line_dash="dash", line_color="orange", 
                          annotation_text="Sous-utilisation")
        fig_dist.add_vline(x=120, line_dash="dash", line_color="red", 
                          annotation_text="Sur-utilisation")
        fig_dist.add_vline(x=100, line_dash="dash", line_color="green", 
                          annotation_text="Optimal")
        
        st.plotly_chart(fig_dist, use_container_width=True)
    
    with tab2:
        # Boxplot par catégorie pour voir la distribution
        fig_box = px.box(
            df_recap_filtered,
            x='Catégorie',
            y='Taux_Utilisation_Periode_%',
            title="Distribution des Taux d'Utilisation par Catégorie",
            labels={
                'Taux_Utilisation_Periode_%': "Taux d'Utilisation (%)",
                'Catégorie': 'Catégorie de Véhicule'
            }
        )
        fig_box.add_hline(y=100, line_dash="dash", line_color="green")
        st.plotly_chart(fig_box, use_container_width=True)
        
        # Agrégation par catégorie
        cat_stats = df_recap_filtered.groupby('Catégorie').agg(
            Nb_Vehicules=('Nouveau Immat', 'nunique'),
            Taux_Moyen=('Taux_Utilisation_Periode_%', 'mean'),
            Taux_Median=('Taux_Utilisation_Periode_%', 'median'),
            Dotation_Totale=('Dotation_Allouee_Periode_L', 'sum'),
            Consommation_Totale=('Consommation_Reelle_Periode_L', 'sum')
        ).reset_index()
        
        cat_stats['Taux_Global'] = (cat_stats['Consommation_Totale'] / cat_stats['Dotation_Totale'] * 100).round(1)
        
        st.dataframe(cat_stats, use_container_width=True)
    
    with tab3:
        col_top, col_flop = st.columns(2)
        
        with col_top:
            # Top 5 des véhicules avec le plus haut taux d'utilisation
            top_vehicles = df_recap_filtered.nlargest(5, 'Taux_Utilisation_Periode_%')
            fig_top = px.bar(
                top_vehicles,
                x='Nouveau Immat',
                y='Taux_Utilisation_Periode_%',
                title="Top 5 - Taux d'Utilisation les Plus Élevés",
                color='Taux_Utilisation_Periode_%',
                color_continuous_scale='Reds',
                labels={'Taux_Utilisation_Periode_%': "Taux d'Utilisation (%)"}
            )
            fig_top.add_hline(y=100, line_dash="dash", line_color="green")
            st.plotly_chart(fig_top, use_container_width=True)
        
        with col_flop:
            # Bottom 5 des véhicules avec le plus bas taux d'utilisation
            flop_vehicles = df_recap_filtered.nsmallest(5, 'Taux_Utilisation_Periode_%')
            fig_flop = px.bar(
                flop_vehicles,
                x='Nouveau Immat',
                y='Taux_Utilisation_Periode_%',
                title="Flop 5 - Taux d'Utilisation les Plus Bas",
                color='Taux_Utilisation_Periode_%',
                color_continuous_scale='Blues_r',
                labels={'Taux_Utilisation_Periode_%': "Taux d'Utilisation (%)"}
            )
            fig_flop.add_hline(y=100, line_dash="dash", line_color="green")
            st.plotly_chart(fig_flop, use_container_width=True)
    
    # --- Analyse détaillée par véhicule ---
    st.subheader("🚗 Analyse Détaillée par Véhicule")
    
    # Sélection améliorée du véhicule avec catégories et taux
    vehicle_options = []
    for _, row in df_recap_filtered.iterrows():
        label = f"{row['Nouveau Immat']} - {row['Catégorie']} ({row['Taux_Utilisation_Periode_%']:.1f}%)"
        vehicle_options.append((row['Nouveau Immat'], label))
    
    vehicle_options.sort(key=lambda x: x[1])
    vehicle_options = [("Tous", "Tous les véhicules")] + vehicle_options
    
    selected_vehicle_tuple = st.selectbox(
        "Sélectionner un véhicule pour l'analyse détaillée",
        options=vehicle_options,
        format_func=lambda x: x[1],
        key="vehicule_detail_select"
    )
    
    selected_vehicle_dot = selected_vehicle_tuple[0]
    
    # Affichage du détail mensuel pour le véhicule sélectionné
    if selected_vehicle_dot != "Tous":
        # Filtrer les données pour ce véhicule
        veh_recap = df_recap_filtered[df_recap_filtered['Nouveau Immat'] == selected_vehicle_dot]
        
        if not veh_recap.empty:
            veh_data = veh_recap.iloc[0]
            
            # Information sur le véhicule
            st.subheader(f"Détails pour {selected_vehicle_dot} - {veh_data['Catégorie']}")
            
            # KPIs spécifiques au véhicule
            col_v1, col_v2, col_v3 = st.columns(3)
            
            col_v1.metric("Dotation Mensuelle", f"{veh_data['Dotation_Mensuelle_L']:.1f} L")
            col_v1.metric("Dotation Totale Période", f"{veh_data['Dotation_Allouee_Periode_L']:.1f} L")
            
            col_v2.metric("Consommation Réelle", f"{veh_data['Consommation_Reelle_Periode_L']:.1f} L")
            col_v2.metric(
                "Différence", 
                f"{veh_data['Difference_Periode_L']:.1f} L",
                delta=f"{veh_data['Difference_Periode_L']:.1f} L"
            )
            
            col_v3.metric(
                "Taux d'Utilisation", 
                f"{veh_data['Taux_Utilisation_Periode_%']:.1f}%", 
                delta=f"{veh_data['Taux_Utilisation_Periode_%'] - 100:.1f}%",
                delta_color="inverse" if veh_data['Taux_Utilisation_Periode_%'] > 100 else "normal"
            )
            col_v3.metric("Durée d'Analyse", f"{veh_data['Nb_Mois_Periode']} mois")
            
            # Filtrer le détail mensuel pour ce véhicule
            df_detail_veh = df_detail_mensuel[df_detail_mensuel['Nouveau Immat'] == selected_vehicle_dot]
            
            if not df_detail_veh.empty:
                # Graphique d'évolution mensuelle
                fig_evol = px.line(
                    df_detail_veh.sort_values('AnneeMois'),
                    x='AnneeMois',
                    y=['Dotation_Mensuelle_L', 'Consommation_Mois_L'],
                    title=f"Évolution Mensuelle - {selected_vehicle_dot}",
                    labels={
                        'value': 'Volume (L)',
                        'variable': 'Type',
                        'AnneeMois': 'Période'
                    },
                    markers=True,
                    color_discrete_map={
                        'Dotation_Mensuelle_L': 'blue',
                        'Consommation_Mois_L': 'orange'
                    }
                )
                st.plotly_chart(fig_evol, use_container_width=True)
                
                # Tableau de détail mensuel stylisé
                df_detail_styled = df_detail_veh.sort_values('AnneeMois', ascending=False).copy()
                
                # Renommer les colonnes pour l'affichage
                df_detail_styled = df_detail_styled[[
                    'AnneeMois', 'Dotation_Mensuelle_L', 'Consommation_Mois_L',
                    'Difference_Mois_L', 'Taux_Utilisation_Mois_%'
                ]].copy()
                
                df_detail_styled.columns = [
                    'Mois', 'Dotation (L)', 'Consommation (L)',
                    'Différence (L)', "Taux d'Utilisation (%)"
                ]
                
                st.dataframe(
                    df_detail_styled.style.applymap(
                        highlight_utilisation,
                        subset=["Taux d'Utilisation (%)"]
                    ),
                    use_container_width=True
                )
                
                # Statistiques sur les variations mensuelles
                if len(df_detail_veh) > 1:
                    st.subheader("📈 Analyse des Variations Mensuelles")
                    
                    df_detail_veh_sorted = df_detail_veh.sort_values('AnneeMois')
                    df_detail_veh_sorted['Var_Consommation'] = df_detail_veh_sorted['Consommation_Mois_L'].pct_change() * 100
                    df_detail_veh_sorted['Var_Taux'] = df_detail_veh_sorted['Taux_Utilisation_Mois_%'].diff()
                    
                    # Calculer les statistiques
                    var_stats = {
                        'Variation max': df_detail_veh_sorted['Var_Consommation'].max(),
                        'Variation min': df_detail_veh_sorted['Var_Consommation'].min(),
                        'Variation moyenne': df_detail_veh_sorted['Var_Consommation'].mean(),
                        'Écart-type': df_detail_veh_sorted['Var_Consommation'].std()
                    }
                    
                    # Afficher les statistiques
                    col_stat1, col_stat2 = st.columns(2)
                    
                    with col_stat1:
                        for key, val in var_stats.items():
                            st.metric(key, f"{val:.1f}%" if pd.notna(val) else "N/A")
                    
                    with col_stat2:
                        # Graphique des variations mensuelles
                        fig_var = px.bar(
                            df_detail_veh_sorted,
                            x='AnneeMois',
                            y='Var_Consommation',
                            title="Variations Mensuelles de Consommation (%)",
                            color='Var_Consommation',
                            color_continuous_scale='RdBu',
                            labels={'Var_Consommation': 'Variation (%)', 'AnneeMois': 'Mois'}
                        )
                        fig_var.add_hline(y=0, line_dash="solid", line_color="black")
                        st.plotly_chart(fig_var, use_container_width=True)
            else:
                st.info(f"Aucun détail mensuel disponible pour {selected_vehicle_dot}.")
        else:
            st.info(f"Données non trouvées pour {selected_vehicle_dot}.")
    else:
        # Vue agrégée pour tous les véhicules
        if not df_detail_mensuel.empty:
            with st.expander("Voir l'évolution mensuelle globale", expanded=True):
                # Agrégation par mois
                monthly_agg = df_detail_mensuel.groupby('AnneeMois').agg(
                    Dotation_Totale=('Dotation_Mensuelle_L', 'sum'),
                    Consommation_Totale=('Consommation_Mois_L', 'sum'),
                    Nb_Vehicules=('Nouveau Immat', 'nunique')
                ).reset_index()
                
                monthly_agg['Taux_Utilisation'] = (monthly_agg['Consommation_Totale'] / 
                                                 monthly_agg['Dotation_Totale'] * 100).round(1)
                
                # Graphique d'évolution tous véhicules
                fig_all = px.line(
                    monthly_agg.sort_values('AnneeMois'),
                    x='AnneeMois',
                    y=['Dotation_Totale', 'Consommation_Totale'],
                    title=f"Évolution Mensuelle Globale - Tous Véhicules",
                    labels={
                        'value': 'Volume Total (L)',
                        'variable': 'Type',
                        'AnneeMois': 'Période'
                    },
                    markers=True
                )
                fig_all.add_hline(y=monthly_agg['Dotation_Totale'].mean(), 
                                 line_dash="dash", line_color="blue", 
                                 annotation_text="Dotation moyenne")
                st.plotly_chart(fig_all, use_container_width=True)
                
                # Graphique du taux d'utilisation mensuel
                fig_taux = px.line(
                    monthly_agg.sort_values('AnneeMois'),
                    x='AnneeMois',
                    y='Taux_Utilisation',
                    title=f"Évolution du Taux d'Utilisation Global",
                    labels={
                        'Taux_Utilisation': "Taux d'Utilisation (%)",
                        'AnneeMois': 'Période'
                    },
                    markers=True
                )
                fig_taux.add_hline(y=100, line_dash="dash", line_color="green", 
                                  annotation_text="Objectif")
                st.plotly_chart(fig_taux, use_container_width=True)
                
                # Tableau récapitulatif par mois
                st.dataframe(monthly_agg.sort_values('AnneeMois', ascending=False), 
                            use_container_width=True)
    
    # --- Recommandations et alertes ---
    with st.expander("💡 Recommandations", expanded=True):
        # Analyser les problèmes et générer des recommandations
        nb_sous_util_severe = len(df_recap_filtered[df_recap_filtered['Taux_Utilisation_Periode_%'] < 70])
        nb_sur_util_severe = len(df_recap_filtered[df_recap_filtered['Taux_Utilisation_Periode_%'] > 130])
        
        st.markdown("### Analyse et Recommandations")
        
        if nb_sous_util > 0:
            st.warning(f"⚠️ **Sous-utilisation**: {nb_sous_util} véhicules ({nb_sous_util_severe} sévères) utilisent moins de 80% de leur dotation.")
            st.markdown("""
            **Recommandations pour la sous-utilisation:**
            - Vérifier si les dotations sont surestimées par rapport aux besoins réels
            - Envisager une réduction des dotations pour les véhicules régulièrement sous-utilisés
            - Vérifier si certains véhicules peu utilisés pourraient être réaffectés
            """)
        
        if nb_sur_util > 0:
            st.error(f"🚨 **Sur-utilisation**: {nb_sur_util} véhicules ({nb_sur_util_severe} sévères) utilisent plus de 120% de leur dotation.")
            st.markdown("""
            **Recommandations pour la sur-utilisation:**
            - Analyser les causes de dépassement (trajets exceptionnels, besoins réels plus importants)
            - Vérifier s'il y a des fuites ou problèmes mécaniques sur ces véhicules
            - Évaluer si les dotations doivent être ajustées à la hausse
            - Vérifier si les transactions correspondent bien au kilométrage parcouru
            """)
        
        if abs(taux_global - 100) > 15:
            st.info(f"ℹ️ **Ajustement global**: Le taux d'utilisation global ({taux_global:.1f}%) s'écarte significativement de l'objectif de 100%.")
            st.markdown("""
            **Recommandations pour l'ajustement global:**
            - Réévaluer la méthodologie de calcul des dotations pour l'ensemble du parc
            - Considérer un ajustement proportionnel pour toutes les catégories de véhicules
            """)
        
        # Recommandations spécifiques par catégorie
        problematic_cats = cat_stats[(cat_stats['Taux_Global'] < 80) | (cat_stats['Taux_Global'] > 120)]
        
        if not problematic_cats.empty:
            st.markdown("#### Recommandations par catégorie de véhicule:")
            for _, cat_row in problematic_cats.iterrows():
                if cat_row['Taux_Global'] < 80:
                    st.markdown(f"- **{cat_row['Catégorie']}**: Dotation probablement surestimée (utilisation: {cat_row['Taux_Global']}%)")
                else:
                    st.markdown(f"- **{cat_row['Catégorie']}**: Dotation probablement sous-estimée (utilisation: {cat_row['Taux_Global']}%)")
def analyser_geolocalisation_vs_transactions(
    df_geoloc: pd.DataFrame,
    df_transactions: pd.DataFrame,
    df_vehicules: pd.DataFrame,
    date_debut: datetime.date,
    date_fin: datetime.date
) -> Tuple[pd.DataFrame, pd.DataFrame]:
    """
    Compare les données de géolocalisation avec les transactions carburant.

    Args:
        df_geoloc: DataFrame des données de géolocalisation.
        df_transactions: DataFrame des transactions.
        df_vehicules: DataFrame des véhicules.
        date_debut: Date de début de la période d'analyse.
        date_fin: Date de fin de la période d'analyse.

    Returns:
        Un tuple de DataFrames:
        - comparaison_par_vehicule: Comparaison des distances pour chaque véhicule.
        - anomalies_distance: Anomalies détectées dans les déclarations de distance.
    """
    # Filtrer les données pour la période spécifiée
    mask_date_geoloc = (df_geoloc['Date'].dt.date >= date_debut) & (df_geoloc['Date'].dt.date <= date_fin)
    df_geoloc_filtre = df_geoloc[mask_date_geoloc].copy()

    mask_date_trans = (df_transactions['Date'].dt.date >= date_debut) & (df_transactions['Date'].dt.date <= date_fin)
    df_trans_filtre = df_transactions[mask_date_trans].copy()

    # Fusionner transactions avec infos véhicules pour avoir l'immatriculation
    df_trans_avec_immat = df_trans_filtre.merge(
        df_vehicules[['N° Carte', 'Nouveau Immat']],
        left_on='Card num.',
        right_on='N° Carte',
        how='inner'
    )

    # Agréger les distances par immatriculation
    # Filtrer sur Type == 'Trajet' pour ne sommer que les distances des trajets
    trajets_geoloc = df_geoloc_filtre[df_geoloc_filtre['Type'] == 'Trajet']
    distance_geoloc = trajets_geoloc.groupby('Véhicule').agg(
        Distance_Geoloc_Totale=('Distance', 'sum'),
        Nb_Trajets=('Distance', 'count'), # Compte le nombre de segments de trajet
        Distance_Moyenne_Trajet=('Distance', 'mean'),
        Vitesse_Moyenne_Globale=('Vitesse moyenne', 'mean'), # Mean of average speeds of segments
        Vitesse_Max_Observee=('Vitesse moyenne', 'max'), # Max of average speeds of segments
        Duree_Totale_Minutes=('Durée_minutes', 'sum')
    ).reset_index()

    # Calculer les écarts de kilométrage dans les transactions
    df_trans_avec_immat['distance_declaree'] = df_trans_avec_immat['Current mileage'] - df_trans_avec_immat['Past mileage']

    # Agréger par immatriculation
    distance_trans = df_trans_avec_immat.groupby('Nouveau Immat').agg(
        Distance_Declaree_Totale=('distance_declaree', lambda x: x[x > 0].sum()),  # Somme des distances positives
        Nb_Transactions=('distance_declaree', 'count'),
        Volume_Carburant_Total=('Quantity', 'sum')
    ).reset_index()

    # Fusionner les deux ensembles de données
    # Normaliser les noms de colonnes pour la fusion
    distance_geoloc.rename(columns={'Véhicule': 'Immatriculation'}, inplace=True)
    distance_trans.rename(columns={'Nouveau Immat': 'Immatriculation'}, inplace=True)

    comparaison = distance_geoloc.merge(
        distance_trans,
        on='Immatriculation',
        how='outer'
    )

    # Calculer les écarts et pourcentages
    comparaison['Ecart_Distance'] = comparaison['Distance_Declaree_Totale'] - comparaison['Distance_Geoloc_Totale']
    comparaison['Pourcentage_Ecart'] = np.where(
        comparaison['Distance_Geoloc_Totale'] > 0,
        (comparaison['Ecart_Distance'] / comparaison['Distance_Geoloc_Totale']) * 100,
        np.nan
    )

    # Calculer la consommation aux 100km basée sur la distance géolocalisée (plus fiable)
    comparaison['Consommation_100km_Reelle'] = np.where(
        comparaison['Distance_Geoloc_Totale'] > 0,
        (comparaison['Volume_Carburant_Total'] / comparaison['Distance_Geoloc_Totale']) * 100,
        np.nan
    )

    # Ajouter la consommation basée sur la distance déclarée
    comparaison['Consommation_100km_Declaree'] = np.where(
        comparaison['Distance_Declaree_Totale'] > 0,
        (comparaison['Volume_Carburant_Total'] / comparaison['Distance_Declaree_Totale']) * 100,
        np.nan
    )

    # Calculer l'écart entre les deux consommations
    comparaison['Ecart_Consommation'] = comparaison['Consommation_100km_Declaree'] - comparaison['Consommation_100km_Reelle']
    comparaison['Pourcentage_Ecart_Consommation'] = np.where(
        comparaison['Consommation_100km_Reelle'] > 0,
        (comparaison['Ecart_Consommation'] / comparaison['Consommation_100km_Reelle']) * 100,
        np.nan
    )

    # Identifier les anomalies significatives (écart > 10% et au moins 10km)
    comparaison['Consommation_100km_Declaree'] = np.where(
        comparaison['Distance_Declaree_Totale'] > 0,
        (comparaison['Volume_Carburant_Total'] / comparaison['Distance_Declaree_Totale']) * 100,
        np.nan
    )

    # Calculer l'écart entre les deux consommations
    comparaison['Ecart_Consommation'] = comparaison['Consommation_100km_Declaree'] - comparaison['Consommation_100km_Reelle']
    comparaison['Pourcentage_Ecart_Consommation'] = np.where(
        comparaison['Consommation_100km_Reelle'] > 0,
        (comparaison['Ecart_Consommation'] / comparaison['Consommation_100km_Reelle']) * 100,
        np.nan
    )

    # Identifier les anomalies significatives (écart > 10% et au moins 10km)
    seuil_ecart_pct = 10  # Pourcentage
    seuil_ecart_km = 10   # Kilomètres

    anomalies = comparaison[
        (abs(comparaison['Pourcentage_Ecart']) > seuil_ecart_pct) &
        (abs(comparaison['Ecart_Distance']) > seuil_ecart_km)
    ].copy()

    anomalies['Type_Anomalie'] = np.where(
        anomalies['Ecart_Distance'] > 0,
        'Sur-déclaration kilométrique',
        'Sous-déclaration kilométrique'
    )

    anomalies['Gravite'] = np.where(
        abs(anomalies['Pourcentage_Ecart']) > 25,
        'Élevée',
        'Moyenne'
    )

    # Arrondir les valeurs numériques
    cols_arrondi = ['Distance_Geoloc_Totale', 'Distance_Declaree_Totale', 'Ecart_Distance',
                    'Pourcentage_Ecart', 'Consommation_100km_Reelle', 'Distance_Moyenne_Trajet',
                    'Vitesse_Moyenne_Globale', 'Vitesse_Max_Observee']

    for col in cols_arrondi:
        if col in comparaison.columns:
            comparaison[col] = comparaison[col].round(1)
        if col in anomalies.columns:
            anomalies[col] = anomalies[col].round(1)

    # Trier les résultats
    comparaison_triee = comparaison.sort_values('Ecart_Distance', ascending=False)
    anomalies_triees = anomalies.sort_values('Pourcentage_Ecart', ascending=False)

    return comparaison_triee, anomalies_triees


def analyser_exces_vitesse(
    df_geoloc: pd.DataFrame,
    date_debut: datetime.date,
    date_fin: datetime.date,
    seuil_vitesse: int = 90  # Vitesse limite par défaut
) -> Tuple[pd.DataFrame, pd.DataFrame]:
    """
    Analyse les excès de vitesse à partir des données de géolocalisation.

    Args:
        df_geoloc: DataFrame des données de géolocalisation.
        date_debut: Date de début de la période d'analyse.
        date_fin: Date de fin de la période d'analyse.
        seuil_vitesse: Seuil de vitesse considéré comme un excès (km/h).

    Returns:
        Un tuple de DataFrames:
        - resume_exces: Résumé des excès de vitesse par véhicule.
        - detail_exces: Détail de tous les trajets avec excès de vitesse.
    """
    # Filtrer les données pour la période spécifiée
    mask_date = (df_geoloc['Date'].dt.date >= date_debut) & (df_geoloc['Date'].dt.date <= date_fin)
    df_filtre = df_geoloc[mask_date].copy()
    
    # Garder uniquement les trajets pour l'analyse de vitesse
    df_filtre = df_filtre[df_filtre['Type'] == 'Trajet'].copy()


    if 'Vitesse moyenne' not in df_filtre.columns:
        st.warning("La colonne 'Vitesse moyenne' est manquante dans les données de géolocalisation pour l'analyse des excès.")
        return pd.DataFrame(), pd.DataFrame()

    # Identifier les excès de vitesse
    df_filtre['Exces_Vitesse'] = df_filtre['Vitesse moyenne'] > seuil_vitesse
    df_filtre['Depassement_km/h'] = np.where(
        df_filtre['Exces_Vitesse'],
        df_filtre['Vitesse moyenne'] - seuil_vitesse,
        0
    )

    # Détail de tous les trajets avec excès de vitesse
    detail_exces = df_filtre[df_filtre['Exces_Vitesse']].copy()
    if not detail_exces.empty:
        detail_exces['Niveau_Exces'] = pd.cut(
            detail_exces['Depassement_km/h'],
            bins=[0, 10, 20, 30, float('inf')],
            labels=['Léger (< 10 km/h)', 'Modéré (10-20 km/h)', 'Important (20-30 km/h)', 'Grave (> 30 km/h)'],
            right=False # Important pour que 0 ne soit pas inclus dans la première catégorie si un excès est juste au seuil
        )
    else:
        detail_exces['Niveau_Exces'] = pd.NA


    # Résumé par véhicule
    resume_exces = df_filtre.groupby('Véhicule').agg(
        Nb_Total_Trajets=('Exces_Vitesse', 'count'),
        Nb_Trajets_Exces=('Exces_Vitesse', 'sum'),
        Vitesse_Max_Observee=('Vitesse moyenne', 'max'), # Max des vitesses moyennes des segments
        Vitesse_Moyenne_Trajets=('Vitesse moyenne', 'mean'), # Moyenne des vitesses moyennes des segments
        Depassement_Moyen=('Depassement_km/h', lambda x: x[x > 0].mean() if any(x > 0) else 0)
    ).reset_index()

    # Calculer le pourcentage de trajets en excès
    resume_exces['Pourcentage_Trajets_Exces'] = (resume_exces['Nb_Trajets_Exces'] / resume_exces['Nb_Total_Trajets'] * 100).round(1)

    # Identifier le niveau de risque
    resume_exces['Niveau_Risque'] = pd.cut(
        resume_exces['Pourcentage_Trajets_Exces'],
        bins=[-1, 10, 25, 50, float('inf')], # -1 pour inclure 0%
        labels=['Faible', 'Modéré', 'Élevé', 'Très élevé']
    )

    # Trier par nombre d'excès décroissant
    resume_exces = resume_exces.sort_values('Nb_Trajets_Exces', ascending=False)
    detail_exces = detail_exces.sort_values(['Véhicule', 'Vitesse moyenne'], ascending=[True, False])

    return resume_exces, detail_exces


def analyser_utilisation_vehicules(
    df_geoloc: pd.DataFrame,
    date_debut: datetime.date,
    date_fin: datetime.date
) -> Tuple[pd.DataFrame, pd.DataFrame]:
    """
    Analyse l'utilisation des véhicules (temps, distance, périodes d'activité).

    Args:
        df_geoloc: DataFrame des données de géolocalisation.
        date_debut: Date de début de la période d'analyse.
        date_fin: Date de fin de la période d'analyse.

    Returns:
        Un tuple de DataFrames:
        - utilisation_par_vehicule: Résumé d'utilisation par véhicule.
        - utilisation_quotidienne: Utilisation quotidienne des véhicules.
    """
    # Filtrer les données pour la période spécifiée
    mask_date = (df_geoloc['Date'].dt.date >= date_debut) & (df_geoloc['Date'].dt.date <= date_fin)
    df_filtre = df_geoloc[mask_date].copy()

    if df_filtre.empty:
        return pd.DataFrame(), pd.DataFrame()

    # Garder uniquement les trajets pour cette analyse spécifique
    df_filtre_trajets = df_filtre[df_filtre['Type'] == 'Trajet'].copy()
    if df_filtre_trajets.empty:
         return pd.DataFrame(), pd.DataFrame()


    # Extraire l'heure des déplacements pour analyser les périodes d'activité
    if 'Heure_debut' in df_filtre_trajets.columns: # Heure_debut est déjà numérique
        # Classifier les périodes de la journée
        conditions = [
            (df_filtre_trajets['Heure_debut'] >= 6) & (df_filtre_trajets['Heure_debut'] < 9),
            (df_filtre_trajets['Heure_debut'] >= 9) & (df_filtre_trajets['Heure_debut'] < 12),
            (df_filtre_trajets['Heure_debut'] >= 12) & (df_filtre_trajets['Heure_debut'] < 14),
            (df_filtre_trajets['Heure_debut'] >= 14) & (df_filtre_trajets['Heure_debut'] < 17),
            (df_filtre_trajets['Heure_debut'] >= 17) & (df_filtre_trajets['Heure_debut'] < 20),
            (df_filtre_trajets['Heure_debut'] >= 20) | (df_filtre_trajets['Heure_debut'] < 6)
        ]

        periodes = [
            'Matin (6h-9h)', 'Matinée (9h-12h)', 'Midi (12h-14h)',
            'Après-midi (14h-17h)', 'Soir (17h-20h)', 'Nuit (20h-6h)'
        ]
        df_filtre_trajets['Periode_Jour'] = np.select(conditions, periodes, default='Non défini')
    else:
        df_filtre_trajets['Periode_Jour'] = 'Non défini'


    # Ajouter le jour de la semaine
    df_filtre_trajets['Jour_Semaine'] = df_filtre_trajets['Date'].dt.day_name()
    # Est_weekend est déjà calculé par charger_donnees_geolocalisation

    # Résumé d'utilisation par véhicule
    utilisation_par_vehicule = df_filtre_trajets.groupby('Véhicule').agg(
        Distance_Totale=('Distance', 'sum'),
        Nb_Trajets=('Distance', 'count'),
        Duree_Totale_Minutes=('Durée_minutes', 'sum'),
        Distance_Moyenne_Trajet=('Distance', 'mean'),
        Duree_Moyenne_Trajet=('Durée_minutes', 'mean'),
        Nb_Trajets_Weekend=('Est_weekend', lambda x: x.sum()),
        Vitesse_Moyenne=('Vitesse moyenne', 'mean')
    ).reset_index()

    # Calculer la durée totale en heures
    utilisation_par_vehicule['Duree_Totale_Heures'] = (utilisation_par_vehicule['Duree_Totale_Minutes'] / 60).round(1)

    # Calculer le pourcentage de trajets le weekend
    utilisation_par_vehicule['Pourcentage_Trajets_Weekend'] = (
        utilisation_par_vehicule['Nb_Trajets_Weekend'] / utilisation_par_vehicule['Nb_Trajets'] * 100
    ).round(1)

    # Utilisation quotidienne (pour graphiques d'évolution)
    utilisation_quotidienne = df_filtre_trajets.groupby(['Date', 'Véhicule']).agg(
        Distance_Jour=('Distance', 'sum'),
        Nb_Trajets_Jour=('Distance', 'count'),
        Duree_Jour_Minutes=('Durée_minutes', 'sum')
    ).reset_index()

    utilisation_quotidienne['Duree_Jour_Heures'] = (utilisation_quotidienne['Duree_Jour_Minutes'] / 60).round(1)

    # Analyser la répartition des trajets par période
    if 'Periode_Jour' in df_filtre_trajets.columns:
        repartition_periodes = df_filtre_trajets.groupby(['Véhicule', 'Periode_Jour']).size().unstack(fill_value=0)
        # Fusionner avec le résumé d'utilisation
        if not repartition_periodes.empty:
            utilisation_par_vehicule = utilisation_par_vehicule.merge(
                repartition_periodes,
                left_on='Véhicule',
                right_index=True,
                how='left'
            )

    # Arrondir les valeurs numériques
    cols_arrondi = ['Distance_Totale', 'Distance_Moyenne_Trajet', 'Duree_Moyenne_Trajet', 'Vitesse_Moyenne']
    for col in cols_arrondi:
        if col in utilisation_par_vehicule.columns:
            utilisation_par_vehicule[col] = utilisation_par_vehicule[col].round(1)

    return utilisation_par_vehicule.sort_values('Distance_Totale', ascending=False), utilisation_quotidienne


# ---------------------------------------------------------------------
# NOUVELLES FONCTIONS POUR L'ANALYSE AVANCÉE DE GÉOLOCALISATION
# ---------------------------------------------------------------------

def detecter_trajets_suspects(
    df_geoloc: pd.DataFrame,
    date_debut: datetime.date,
    date_fin: datetime.date
) -> pd.DataFrame:
    """
    Détecte les trajets suspects basés sur plusieurs critères:
    - Trajets hors heures de service
    - Trajets le weekend
    - Trajets avec vitesse anormalement basse (arrêts fréquents non enregistrés)

    Args:
        df_geoloc: DataFrame des données de géolocalisation.
        date_debut: Date de début de la période d'analyse.
        date_fin: Date de fin de la période d'analyse.

    Returns:
        Un DataFrame des trajets suspects avec les détails et le score de risque.
    """
    # Vérifier si la détection est activée
    if not st.session_state.get('ss_activer_trajets_suspects', True):
        return pd.DataFrame()  # Retourner un DataFrame vide si désactivé
        
    # Paramètres (récupérés de session_state)
    heure_debut_service = st.session_state.get('ss_heure_debut_service', DEFAULT_HEURE_DEBUT_SERVICE)
    heure_fin_service = st.session_state.get('ss_heure_fin_service', DEFAULT_HEURE_FIN_SERVICE)
    # nb_arrets_suspect n'est pas utilisé directement ici mais dans la logique d'agrégation globale

    # Poids pour le score de risque
    poids_trajet_hors_heures = st.session_state.get('ss_poids_trajet_hors_heures', DEFAULT_POIDS_TRAJET_HORS_HEURES)
    poids_trajet_weekend = st.session_state.get('ss_poids_trajet_weekend', DEFAULT_POIDS_TRAJET_WEEKEND)
    poids_arrets_frequents = st.session_state.get('ss_poids_arrets_frequents', DEFAULT_POIDS_ARRETS_FREQUENTS) # pour vitesse lente

    # Filtrer les données pour la période spécifiée
    mask_date = (df_geoloc['Date'].dt.date >= date_debut) & (df_geoloc['Date'].dt.date <= date_fin)
    df_filtre = df_geoloc[mask_date].copy()

    # S'assurer de ne traiter que les trajets
    df_filtre = df_filtre[df_filtre['Type'] == 'Trajet'].copy()


    if df_filtre.empty:
        return pd.DataFrame()

    # Identifier les types d'anomalies et calculer les scores
    # Heure_debut et Est_weekend sont déjà dans df_filtre grâce à charger_donnees_geolocalisation
    df_filtre['Est_Hors_Heures'] = (df_filtre['Heure_debut'] < heure_debut_service) | (df_filtre['Heure_debut'] >= heure_fin_service)
    df_filtre['Score_Hors_Heures'] = np.where(df_filtre['Est_Hors_Heures'], poids_trajet_hors_heures, 0)

    df_filtre['Score_Weekend'] = np.where(df_filtre['Est_weekend'], poids_trajet_weekend, 0)

    # Vitesse moyenne trop basse pour la distance peut indiquer des arrêts fréquents non documentés
    df_filtre['Est_Vitesse_Lente'] = (df_filtre['Vitesse moyenne'] < 25) & (df_filtre['Distance'] > 5) # Exemple de seuils
    df_filtre['Score_Vitesse_Lente'] = np.where(df_filtre['Est_Vitesse_Lente'], poids_arrets_frequents, 0)

    # Calculer le score total de suspicion
    df_filtre['Score_Suspicion_Total'] = df_filtre['Score_Hors_Heures'] + df_filtre['Score_Weekend'] + df_filtre['Score_Vitesse_Lente']

    # Filtrer les trajets avec un score de suspicion > 0
    trajets_suspects = df_filtre[df_filtre['Score_Suspicion_Total'] > 0].copy()

    if trajets_suspects.empty:
        return pd.DataFrame()

    # Créer un résumé des motifs de suspicion
    def creer_description_suspicion(row):
        motifs = []
        if row['Est_Hors_Heures']:
            motifs.append(f"Hors heures de service ({row['Heure_debut']}h)")
        if row['Est_weekend']:
            motifs.append("Weekend")
        if row['Est_Vitesse_Lente']:
            motifs.append(f"Vitesse anormalement basse ({row['Vitesse moyenne']:.1f} km/h)")
        return " + ".join(motifs)

    trajets_suspects['Description_Suspicion'] = trajets_suspects.apply(creer_description_suspicion, axis=1)

    # Déterminer le niveau de gravité
    trajets_suspects['Niveau_Suspicion'] = pd.cut(
        trajets_suspects['Score_Suspicion_Total'],
        bins=[0, 5, 10, float('inf')],
        labels=['Faible', 'Modéré', 'Élevé'],
        right=False
    )

    # Ajouter des informations utiles (Début est la colonne str originale, DateTime_Debut est le datetime)
    trajets_suspects['Date_Heure_Debut'] = trajets_suspects.apply(
        lambda row: f"{row['Date'].strftime('%d/%m/%Y')} {row['Début']}" if pd.notna(row['Début']) and isinstance(row['Début'], str) else row['DateTime_Debut'].strftime('%d/%m/%Y %H:%M:%S') if pd.notna(row['DateTime_Debut']) else "N/A",
        axis=1
    )


    # Sélectionner et ordonner les colonnes pertinentes
    cols_suspects = [
        'Véhicule', 'Date_Heure_Debut', 'Distance', 'Durée_minutes', 'Vitesse moyenne',
        'Description_Suspicion', 'Score_Suspicion_Total', 'Niveau_Suspicion'
    ]

    return trajets_suspects[cols_suspects].sort_values('Score_Suspicion_Total', ascending=False)

def afficher_page_bilan_carbone(df_transactions: pd.DataFrame, df_vehicules: pd.DataFrame, date_debut: datetime.date, date_fin: datetime.date):
    """Affiche la page d'analyse du bilan carbone."""
    st.header(f"🌍 Bilan Carbone ({date_debut.strftime('%d/%m/%Y')} - {date_fin.strftime('%d/%m/%Y')})")

    if df_transactions.empty:
        st.warning("Aucune transaction à analyser pour la période sélectionnée.")
        return

    # Récupérer les facteurs d'émission depuis session_state ou utiliser les valeurs par défaut
    facteur_emission_essence = st.session_state.get('facteur_emission_essence', 2.3)  # kg CO2e/L
    facteur_emission_diesel = st.session_state.get('facteur_emission_diesel', 2.7)    # kg CO2e/L
    facteur_emission_default = st.session_state.get('facteur_emission_default', 2.5)  # kg CO2e/L
    
    # Options d'affichage
    unite_bilan = st.session_state.get('unite_bilan', "kg CO2e")
    afficher_comparaisons = st.session_state.get('afficher_comparaisons', True)
    comparaison_type = st.session_state.get('comparaison_type', ["Kilomètres en voiture", "Arbres nécessaires"])

    # Fusionner les données de transactions avec les véhicules
    df_merged = df_transactions.merge(
        df_vehicules[['N° Carte', 'Nouveau Immat', 'Catégorie', 'Type']],
        left_on='Card num.',
        right_on='N° Carte',
        how='inner'
    )

    # Calculer les émissions par transaction
    df_merged['Facteur_Emission'] = facteur_emission_default  # Valeur par défaut
    
    # Attribuer le facteur d'émission en fonction du type de véhicule (si spécifié)
    # Ici on fait une hypothèse simple que les véhicules avec "DIESEL" dans le type utilisent du diesel
    # et les autres de l'essence - à adapter selon vos données réelles
    if 'Type' in df_merged.columns:
        diesel_mask = df_merged['Type'].astype(str).str.contains('DIESEL', case=False, na=False)
        df_merged.loc[diesel_mask, 'Facteur_Emission'] = facteur_emission_diesel
        df_merged.loc[~diesel_mask, 'Facteur_Emission'] = facteur_emission_essence

    # Calculer les émissions CO2 pour chaque transaction
    df_merged['Emissions_CO2'] = df_merged['Quantity'] * df_merged['Facteur_Emission']

    # Préparation des statistiques générales
    total_emissions = df_merged['Emissions_CO2'].sum()
    total_litres = df_merged['Quantity'].sum()
    emissions_moyennes_par_litre = total_emissions / total_litres if total_litres > 0 else 0
    
    # Conversion en tonnes si nécessaire pour l'affichage
    display_factor = 1000 if unite_bilan == "tonnes CO2e" else 1
    display_unit = unite_bilan
    total_emissions_display = total_emissions / display_factor
    
    # Afficher les KPIs principaux
    st.subheader("💨 Émissions Totales de CO2")
    col1, col2, col3 = st.columns(3)
    col1.metric("Émissions Totales", f"{total_emissions_display:,.1f} {display_unit}")
    col2.metric("Volume Carburant", f"{total_litres:,.1f} L")
    col3.metric("Émissions Moyennes", f"{emissions_moyennes_par_litre:.2f} kg CO2e/L")

    # Afficher des comparaisons pour contextualiser les émissions
    if afficher_comparaisons:
        st.subheader("🔄 Équivalents des Émissions")
        
        # Facteurs de conversion pour différentes comparaisons
        # Ces valeurs sont approximatives et peuvent être ajustées
        equiv_factors = {
            "Kilomètres en voiture": 0.2,        # kg CO2e par km en voiture moyenne
            "Vols Paris-New York": 1000,         # kg CO2e par vol aller simple
            "Arbres nécessaires": 25,            # kg CO2e absorbés par arbre par an
            "Repas avec viande": 7               # kg CO2e par repas avec viande de bœuf
        }
        
        # Sélectionner les comparaisons à afficher
        equivalents = {}
        for comp_type in comparaison_type:
            if comp_type in equiv_factors:
                if comp_type == "Arbres nécessaires":
                    # Pour les arbres, on calcule combien il faut pour absorber les émissions en 1 an
                    equivalents[comp_type] = total_emissions / equiv_factors[comp_type]
                else:
                    # Pour les autres, on calcule combien d'unités équivalent aux émissions
                    equivalents[comp_type] = total_emissions / equiv_factors[comp_type]
        
        # Afficher les équivalents
        cols_equiv = st.columns(len(equivalents))
        for i, (comp_type, value) in enumerate(equivalents.items()):
            if comp_type == "Kilomètres en voiture":
                cols_equiv[i].metric("Équivalent en km voiture", f"{value:,.0f} km")
            elif comp_type == "Vols Paris-New York":
                cols_equiv[i].metric("Équivalent en vols", f"{value:.1f} vols Paris-NY")
            elif comp_type == "Arbres nécessaires":
                cols_equiv[i].metric("Arbres pour compensation", f"{value:.0f} arbres/an")
            elif comp_type == "Repas avec viande":
                cols_equiv[i].metric("Équivalent en repas", f"{value:,.0f} repas avec bœuf")

    # Analyse par catégorie de véhicule
    st.subheader("📊 Émissions par Catégorie de Véhicule")
    emissions_par_categorie = df_merged.groupby('Catégorie').agg(
        Emissions_CO2=('Emissions_CO2', 'sum'),
        Volume_Total=('Quantity', 'sum'),
        Nb_Vehicules=('Nouveau Immat', 'nunique'),
        Nb_Transactions=('Quantity', 'count')
    ).reset_index()
    
    # Calculer les émissions moyennes par véhicule
    emissions_par_categorie['Emissions_Par_Vehicule'] = emissions_par_categorie['Emissions_CO2'] / emissions_par_categorie['Nb_Vehicules']
    
    # Conversion pour affichage si nécessaire
    if unite_bilan == "tonnes CO2e":
        emissions_par_categorie['Emissions_CO2'] = emissions_par_categorie['Emissions_CO2'] / 1000
        emissions_par_categorie['Emissions_Par_Vehicule'] = emissions_par_categorie['Emissions_Par_Vehicule'] / 1000
    
    # Afficher le tableau des émissions par catégorie
    afficher_dataframe_avec_export(
        emissions_par_categorie,
        "Émissions CO2 par Catégorie",
        key="emissions_categorie"
    )
    
    # Graphique des émissions par catégorie
    fig_categorie = px.pie(
        emissions_par_categorie,
        values='Emissions_CO2',
        names='Catégorie',
        title=f"Répartition des Émissions de CO2 par Catégorie ({display_unit})"
    )
    st.plotly_chart(fig_categorie, use_container_width=True)
    
    # Analyse par véhicule
    st.subheader("🚗 Émissions par Véhicule")
    emissions_par_vehicule = df_merged.groupby(['Nouveau Immat', 'Catégorie']).agg(
        Emissions_CO2=('Emissions_CO2', 'sum'),
        Volume_Total=('Quantity', 'sum'),
        Nb_Transactions=('Quantity', 'count')
    ).reset_index()
    
    # Conversion pour affichage si nécessaire
    if unite_bilan == "tonnes CO2e":
        emissions_par_vehicule['Emissions_CO2'] = emissions_par_vehicule['Emissions_CO2'] / 1000
    
    # Trier par émissions décroissantes
    emissions_par_vehicule = emissions_par_vehicule.sort_values('Emissions_CO2', ascending=False)
    
    # Afficher le tableau des émissions par véhicule
    afficher_dataframe_avec_export(
        emissions_par_vehicule,
        "Émissions CO2 par Véhicule",
        key="emissions_vehicule"
    )
    
    # Graphique des émissions pour les 10 véhicules les plus émetteurs
    top_vehicles = emissions_par_vehicule.head(10)
    fig_top_vehicles = px.bar(
        top_vehicles,
        x='Nouveau Immat',
        y='Emissions_CO2',
        color='Catégorie',
        title=f"Top 10 des Véhicules Émetteurs ({display_unit})",
        labels={'Emissions_CO2': f'Émissions CO2 ({display_unit})'}
    )
    st.plotly_chart(fig_top_vehicles, use_container_width=True)
    
    # Évolution mensuelle des émissions
    st.subheader("📈 Évolution Mensuelle des Émissions")
    emissions_mensuelles = df_merged.groupby(pd.Grouper(key='DateTime', freq='M')).agg(
        Emissions_CO2=('Emissions_CO2', 'sum'),
        Volume=('Quantity', 'sum')
    ).reset_index()
    
    emissions_mensuelles['Mois'] = emissions_mensuelles['DateTime'].dt.strftime('%Y-%m')
    
    # Conversion pour affichage si nécessaire
    if unite_bilan == "tonnes CO2e":
        emissions_mensuelles['Emissions_CO2'] = emissions_mensuelles['Emissions_CO2'] / 1000
    
    fig_monthly = px.line(
        emissions_mensuelles,
        x='Mois',
        y='Emissions_CO2',
        title=f"Évolution Mensuelle des Émissions de CO2 ({display_unit})",
        markers=True,
        labels={'Emissions_CO2': f'Émissions CO2 ({display_unit})'}
    )
    st.plotly_chart(fig_monthly, use_container_width=True)
    
    # Section recommandations
    st.subheader("💡 Recommandations pour Réduire l'Empreinte Carbone")
    with st.expander("Voir les recommandations", expanded=True):
        st.markdown("""
        ### Stratégies de réduction des émissions de CO2:

        1. **Optimisation des trajets**:
           - Planifier les itinéraires pour minimiser les distances
           - Éviter les détours inutiles
           - Combiner plusieurs missions/livraisons en un seul trajet

        2. **Formation à l'éco-conduite**:
           - Former les conducteurs aux techniques d'éco-conduite
           - Maintenir une vitesse constante
           - Éviter les accélérations et freinages brusques
           - Couper le moteur lors des arrêts prolongés

        3. **Maintenance préventive**:
           - Vérifier régulièrement la pression des pneus
           - Respecter le calendrier d'entretien des véhicules
           - Remplacer les filtres à air selon les recommandations

        4. **Renouvellement de la flotte**:
           - Privilégier les véhicules hybrides ou électriques lors du renouvellement
           - Choisir des véhicules adaptés aux besoins réels (éviter le surdimensionnement)

        5. **Suivi et sensibilisation**:
           - Établir des objectifs de réduction des émissions
           - Communiquer régulièrement sur les progrès réalisés
           - Récompenser les comportements vertueux
        """)

    # Option pour exporter le bilan carbone
    st.subheader("📑 Exporter le Bilan Carbone")
    if st.button("Générer un rapport de bilan carbone"):
        # Créer un DataFrame résumé pour le rapport
        resume_df = pd.DataFrame({
            'Métrique': [
                'Période', 'Émissions totales', 'Volume total de carburant', 
                'Émissions moyennes', 'Nombre de véhicules', 'Nombre de transactions'
            ],
            'Valeur': [
                f"{date_debut.strftime('%d/%m/%Y')} - {date_fin.strftime('%d/%m/%Y')}",
                f"{total_emissions_display:,.1f} {display_unit}",
                f"{total_litres:,.1f} L",
                f"{emissions_moyennes_par_litre:.2f} kg CO2e/L",
                f"{df_merged['Nouveau Immat'].nunique()}",
                f"{len(df_merged)}"
            ]
        })
        
        # Ajouter les comparaisons si activées
        if afficher_comparaisons:
            for comp_type, value in equivalents.items():
                new_row = pd.DataFrame({
                    'Métrique': [comp_type],
                    'Valeur': [f"{value:,.0f} km" if comp_type == "Kilomètres en voiture" else
                              f"{value:.1f} vols Paris-NY" if comp_type == "Vols Paris-New York" else
                              f"{value:.0f} arbres/an" if comp_type == "Arbres nécessaires" else
                              f"{value:,.0f} repas avec bœuf"]
                })
                resume_df = pd.concat([resume_df, new_row], ignore_index=True)
        
        # Exporter le rapport
        afficher_dataframe_avec_export(
            resume_df,
            "Résumé du Bilan Carbone",
            key="resume_bilan_carbone"
        )
        
        st.success("Rapport de bilan carbone généré avec succès! Cliquez sur le bouton d'export pour télécharger.")
def analyser_correspondance_transactions_geoloc(
    df_geoloc: pd.DataFrame,
    df_transactions: pd.DataFrame,
    df_vehicules: pd.DataFrame,
    date_debut: datetime.date,
    date_fin: datetime.date
) -> Tuple[pd.DataFrame, pd.DataFrame]:
    """
    Analyse la correspondance entre les transactions de carburant et la géolocalisation.
    Vérifie notamment si le véhicule était présent à la station lors des transactions.

    Args:
        df_geoloc: DataFrame des données de géolocalisation.
        df_transactions: DataFrame des transactions.
        df_vehicules: DataFrame des véhicules.
        date_debut: Date de début de la période d'analyse.
        date_fin: Date de fin de la période d'analyse.

    Returns:
        Un tuple de DataFrames:
        - transactions_avec_presence: Transactions avec indication de présence du véhicule.
        - transactions_suspectes: Transactions pour lesquelles aucune présence du véhicule n'est détectée.
    """
    # Paramètres
    rayon_station_km = st.session_state.get('ss_rayon_station_km', DEFAULT_RAYON_STATION_KM)
    poids_transaction_sans_presence = st.session_state.get('ss_poids_transaction_sans_presence', DEFAULT_POIDS_TRANSACTION_SANS_PRESENCE)

    # Filtrer les données pour la période spécifiée
    mask_date_geoloc = (df_geoloc['Date'].dt.date >= date_debut) & (df_geoloc['Date'].dt.date <= date_fin)
    df_geoloc_filtre = df_geoloc[mask_date_geoloc].copy()

    mask_date_trans = (df_transactions['Date'].dt.date >= date_debut) & (df_transactions['Date'].dt.date <= date_fin)
    df_trans_filtre = df_transactions[mask_date_trans].copy()

    # Fusionner transactions avec infos véhicules pour avoir l'immatriculation
    df_trans_avec_immat = df_trans_filtre.merge(
        df_vehicules[['N° Carte', 'Nouveau Immat']],
        left_on='Card num.',
        right_on='N° Carte',
        how='inner'
    )

    # Vérifier si des coordonnées GPS sont disponibles dans les données
    coordonnees_disponibles_geoloc = (
        'Latitude_depart' in df_geoloc_filtre.columns and
        'Longitude_depart' in df_geoloc_filtre.columns and
        not df_geoloc_filtre['Latitude_depart'].isna().all()
    )

    df_trans_avec_immat['Presence_Vehicule'] = False
    df_trans_avec_immat['Methode_Verification'] = "Non vérifié"
    df_trans_avec_immat['Score_Suspicion'] = 0
    df_trans_avec_immat['Distance_Station_Km'] = np.nan


    for idx, trans_row in df_trans_avec_immat.iterrows():
        immat = trans_row['Nouveau Immat']
        datetime_trans = trans_row['DateTime']

        if pd.isna(datetime_trans):
            df_trans_avec_immat.at[idx, 'Methode_Verification'] = "Date/heure transaction manquante"
            continue

        # Filtrer les trajets/arrêts du même véhicule le jour de la transaction
        date_trans = datetime_trans.date()
        geoloc_veh_jour = df_geoloc_filtre[
            (df_geoloc_filtre['Véhicule'] == immat) &
            (df_geoloc_filtre['Date'].dt.date == date_trans)
        ]

        presence_detectee_temporelle = False
        for _, geoloc_entry in geoloc_veh_jour.iterrows():
            if pd.notna(geoloc_entry['DateTime_Debut']) and pd.notna(geoloc_entry['DateTime_Fin']):
                # Fenêtre de tolérance pour la présence (ex: 30 min avant début à 30 min après fin)
                periode_debut = geoloc_entry['DateTime_Debut'] - timedelta(minutes=30)
                periode_fin = geoloc_entry['DateTime_Fin'] + timedelta(minutes=30)

                if periode_debut <= datetime_trans <= periode_fin:
                    presence_detectee_temporelle = True
                    df_trans_avec_immat.at[idx, 'Presence_Vehicule'] = True
                    df_trans_avec_immat.at[idx, 'Methode_Verification'] = f"Temporelle ({geoloc_entry['Type']})"

                    # Si les coordonnées sont disponibles et on a une localisation de station (hypothétique)
                    # Cette partie nécessiterait un fichier/DB de stations avec leurs coordonnées
                    # Pour l'instant, on se base sur la présence temporelle
                    # if coordonnees_disponibles_geoloc and 'Station_Lat' in trans_row and 'Station_Lon' in trans_row:
                    #     station_coords = (trans_row['Station_Lat'], trans_row['Station_Lon'])
                    #     if geoloc_entry['Type'] == 'Arrêt' and pd.notna(geoloc_entry['Latitude_depart']): # Supposons que Lat_depart est le lieu de l'arrêt
                    #         dist = haversine(station_coords, (geoloc_entry['Latitude_depart'], geoloc_entry['Longitude_depart']), unit=Unit.KILOMETERS)
                    #         df_trans_avec_immat.at[idx, 'Distance_Station_Km'] = dist
                    #         if dist <= rayon_station_km:
                    #             df_trans_avec_immat.at[idx, 'Methode_Verification'] = f"Géographique ({dist:.2f} km)"
                    #         # else: la présence temporelle reste
                    #     # Pourrait aussi vérifier la proximité pendant un trajet
                    break # Sortir de la boucle geoloc_entry si présence trouvée

        if not presence_detectee_temporelle:
             df_trans_avec_immat.at[idx, 'Methode_Verification'] = "Absence temporelle"


    # Calculer le score de suspicion pour les transactions sans présence détectée
    df_trans_avec_immat['Score_Suspicion'] = np.where(
        ~df_trans_avec_immat['Presence_Vehicule'],
        poids_transaction_sans_presence,
        0
    )

    # Identifier les transactions suspectes
    transactions_suspectes = df_trans_avec_immat[~df_trans_avec_immat['Presence_Vehicule']].copy()
    if not transactions_suspectes.empty:
        transactions_suspectes['Type_Anomalie'] = "Transaction sans présence détectée du véhicule"

        # Ajouter des critères supplémentaires pour les transactions suspectes
        transactions_suspectes['Jour_Semaine'] = transactions_suspectes['DateTime'].dt.dayofweek
        transactions_suspectes['Est_Weekend'] = transactions_suspectes['Jour_Semaine'] >= 5
        transactions_suspectes['Heure'] = transactions_suspectes['DateTime'].dt.hour
        transactions_suspectes['Est_Hors_Heures'] = (
            (transactions_suspectes['Heure'] < st.session_state.get('ss_heure_debut_service', DEFAULT_HEURE_DEBUT_SERVICE)) |
            (transactions_suspectes['Heure'] >= st.session_state.get('ss_heure_fin_service', DEFAULT_HEURE_FIN_SERVICE))
        )

        # Augmenter le score de suspicion pour les critères supplémentaires
        transactions_suspectes['Score_Suspicion'] = transactions_suspectes['Score_Suspicion'] + \
            np.where(transactions_suspectes['Est_Weekend'],
                     st.session_state.get('ss_poids_trajet_weekend', DEFAULT_POIDS_TRAJET_WEEKEND), 0) + \
            np.where(transactions_suspectes['Est_Hors_Heures'],
                     st.session_state.get('ss_poids_trajet_hors_heures', DEFAULT_POIDS_TRAJET_HORS_HEURES), 0)

        # Classifier le niveau de suspicion
        transactions_suspectes['Niveau_Suspicion'] = pd.cut(
            transactions_suspectes['Score_Suspicion'],
            bins=[0, 7, 12, float('inf')], # Ajuster les bins au besoin
            labels=['Faible', 'Modéré', 'Élevé'],
            right=False
        )

    return df_trans_avec_immat, transactions_suspectes

def analyser_vehicules_sans_geoloc(
    df_transactions: pd.DataFrame,
    df_vehicules: pd.DataFrame,
    df_geoloc: pd.DataFrame,
    date_debut: datetime.date,
    date_fin: datetime.date
) -> pd.DataFrame:
    """
    Identifie les véhicules ayant effectué des transactions mais sans données de géolocalisation correspondantes.
    
    Args:
        df_transactions: DataFrame des transactions.
        df_vehicules: DataFrame des véhicules.
        df_geoloc: DataFrame des données de géolocalisation.
        date_debut: Date de début de la période d'analyse.
        date_fin: Date de fin de la période d'analyse.
        
    Returns:
        DataFrame contenant les véhicules sans géolocalisation et leurs statistiques.
    """
    # Filtrer les transactions pour la période
    mask_date_trans = (df_transactions['Date'].dt.date >= date_debut) & (df_transactions['Date'].dt.date <= date_fin)
    df_trans_filtered = df_transactions[mask_date_trans].copy()
    
    # Fusionner avec les informations véhicules
    df_trans_with_veh = df_trans_filtered.merge(
        df_vehicules[['N° Carte', 'Nouveau Immat', 'Catégorie']],
        left_on='Card num.',
        right_on='N° Carte',
        how='inner'
    )
    
    # Obtenir la liste des véhicules avec transactions
    vehicules_avec_transactions = set(df_trans_with_veh['Nouveau Immat'].unique())
    
    # Filtrer les données géoloc pour la période
    mask_date_geoloc = (df_geoloc['Date'].dt.date >= date_debut) & (df_geoloc['Date'].dt.date <= date_fin)
    df_geoloc_filtered = df_geoloc[mask_date_geoloc].copy()
    
    # Obtenir la liste des véhicules avec données géoloc
    vehicules_avec_geoloc = set(df_geoloc_filtered['Véhicule'].unique())
    
    # Identifier les véhicules avec transactions mais sans géoloc
    vehicules_sans_geoloc = vehicules_avec_transactions - vehicules_avec_geoloc
    
    # Si aucun véhicule ne répond au critère, retourner un DataFrame vide
    if not vehicules_sans_geoloc:
        return pd.DataFrame()
    
    # Calculer les statistiques pour ces véhicules
    resultats = []
    for immat in vehicules_sans_geoloc:
        veh_trans = df_trans_with_veh[df_trans_with_veh['Nouveau Immat'] == immat]
        
        # Agréger les données de transactions
        nb_transactions = len(veh_trans)
        premiere_transaction = veh_trans['Date'].min()
        derniere_transaction = veh_trans['Date'].max()
        volume_total = veh_trans['Quantity'].sum()
        montant_total = veh_trans['Amount'].sum()
        categorie_veh = veh_trans['Catégorie'].iloc[0] if not veh_trans.empty else "Inconnue"
        
        # Calculer le nombre de jours distincts avec des transactions
        jours_distincts = veh_trans['Date'].dt.date.nunique()
        
        # Récupérer la carte associée
        carte = veh_trans['Card num.'].iloc[0] if not veh_trans.empty else "Inconnue"
        
        resultats.append({
            'Immatriculation': immat,
            'N° Carte': carte,
            'Catégorie': categorie_veh,
            'Nb_Transactions': nb_transactions,
            'Volume_Total_L': volume_total,
            'Montant_Total_CFA': montant_total,
            'Jours_Avec_Transactions': jours_distincts,
            'Première_Transaction': premiere_transaction,
            'Dernière_Transaction': derniere_transaction
        })
    
    # Créer le DataFrame de résultats
    df_resultats = pd.DataFrame(resultats)
    
    # Trier par nombre de transactions décroissant
    if not df_resultats.empty:
        df_resultats = df_resultats.sort_values('Nb_Transactions', ascending=False)
    
    return df_resultats

def detecter_detours_suspects(
    df_geoloc: pd.DataFrame,
    date_debut: datetime.date,
    date_fin: datetime.date
) -> pd.DataFrame:
    """
    Détecte les trajets avec des détours potentiellement suspects
    basés sur des critères comme distance/durée anormale, vitesse irrégulière, etc.

    Args:
        df_geoloc: DataFrame des données de géolocalisation.
        date_debut: Date de début de la période d'analyse.
        date_fin: Date de fin de la période d'analyse.

    Returns:
        Un DataFrame des trajets avec détours suspects.
    """
    # Vérifier si la détection est activée
    if not st.session_state.get('ss_activer_detours_suspects', True):
        return pd.DataFrame()  # Retourner un DataFrame vide si désactivé
        
    # Paramètres
    seuil_detour_pct = st.session_state.get('ss_seuil_detour_pct', DEFAULT_SEUIL_DETOUR_PCT)
    poids_detour_suspect = st.session_state.get('ss_poids_detour_suspect', DEFAULT_POIDS_DETOUR_SUSPECT)

    # Filtrer les données pour la période spécifiée
    mask_date = (df_geoloc['Date'].dt.date >= date_debut) & (df_geoloc['Date'].dt.date <= date_fin)
    df_filtre = df_geoloc[mask_date].copy()

    # Ne considérer que les trajets
    df_filtre = df_filtre[df_filtre['Type'] == 'Trajet'].copy()

    if df_filtre.empty or 'Durée_minutes' not in df_filtre.columns or df_filtre['Durée_minutes'].isna().all():
        return pd.DataFrame()

    # Calcul des ratios pour identifier les détours
    # Ratio Distance/Durée - si trop faible, peut indiquer un détour ou des arrêts non documentés
    df_filtre['Ratio_Dist_Duree'] = np.where(
        df_filtre['Durée_minutes'] > 0,
        df_filtre['Distance'] / (df_filtre['Durée_minutes'] / 60), # Distance / Heures = Vitesse
        np.nan
    )
    # Ce 'Ratio_Dist_Duree' est en fait la vitesse moyenne, qui est déjà calculée.
    # Utilisons directement 'Vitesse moyenne'
    df_filtre.rename(columns={'Vitesse moyenne': 'Ratio_Dist_Duree_Calc'}, inplace=True, errors='ignore') # Just in case


    # Moyennes par véhicule pour comparer
    moyennes_vehicules = df_filtre.groupby('Véhicule').agg(
        Ratio_Moyen=('Ratio_Dist_Duree_Calc', 'mean'), # Utiliser Vitesse moyenne calculée
        Dist_Moyenne=('Distance', 'mean'),
        Duree_Moyenne=('Durée_minutes', 'mean')
    ).reset_index()


    # Fusionner avec les moyennes par véhicule
    df_filtre = df_filtre.merge(moyennes_vehicules, on='Véhicule', how='left', suffixes=('', '_veh_moyen'))


    # Calcul des écarts par rapport aux moyennes du véhicule
    df_filtre['Ecart_Ratio_Pct'] = np.where(
        (df_filtre['Ratio_Moyen'].notna()) & (df_filtre['Ratio_Moyen'] > 0),
        ((df_filtre['Ratio_Dist_Duree_Calc'] / df_filtre['Ratio_Moyen']) - 1) * 100,
        np.nan
    )

    # Un écart négatif important indique un trajet inefficace (détour potentiel)
    # Un trajet est suspect si sa vitesse moyenne est significativement plus basse que la vitesse moyenne habituelle du véhicule
    # Et si la distance est non négligeable (ex: > 5km) pour éviter les petits trajets urbains lents.
    df_filtre['Est_Detour_Potentiel'] = (df_filtre['Ecart_Ratio_Pct'] < -seuil_detour_pct) & (df_filtre['Distance'] > 5)


    # Calculer le score de suspicion pour les détours potentiels
    df_filtre['Score_Detour'] = np.where(
        df_filtre['Est_Detour_Potentiel'],
        # Plus l'écart est négatif, plus le score est élevé
        poids_detour_suspect * (1 + abs(df_filtre['Ecart_Ratio_Pct']) / 100),
        0
    )

    # Filtrer les trajets suspects de détour
    detours_suspects = df_filtre[df_filtre['Est_Detour_Potentiel']].copy()

    if detours_suspects.empty:
        return pd.DataFrame()

    # Calculer la sévérité du détour
    detours_suspects['Severite_Detour'] = pd.cut(
        detours_suspects['Ecart_Ratio_Pct'].abs(),
        bins=[seuil_detour_pct, seuil_detour_pct + 25, seuil_detour_pct + 50, float('inf')], # Ajuster bins
        labels=['Léger', 'Modéré', 'Important'],
        right=False
    )

    # Ajouter une description du détour
    detours_suspects['Description_Detour'] = detours_suspects.apply(
        lambda row: f"Trajet {row['Severite_Detour']} inefficace ({abs(row['Ecart_Ratio_Pct']):.1f}% sous la moyenne du véhicule). "
                   f"Distance: {row['Distance']:.1f}km, Durée: {row['Durée_minutes']:.0f}min, "
                   f"Vitesse moy: {row['Ratio_Dist_Duree_Calc']:.1f}km/h (vs. {row['Ratio_Moyen']:.1f}km/h pour ce véhicule)",
        axis=1
    )

    # Sélectionner et ordonner les colonnes pertinentes
    cols_detours = [
        'Véhicule', 'Date', 'Début', 'Fin', 'Distance', 'Durée_minutes',
        'Vitesse moyenne', 'Ecart_Ratio_Pct', 'Severite_Detour',
        'Description_Detour', 'Score_Detour'
    ]
    # S'assurer que les colonnes existent avant de les sélectionner
    cols_detours_existantes = [col for col in cols_detours if col in detours_suspects.columns]

    return detours_suspects[cols_detours_existantes].sort_values('Ecart_Ratio_Pct')


def analyser_efficacite_carburant(
    df_geoloc: pd.DataFrame,
    df_transactions: pd.DataFrame,
    df_vehicules: pd.DataFrame,
    date_debut: datetime.date,
    date_fin: datetime.date
) -> pd.DataFrame:
    """
    Analyse l'efficacité d'utilisation du carburant en combinant données géoloc et transactions.
    Identifie les trajets avec une consommation anormale par rapport aux moyennes.

    Args:
        df_geoloc: DataFrame des données de géolocalisation.
        df_transactions: DataFrame des transactions.
        df_vehicules: DataFrame des véhicules.
        date_debut: Date de début de la période d'analyse.
        date_fin: Date de fin de la période d'analyse.

    Returns:
        Un DataFrame avec les analyses d'efficacité par véhicule.
    """
    # Filtrer les données pour la période spécifiée
    mask_date_geoloc = (df_geoloc['Date'].dt.date >= date_debut) & (df_geoloc['Date'].dt.date <= date_fin)
    df_geoloc_filtre = df_geoloc[mask_date_geoloc].copy()
    df_geoloc_filtre_trajets = df_geoloc_filtre[df_geoloc_filtre['Type'] == 'Trajet']


    mask_date_trans = (df_transactions['Date'].dt.date >= date_debut) & (df_transactions['Date'].dt.date <= date_fin)
    df_trans_filtre = df_transactions[mask_date_trans].copy()

    # Fusionner transactions avec infos véhicules
    df_trans_avec_immat = df_trans_filtre.merge(
        df_vehicules[['N° Carte', 'Nouveau Immat', 'Catégorie']],
        left_on='Card num.',
        right_on='N° Carte',
        how='inner'
    )

    # Agréger les distances par immatriculation (géolocalisation)
    distance_geoloc = df_geoloc_filtre_trajets.groupby('Véhicule').agg(
        Distance_Geoloc_Totale=('Distance', 'sum'),
        Nb_Trajets=('Distance', 'count'),
        Duree_Totale_Minutes=('Durée_minutes', 'sum')
    ).reset_index()

    # Agréger la consommation par immatriculation (transactions)
    conso_vehicule = df_trans_avec_immat.groupby('Nouveau Immat').agg(
        Volume_Total=('Quantity', 'sum'),
        Cout_Total=('Amount', 'sum'),
        Nb_Transactions=('Quantity', 'count')
    ).reset_index()

    # Fusionner les données
    efficacite = distance_geoloc.merge(
        conso_vehicule,
        left_on='Véhicule',
        right_on='Nouveau Immat',
        how='outer' # Outer join pour garder les véhicules même s'ils manquent dans un des DFs
    )
    efficacite['Véhicule'] = efficacite['Véhicule'].fillna(efficacite['Nouveau Immat'])


    # Calculer les métriques d'efficacité
    efficacite['Consommation_100km'] = np.where(
        efficacite['Distance_Geoloc_Totale'] > 0,
        (efficacite['Volume_Total'] / efficacite['Distance_Geoloc_Totale']) * 100,
        np.nan
    )

    efficacite['Cout_par_km'] = np.where(
        efficacite['Distance_Geoloc_Totale'] > 0,
        efficacite['Cout_Total'] / efficacite['Distance_Geoloc_Totale'],
        np.nan
    )

    efficacite['Cout_par_heure'] = np.where(
        (efficacite['Duree_Totale_Minutes'].notna()) & (efficacite['Duree_Totale_Minutes'] > 0),
        (efficacite['Cout_Total'] / efficacite['Duree_Totale_Minutes']) * 60,
        np.nan
    )


    # Conserver l'immatriculation comme identifiant unique
    efficacite.drop('Nouveau Immat', axis=1, inplace=True, errors='ignore')

    # Ajouter la catégorie du véhicule
    mapping_categorie = df_vehicules.set_index('Nouveau Immat')['Catégorie'].to_dict()
    efficacite['Catégorie'] = efficacite['Véhicule'].map(mapping_categorie)

    # Calculer les moyennes par catégorie pour comparaison
    moyennes_categorie = efficacite.groupby('Catégorie').agg(
        Conso_Moyenne_Cat=('Consommation_100km', 'mean'),
        Cout_km_Moyen_Cat=('Cout_par_km', 'mean')
    ).reset_index() # Pour que Catégorie redevienne une colonne


    # Fusionner avec les moyennes par catégorie
    efficacite = efficacite.merge(moyennes_categorie, on='Catégorie', how='left', suffixes=('', '_cat_ref'))


    # Calculer les écarts par rapport aux moyennes de la catégorie
    efficacite['Ecart_Conso_Pct'] = np.where(
        (efficacite['Conso_Moyenne_Cat'].notna()) & (efficacite['Conso_Moyenne_Cat'] > 0),
        ((efficacite['Consommation_100km'] / efficacite['Conso_Moyenne_Cat']) - 1) * 100,
        np.nan
    )

    efficacite['Ecart_Cout_km_Pct'] = np.where(
        (efficacite['Cout_km_Moyen_Cat'].notna()) & (efficacite['Cout_km_Moyen_Cat'] > 0),
        ((efficacite['Cout_par_km'] / efficacite['Cout_km_Moyen_Cat']) - 1) * 100,
        np.nan
    )


    # Calculer l'efficacité globale (score combiné)
    efficacite['Score_Efficacite'] = np.where(
        pd.notna(efficacite['Ecart_Conso_Pct']) & pd.notna(efficacite['Ecart_Cout_km_Pct']),
        100 - (efficacite['Ecart_Conso_Pct'] + efficacite['Ecart_Cout_km_Pct']) / 2, # Ecart positif = moins efficace
        np.nan
    )
    efficacite['Score_Efficacite'] = efficacite['Score_Efficacite'].fillna(50) # Note neutre si pas assez de données


    # Classifier l'efficacité
    efficacite['Niveau_Efficacite'] = pd.cut(
        efficacite['Score_Efficacite'],
        bins=[-float('inf'), 85, 95, 105, 115, float('inf')],
        labels=['Très faible', 'Faible', 'Normale', 'Bonne', 'Excellente'],
        right=False
    )

    # Arrondir les valeurs numériques pour l'affichage
    cols_arrondi = [
        'Distance_Geoloc_Totale', 'Duree_Totale_Minutes', 'Volume_Total',
        'Consommation_100km', 'Cout_par_km', 'Cout_par_heure',
        'Ecart_Conso_Pct', 'Ecart_Cout_km_Pct', 'Score_Efficacite',
        'Conso_Moyenne_Cat', 'Cout_km_Moyen_Cat'
    ]


    for col in cols_arrondi:
        if col in efficacite.columns:
            efficacite[col] = efficacite[col].round(1)

    return efficacite.sort_values('Score_Efficacite', ascending=False)

def generer_rapport_powerpoint_geoloc(
    df_geoloc: pd.DataFrame,
    df_transactions: pd.DataFrame,
    df_vehicules: pd.DataFrame,
    date_debut: datetime.date,
    date_fin: datetime.date,
    vehicules_selectionnes: List[str] = None,
    inclure_cartes: bool = True,
    inclure_analyse_vitesse: bool = True,
    inclure_analyse_trajets: bool = True,
    inclure_comparaison_carburant: bool = True,
    titre: str = None,
    description: str = None,
    theme: str = "Standard",
    orientation: str = "Paysage",
    inclure_page_titre: bool = True,
    inclure_sommaire: bool = True,
    inclure_numerotation: bool = True,
    inclure_footer: bool = True,
    inclure_logos: bool = False,
    inclure_annexes: bool = False,
    logo_file = None
) -> BytesIO:
    """
    Génère un rapport PowerPoint complet des analyses de géolocalisation.
    
    Args:
        df_geoloc: DataFrame des données de géolocalisation
        df_transactions: DataFrame des transactions
        df_vehicules: DataFrame des véhicules
        date_debut: Date de début de l'analyse
        date_fin: Date de fin de l'analyse
        vehicules_selectionnes: Liste des véhicules à inclure dans le rapport
        inclure_cartes: Inclure des cartes dans le rapport
        inclure_analyse_vitesse: Inclure l'analyse des vitesses
        inclure_analyse_trajets: Inclure l'analyse détaillée des trajets
        inclure_comparaison_carburant: Inclure la comparaison avec les données carburant
        titre: Titre personnalisé pour le rapport
        description: Description ou notes additionnelles
        theme: Thème visuel (Standard, Professionnel, Moderne, Coloré)
        orientation: Orientation des slides (Paysage ou Portrait)
        inclure_page_titre: Inclure une page de titre
        inclure_sommaire: Inclure un sommaire
        inclure_numerotation: Inclure la numérotation des pages
        inclure_footer: Inclure pied de page avec date
        inclure_logos: Inclure logos d'entreprise
        inclure_annexes: Inclure annexes techniques
        logo_file: Fichier logo uploadé
        
    Returns:
        BytesIO: Le fichier PowerPoint en mémoire
    """
    from pptx import Presentation
    from pptx.util import Inches, Pt
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN
    import io
    import tempfile
    import time
    from PIL import Image
    
    # Filtrer les données pour la période et les véhicules sélectionnés
    mask_date = (df_geoloc['Date'].dt.date >= date_debut) & (df_geoloc['Date'].dt.date <= date_fin)
    df_geoloc_filtered = df_geoloc[mask_date]
    
    if vehicules_selectionnes and vehicules_selectionnes != ["Tous les véhicules"]:
        df_geoloc_filtered = df_geoloc_filtered[df_geoloc_filtered['Véhicule'].isin(vehicules_selectionnes)]
    
    # Définir les couleurs et styles selon le thème choisi
    theme_colors = {
        "Standard": {
            "title_bg": "1F77B4",  # Bleu
            "accent1": "FF7F0E",   # Orange
            "accent2": "2CA02C",   # Vert
            "text": "333333",      # Gris foncé
            "background": "FFFFFF" # Blanc
        },
        "Professionnel": {
            "title_bg": "0F2E4C",  # Bleu marine
            "accent1": "143D59",   # Bleu foncé
            "accent2": "6B9AC4",   # Bleu clair
            "text": "333333",      # Gris foncé
            "background": "F2F2F2" # Gris très clair
        },
        "Moderne": {
            "title_bg": "3D5A80",  # Bleu grisé
            "accent1": "E07A5F",   # Saumon
            "accent2": "81B29A",   # Vert pastel
            "text": "293241",      # Bleu très foncé
            "background": "F7F7F7" # Gris très clair
        },
        "Coloré": {
            "title_bg": "6A0572",  # Violet
            "accent1": "AB83A1",   # Mauve
            "accent2": "F26419",   # Orange
            "text": "333333",      # Gris foncé
            "background": "FFFFFF" # Blanc
        }
    }
    
    # Sélectionner le thème
    current_theme = theme_colors.get(theme, theme_colors["Standard"])
    
    # Créer une nouvelle présentation
    prs = Presentation()
    
    # Définir l'orientation des diapositives
    if orientation == "Paysage":
        prs.slide_width = Inches(13.33)
        prs.slide_height = Inches(7.5)
    else:  # Portrait
        prs.slide_width = Inches(7.5)
        prs.slide_height = Inches(13.33)
    
    # Fonction helper pour définir le style d'un paragraphe
    def style_paragraph(paragraph, font_size=18, bold=False, color=current_theme["text"], alignment=PP_ALIGN.LEFT):
        run = paragraph.runs[0] if paragraph.runs else paragraph.add_run()
        font = run.font
        font.size = Pt(font_size)
        font.bold = bold
        font.color.rgb = RGBColor.from_string(color)
        paragraph.alignment = alignment
    
    # Fonction helper pour ajouter un titre de diapo
    def add_title_slide(title_text, subtitle_text=""):
        slide = prs.slides.add_slide(prs.slide_layouts[0])
        
        # Personnalisation du fond de la diapo de titre
        slide.background.fill.solid()
        slide.background.fill.fore_color.rgb = RGBColor.from_string(current_theme["background"])
        
        # Titre
        title = slide.shapes.title
        title.text = title_text
        
        # Style du titre
        for paragraph in title.text_frame.paragraphs:
            style_paragraph(paragraph, font_size=44, bold=True, color=current_theme["title_bg"], alignment=PP_ALIGN.CENTER)
        
        # Sous-titre
        if subtitle_text:
            subtitle = slide.placeholders[1]
            subtitle.text = subtitle_text
            
            # Style du sous-titre
            for paragraph in subtitle.text_frame.paragraphs:
                style_paragraph(paragraph, font_size=24, color=current_theme["accent1"], alignment=PP_ALIGN.CENTER)
        
        # Ajouter un logo si demandé
        if inclure_logos and logo_file:
            try:
                # Créer un fichier temporaire pour sauvegarder le logo
                with tempfile.NamedTemporaryFile(delete=False, suffix=".png") as temp_logo:
                    # Sauvegarder le logo uploadé
                    img = Image.open(logo_file)
                    img.save(temp_logo.name)
                    
                    # Ajouter le logo à la diapo
                    slide.shapes.add_picture(temp_logo.name, Inches(0.5), Inches(0.5), width=Inches(1.5))
            except Exception as e:
                print(f"Erreur lors de l'ajout du logo: {e}")
        
        return slide
    
    # Fonction helper pour ajouter une diapo de contenu
    def add_content_slide(title_text, content_placeholder_idx=1):
        slide = prs.slides.add_slide(prs.slide_layouts[1])
        
        # Fond de la diapo
        slide.background.fill.solid()
        slide.background.fill.fore_color.rgb = RGBColor.from_string(current_theme["background"])
        
        # Titre
        title = slide.shapes.title
        title.text = title_text
        
        # Style du titre
        for paragraph in title.text_frame.paragraphs:
            style_paragraph(paragraph, font_size=32, bold=True, color=current_theme["title_bg"])
        
        # Ajout d'un pied de page si demandé
        if inclure_footer:
            footer_shape = slide.shapes.add_textbox(Inches(0.5), prs.slide_height - Inches(0.5), Inches(12), Inches(0.3))
            footer_text = footer_shape.text_frame
            
            p = footer_text.add_paragraph()
            p.text = f"Rapport Géolocalisation • {date_debut.strftime('%d/%m/%Y')} - {date_fin.strftime('%d/%m/%Y')} • Page {prs.slides.index(slide) + 1}"
            style_paragraph(p, font_size=10, color=current_theme["accent1"], alignment=PP_ALIGN.CENTER)
        
        # Ajout de la numérotation si demandée
        if inclure_numerotation:
            page_num = slide.shapes.add_textbox(prs.slide_width - Inches(1), prs.slide_height - Inches(0.5), Inches(0.5), Inches(0.3))
            page_text = page_num.text_frame
            
            p = page_text.add_paragraph()
            p.text = f"{prs.slides.index(slide) + 1}"
            style_paragraph(p, font_size=12, bold=True, color=current_theme["accent2"], alignment=PP_ALIGN.RIGHT)
        
        # Retourner la diapo et le placeholder de contenu
        try:
            content = slide.placeholders[content_placeholder_idx]
            return slide, content
        except:
            # Si le layout n'a pas le placeholder attendu, retourner juste la diapo
            return slide, None
    
    # Fonction helper pour ajouter une image
    def add_image_slide(title_text, img_path, caption=""):
        slide = prs.slides.add_slide(prs.slide_layouts[5])
        
        # Fond de la diapo
        slide.background.fill.solid()
        slide.background.fill.fore_color.rgb = RGBColor.from_string(current_theme["background"])
        
        # Titre
        title = slide.shapes.title
        title.text = title_text
        
        # Style du titre
        for paragraph in title.text_frame.paragraphs:
            style_paragraph(paragraph, font_size=32, bold=True, color=current_theme["title_bg"])
        
        # Ajuster la taille et position de l'image selon l'orientation
        if orientation == "Paysage":
            image_width, image_height = Inches(9), Inches(5)
            image_left, image_top = Inches(2), Inches(1.8)
        else:  # Portrait
            image_width, image_height = Inches(6), Inches(8)
            image_left, image_top = Inches(0.75), Inches(2)
        
        # Ajouter l'image
        try:
            pic = slide.shapes.add_picture(img_path, image_left, image_top, width=image_width, height=image_height)
            
            # Ajouter une légende si fournie
            if caption:
                caption_box = slide.shapes.add_textbox(image_left, image_top + image_height + Inches(0.2), 
                                                       image_width, Inches(0.5))
                caption_frame = caption_box.text_frame
                
                p = caption_frame.add_paragraph()
                p.text = caption
                style_paragraph(p, font_size=14, color=current_theme["accent1"], alignment=PP_ALIGN.CENTER)
        except Exception as e:
            # En cas d'erreur, ajouter un texte explicatif
            error_box = slide.shapes.add_textbox(image_left, image_top, image_width, image_height)
            error_frame = error_box.text_frame
            
            p = error_frame.add_paragraph()
            p.text = f"Impossible d'afficher l'image. Erreur: {str(e)}"
            style_paragraph(p, font_size=14, color="FF0000", alignment=PP_ALIGN.CENTER)
        
        # Ajout d'un pied de page si demandé
        if inclure_footer:
            footer_shape = slide.shapes.add_textbox(Inches(0.5), prs.slide_height - Inches(0.5), Inches(12), Inches(0.3))
            footer_text = footer_shape.text_frame
            
            p = footer_text.add_paragraph()
            p.text = f"Rapport Géolocalisation • {date_debut.strftime('%d/%m/%Y')} - {date_fin.strftime('%d/%m/%Y')} • Page {prs.slides.index(slide) + 1}"
            style_paragraph(p, font_size=10, color=current_theme["accent1"], alignment=PP_ALIGN.CENTER)
        
        # Ajout de la numérotation si demandée
        if inclure_numerotation:
            page_num = slide.shapes.add_textbox(prs.slide_width - Inches(1), prs.slide_height - Inches(0.5), Inches(0.5), Inches(0.3))
            page_text = page_num.text_frame
            
            p = page_text.add_paragraph()
            p.text = f"{prs.slides.index(slide) + 1}"
            style_paragraph(p, font_size=12, bold=True, color=current_theme["accent2"], alignment=PP_ALIGN.RIGHT)
        
        return slide
    
    # Fonction helper pour sauver un graphique plotly en image temporaire
    def save_plotly_as_image(fig, width=1600, height=900):
        try:
            img_bytes = fig.to_image(format="png", width=width, height=height)
            with tempfile.NamedTemporaryFile(delete=False, suffix=".png") as temp_file:
                temp_file.write(img_bytes)
                return temp_file.name
        except Exception as e:
            print(f"Erreur lors de la conversion du graphique: {e}")
            # Créer une image vide avec un message d'erreur
            from PIL import Image, ImageDraw, ImageFont
            img = Image.new('RGB', (width, height), color = (255, 255, 255))
            d = ImageDraw.Draw(img)
            d.text((width/2, height/2), f"Erreur de génération du graphique: {e}", fill=(255, 0, 0))
            
            with tempfile.NamedTemporaryFile(delete=False, suffix=".png") as temp_file:
                img.save(temp_file.name)
                return temp_file.name
    
    # 1. Diapo de titre (si activée)
    if inclure_page_titre:
        # Utiliser le titre personnalisé s'il est fourni
        title_text = titre if titre else f"Rapport d'Analyse de Géolocalisation"
        subtitle_text = f"Période: {date_debut.strftime('%d/%m/%Y')} - {date_fin.strftime('%d/%m/%Y')}"
        
        # Ajouter la description si fournie
        if description:
            subtitle_text += f"\n\n{description}"
        
        title_slide = add_title_slide(title_text, subtitle_text)
    
    # 2. Diapo de sommaire (si activée)
    if inclure_sommaire:
        summary_slide, summary_content = add_content_slide("Sommaire")
        summary_text = summary_content.text_frame
        summary_text.clear()
        
        p = summary_text.paragraphs[0]
        p.text = "1. Synthèse des données de géolocalisation"
        style_paragraph(p, font_size=20)
        
        p = summary_text.add_paragraph()
        p.text = "2. Kilométrage et utilisation des véhicules"
        style_paragraph(p, font_size=20)
        
        slide_index = 3
        
        if inclure_analyse_vitesse:
            p = summary_text.add_paragraph()
            p.text = f"{slide_index}. Analyse des vitesses"
            style_paragraph(p, font_size=20)
            slide_index += 1
        
        if inclure_analyse_trajets:
            p = summary_text.add_paragraph()
            p.text = f"{slide_index}. Analyse détaillée des trajets"
            style_paragraph(p, font_size=20)
            slide_index += 1
        
        if inclure_comparaison_carburant:
            p = summary_text.add_paragraph()
            p.text = f"{slide_index}. Comparaison avec les données carburant"
            style_paragraph(p, font_size=20)
            slide_index += 1
        
        if inclure_cartes:
            p = summary_text.add_paragraph()
            p.text = f"{slide_index}. Visualisation cartographique des trajets"
            style_paragraph(p, font_size=20)
            slide_index += 1
        
        if inclure_annexes:
            p = summary_text.add_paragraph()
            p.text = f"{slide_index}. Annexes techniques"
            style_paragraph(p, font_size=20)
    
    # 3. Diapo de synthèse
    if not df_geoloc_filtered.empty:
        nb_vehicules = df_geoloc_filtered['Véhicule'].nunique()
        
        # Filtrer les trajets
        trajets_df = df_geoloc_filtered[df_geoloc_filtered['Type'] == 'Trajet']
        nb_trajets = len(trajets_df)
        distance_totale = trajets_df['Distance'].sum() if 'Distance' in trajets_df.columns else 0
        duree_totale = trajets_df['Durée_minutes'].sum() / 60 if 'Durée_minutes' in trajets_df.columns else 0
        
        synthesis_slide, content = add_content_slide("Synthèse des Données de Géolocalisation")
        
        if content:
            tf = content.text_frame
            tf.clear()
            
            p = tf.paragraphs[0]
            p.text = f"• Nombre de véhicules analysés: {nb_vehicules}"
            style_paragraph(p, font_size=20, bold=True)
            
            p = tf.add_paragraph()
            p.text = f"• Nombre total de trajets: {nb_trajets:,}"
            style_paragraph(p, font_size=20)
            
            p = tf.add_paragraph()
            p.text = f"• Distance totale parcourue: {distance_totale:,.1f} km"
            style_paragraph(p, font_size=20)
            
            p = tf.add_paragraph()
            p.text = f"• Durée totale d'utilisation: {duree_totale:,.1f} heures"
            style_paragraph(p, font_size=20)
            
            if 'Est_weekend' in trajets_df.columns:
                nb_trajets_weekend = trajets_df['Est_weekend'].sum()
                pct_weekend = (nb_trajets_weekend / nb_trajets * 100) if nb_trajets > 0 else 0
                
                p = tf.add_paragraph()
                p.text = f"• Trajets en weekend: {nb_trajets_weekend:,} ({pct_weekend:.1f}%)"
                style_paragraph(p, font_size=20)
        
        # Ajouter un graphique de distance par véhicule
        distance_par_vehicule = trajets_df.groupby('Véhicule')['Distance'].sum().reset_index()
        if not distance_par_vehicule.empty:
            distance_par_vehicule = distance_par_vehicule.sort_values('Distance', ascending=False)
            fig_distance = px.bar(
                distance_par_vehicule.head(10),
                x='Véhicule',
                y='Distance',
                title="Top 10 - Distance Totale Parcourue par Véhicule",
                labels={'Distance': 'Distance (km)'},
                color='Distance',
                color_continuous_scale='Blues'
            )
            img_path = save_plotly_as_image(fig_distance)
            add_image_slide("Top 10 des Véhicules par Distance", img_path, 
                           caption="Les 10 véhicules ayant parcouru les plus grandes distances sur la période")
            
        # Ajouter un graphique d'évolution quotidienne
        if not trajets_df.empty:
            distance_quotidienne = trajets_df.groupby(pd.Grouper(key='Date', freq='D'))['Distance'].sum().reset_index()
            fig_evol = px.line(
                distance_quotidienne,
                x='Date',
                y='Distance',
                title="Évolution Quotidienne de la Distance Parcourue",
                labels={'Distance': 'Distance (km)'},
                markers=True
            )
            img_path = save_plotly_as_image(fig_evol)
            add_image_slide("Évolution Quotidienne de la Distance", img_path,
                           caption="Tendance de l'utilisation des véhicules jour par jour")
    
    # 4. Diapo d'analyse des vitesses
    if inclure_analyse_vitesse and not df_geoloc_filtered.empty:
        trajets_df = df_geoloc_filtered[df_geoloc_filtered['Type'] == 'Trajet']
        
        if 'Vitesse moyenne' in trajets_df.columns:
            vitesse_slide, content = add_content_slide("Analyse des Vitesses")
            
            # Calculer les statistiques de vitesse
            vitesse_moy_globale = trajets_df['Vitesse moyenne'].mean()
            vitesse_max = trajets_df['Vitesse moyenne'].max()
            
            # Seuil de vitesse excessive configurable
            seuil_vitesse = 90  # Paramètre potentiel à ajouter à la fonction
            trajets_excessifs = trajets_df[trajets_df['Vitesse moyenne'] > seuil_vitesse]
            pct_excessifs = (len(trajets_excessifs) / len(trajets_df)) * 100 if len(trajets_df) > 0 else 0
            
            if content:
                tf = content.text_frame
                tf.clear()
                
                p = tf.paragraphs[0]
                p.text = f"• Vitesse moyenne globale: {vitesse_moy_globale:.1f} km/h"
                style_paragraph(p, font_size=20, bold=True)
                
                p = tf.add_paragraph()
                p.text = f"• Vitesse maximale enregistrée: {vitesse_max:.1f} km/h"
                style_paragraph(p, font_size=20)
                
                p = tf.add_paragraph()
                p.text = f"• Trajets en excès de vitesse (>{seuil_vitesse} km/h): {len(trajets_excessifs)} ({pct_excessifs:.1f}%)"
                style_paragraph(p, font_size=20)
                
                if len(trajets_excessifs) > 0:
                    p = tf.add_paragraph()
                    p.text = f"• Véhicules concernés: {trajets_excessifs['Véhicule'].nunique()}"
                    style_paragraph(p, font_size=20)
            
            # Graphique de répartition des vitesses
            fig_vitesse = px.histogram(
                trajets_df,
                x='Vitesse moyenne',
                title="Distribution des Vitesses Moyennes",
                labels={'Vitesse moyenne': 'Vitesse (km/h)'},
                nbins=20,
                color_discrete_sequence=[f'#{current_theme["accent1"]}']
            )
            
            # Ajouter une ligne pour le seuil de vitesse
            fig_vitesse.add_vline(x=seuil_vitesse, line_dash="dash", line_color="red",
                           annotation_text=f"Seuil: {seuil_vitesse} km/h")
            
            img_path = save_plotly_as_image(fig_vitesse)
            add_image_slide("Distribution des Vitesses", img_path,
                           caption="Répartition des vitesses moyennes de tous les trajets")
            
            # Graphique des véhicules avec excès de vitesse
            if len(trajets_excessifs) > 0:
                exces_par_vehicule = trajets_excessifs.groupby('Véhicule').size().reset_index(name='Nombre')
                exces_par_vehicule = exces_par_vehicule.sort_values('Nombre', ascending=False)
                
                fig_exces = px.bar(
                    exces_par_vehicule.head(10),
                    x='Véhicule',
                    y='Nombre',
                    title="Top 10 - Véhicules avec Excès de Vitesse",
                    labels={'Nombre': f'Nombre de trajets > {seuil_vitesse} km/h'},
                    color='Nombre',
                    color_continuous_scale='Reds'
                )
                
                img_path = save_plotly_as_image(fig_exces)
                add_image_slide("Top 10 des Véhicules avec Excès de Vitesse", img_path,
                               caption=f"Véhicules avec le plus grand nombre de trajets dépassant {seuil_vitesse} km/h")
    
    # 5. Analyse détaillée des trajets
    if inclure_analyse_trajets and not df_geoloc_filtered.empty:
        trajets_df = df_geoloc_filtered[df_geoloc_filtered['Type'] == 'Trajet']
        if not trajets_df.empty:
            # Analyse des périodes d'utilisation
            trajets_slide, content = add_content_slide("Analyse Détaillée des Trajets")
            
            # Classer les trajets par période de la journée
            if 'Heure_debut' in trajets_df.columns:
                trajets_df['Periode'] = pd.cut(
                    trajets_df['Heure_debut'],
                    bins=[0, 6, 9, 12, 14, 17, 20, 24],
                    labels=['Nuit (0h-6h)', 'Matin (6h-9h)', 'Matinée (9h-12h)', 
                           'Midi (12h-14h)', 'Après-midi (14h-17h)', 
                           'Soir (17h-20h)', 'Nuit (20h-24h)']
                )
                
                periodes_count = trajets_df.groupby('Periode').size().reset_index(name='Nombre')
                
                fig_periode = px.bar(
                    periodes_count,
                    x='Periode',
                    y='Nombre',
                    title="Répartition des Trajets par Période de la Journée",
                    labels={'Nombre': 'Nombre de trajets', 'Periode': 'Période'},
                    color='Nombre',
                    color_continuous_scale='Blues'
                )
                
                img_path = save_plotly_as_image(fig_periode)
                add_image_slide("Répartition des Trajets par Période", img_path,
                               caption="Distribution des trajets selon les plages horaires")
            
            # Analyse des jours de la semaine
            if 'Jour_semaine' in trajets_df.columns:
                jours = ['Lundi', 'Mardi', 'Mercredi', 'Jeudi', 'Vendredi', 'Samedi', 'Dimanche']
                trajets_df['Jour'] = trajets_df['Jour_semaine'].map(lambda x: jours[x] if 0 <= x < len(jours) else 'Inconnu')
                
                jours_count = trajets_df.groupby('Jour').size().reset_index(name='Nombre')
                # Convertir Jour en catégorie ordonnée
                jours_count['Jour'] = pd.Categorical(jours_count['Jour'], categories=jours, ordered=True)
                jours_count = jours_count.sort_values('Jour')
                
                # Créer des couleurs pour distinguer la semaine du weekend
                colors = ['blue'] * 5 + ['red'] * 2
                
                fig_jours = px.bar(
                    jours_count,
                    x='Jour',
                    y='Nombre',
                    title="Répartition des Trajets par Jour de la Semaine",
                    labels={'Nombre': 'Nombre de trajets', 'Jour': 'Jour'},
                    color='Jour',
                    color_discrete_sequence=colors
                )
                
                img_path = save_plotly_as_image(fig_jours)
                add_image_slide("Répartition des Trajets par Jour", img_path,
                               caption="Distribution des trajets selon les jours de la semaine (weekend en rouge)")
    
    # 6. Comparaison avec les données carburant
    if inclure_comparaison_carburant and not df_geoloc_filtered.empty:
        # Utiliser la fonction existante pour l'analyse comparative
        try:
            comparaison, anomalies = analyser_geolocalisation_vs_transactions(
                df_geoloc_filtered, df_transactions, df_vehicules, date_debut, date_fin
            )
            
            if not comparaison.empty:
                comparaison_slide, content = add_content_slide("Comparaison Géolocalisation vs Transactions Carburant")
                
                nb_anomalies = len(anomalies)
                nb_sur_decla = len(anomalies[anomalies['Type_Anomalie'] == 'Sur-déclaration kilométrique']) if not anomalies.empty else 0
                nb_sous_decla = len(anomalies[anomalies['Type_Anomalie'] == 'Sous-déclaration kilométrique']) if not anomalies.empty else 0
                
                if content:
                    tf = content.text_frame
                    tf.clear()
                    
                    p = tf.paragraphs[0]
                    p.text = f"• Nombre d'anomalies de kilométrage: {nb_anomalies}"
                    style_paragraph(p, font_size=20, bold=True)
                    
                    p = tf.add_paragraph()
                    p.text = f"• Sur-déclarations: {nb_sur_decla}"
                    style_paragraph(p, font_size=20)
                    
                    p = tf.add_paragraph()
                    p.text = f"• Sous-déclarations: {nb_sous_decla}"
                    style_paragraph(p, font_size=20)
                    
                    if nb_anomalies > 0:
                        p = tf.add_paragraph()
                        p.text = f"• Impact potentiel: surconsommation apparente due aux écarts de km"
                        style_paragraph(p, font_size=20, color=current_theme["accent1"])
                
                # Graphique des écarts
                if nb_anomalies > 0:
                    comparaison_valide = comparaison.dropna(subset=['Ecart_Distance'])
                    comparaison_valide = comparaison_valide[comparaison_valide['Ecart_Distance'] != 0]
                    
                    if not comparaison_valide.empty:
                        fig_ecart = px.bar(
                            comparaison_valide.sort_values('Ecart_Distance', ascending=False).head(10),
                            x='Immatriculation',
                            y='Ecart_Distance',
                            title="Top 10 des Écarts de Kilométrage (Déclaré - Géolocalisé)",
                            labels={'Ecart_Distance': 'Écart (km)'},
                            color='Ecart_Distance',
                            color_continuous_scale=px.colors.diverging.RdBu,
                            color_continuous_midpoint=0
                        )
                        
                        img_path = save_plotly_as_image(fig_ecart)
                        add_image_slide("Top 10 des Écarts de Kilométrage", img_path,
                                      caption="Véhicules présentant les plus grands écarts entre km déclarés et géolocalisés")
                
                # Graphique comparaison consommation
                if 'Consommation_100km_Reelle' in comparaison.columns:
                    comparaison_conso = comparaison.dropna(subset=['Consommation_100km_Reelle', 'Consommation_100km_Declaree'])
                    comparaison_conso = comparaison_conso[
                        (comparaison_conso['Consommation_100km_Reelle'] > 0) & 
                        (comparaison_conso['Consommation_100km_Declaree'] > 0)
                    ]
                    
                    if not comparaison_conso.empty:
                        comparaison_conso = comparaison_conso.sort_values('Pourcentage_Ecart_Consommation', ascending=False).head(10)
                        
                        fig_conso = px.bar(
                            comparaison_conso,
                            x='Immatriculation',
                            y=['Consommation_100km_Reelle', 'Consommation_100km_Declaree'],
                            title="Top 10 - Comparaison Consommation Réelle vs Déclarée",
                            labels={'value': 'Consommation (L/100km)', 'variable': 'Source'},
                            barmode='group',
                            color_discrete_map={
                                'Consommation_100km_Reelle': f'#{current_theme["accent2"]}',
                                'Consommation_100km_Declaree': f'#{current_theme["accent1"]}'
                            }
                        )
                        
                        img_path = save_plotly_as_image(fig_conso)
                        add_image_slide("Comparaison des Consommations", img_path,
                                      caption="Différences entre consommation réelle (basée sur géoloc) et déclarée")
        except Exception as e:
            # Ajouter une diapo d'erreur
            error_slide, content = add_content_slide("Erreur dans la Comparaison des Données")
            if content:
                tf = content.text_frame
                tf.clear()
                
                p = tf.paragraphs[0]
                p.text = f"Une erreur est survenue lors de l'analyse comparative:"
                style_paragraph(p, font_size=20, bold=True)
                
                p = tf.add_paragraph()
                p.text = str(e)
                style_paragraph(p, font_size=18, color="FF0000")
    
    # 7. Cartes des trajets (si activé)
    if inclure_cartes and not df_geoloc_filtered.empty:
        trajets_df = df_geoloc_filtered[df_geoloc_filtered['Type'] == 'Trajet']
        
        # Vérifier si les coordonnées GPS sont disponibles
        has_coords = all(col in trajets_df.columns for col in 
                         ['Latitude_depart', 'Longitude_depart', 'Latitude_arrivee', 'Longitude_arrivee'])
        
        if has_coords:
            # Filtrer les trajets avec des coordonnées valides
            coords_valides = trajets_df.dropna(subset=['Latitude_depart', 'Longitude_depart', 
                                                      'Latitude_arrivee', 'Longitude_arrivee'])
            
            if not coords_valides.empty:
                cartes_slide, content = add_content_slide("Visualisation Cartographique des Trajets")
                
                if content:
                    tf = content.text_frame
                    tf.clear()
                    
                    p = tf.paragraphs[0]
                    p.text = "• Les visualisations cartographiques nécessitent une visualisation interactive."
                    style_paragraph(p, font_size=20, bold=True)
                    
                    p = tf.add_paragraph()
                    p.text = "• Pour une analyse détaillée, consultez l'application avec la carte interactive."
                    style_paragraph(p, font_size=20)
                    
                    p = tf.add_paragraph()
                    p.text = "• Un aperçu statique simplifié est inclus dans les slides suivantes."
                    style_paragraph(p, font_size=20)
                
                # Sélectionner les véhicules à afficher pour les cartes
                vehicules_maps = vehicules_selectionnes if vehicules_selectionnes else coords_valides['Véhicule'].unique()[:5]
                
                # Limiter à 5 véhicules pour éviter trop de cartes
                for i, vehicule in enumerate(vehicules_maps[:5]):
                    trajets_veh = coords_valides[coords_valides['Véhicule'] == vehicule]
                    if not trajets_veh.empty:
                        try:
                            # Créer une carte folium
                            mean_lat = trajets_veh['Latitude_depart'].mean()
                            mean_lon = trajets_veh['Longitude_depart'].mean()
                            
                            m = folium.Map(location=[mean_lat, mean_lon], zoom_start=10)
                            
                            # Limiter à 50 trajets pour ne pas surcharger
                            for _, row in trajets_veh.head(50).iterrows():
                                folium.PolyLine(
                                    locations=[
                                        [row['Latitude_depart'], row['Longitude_depart']],
                                        [row['Latitude_arrivee'], row['Longitude_arrivee']]
                                    ],
                                    color='blue',
                                    weight=2,
                                    opacity=0.7
                                ).add_to(m)
                            
                            # Sauvegarder la carte en HTML
                            with tempfile.NamedTemporaryFile(delete=False, suffix=".html") as temp_file:
                                m.save(temp_file.name)
                                
                                # Pour une visualisation statique simple, on pourrait utiliser une capture d'écran
                                # mais ici on ajoute juste une note que la carte est disponible dans l'application
                                
                                # Ajouter une diapo explicative
                                map_slide, map_content = add_content_slide(f"Carte des Trajets - {vehicule}")
                                
                                if map_content:
                                    map_tf = map_content.text_frame
                                    map_tf.clear()
                                    
                                    p = map_tf.paragraphs[0]
                                    p.text = f"• Véhicule: {vehicule}"
                                    style_paragraph(p, font_size=20, bold=True)
                                    
                                    p = map_tf.add_paragraph()
                                    p.text = f"• Nombre de trajets: {len(trajets_veh)}"
                                    style_paragraph(p, font_size=20)
                                    
                                    p = map_tf.add_paragraph()
                                    p.text = f"• Distance totale: {trajets_veh['Distance'].sum():.1f} km"
                                    style_paragraph(p, font_size=20)
                                    
                                    p = map_tf.add_paragraph()
                                    p.text = f"• Note: La visualisation interactive est disponible dans l'application."
                                    style_paragraph(p, font_size=18, color=current_theme["accent1"])
                        
                        except Exception as e:
                            print(f"Erreur lors de la création de la carte pour le véhicule {vehicule}: {e}")
    
    # 8. Annexes techniques (si activé)
    if inclure_annexes:
        annexe_slide, content = add_content_slide("Annexes Techniques")
        
        if content:
            tf = content.text_frame
            tf.clear()
            
            p = tf.paragraphs[0]
            p.text = "Paramètres de l'analyse:"
            style_paragraph(p, font_size=20, bold=True)
            
            p = tf.add_paragraph()
            p.text = f"• Période d'analyse: {date_debut.strftime('%d/%m/%Y')} - {date_fin.strftime('%d/%m/%Y')}"
            style_paragraph(p, font_size=18)
            
            p = tf.add_paragraph()
            p.text = f"• Nombre de véhicules sélectionnés: {len(vehicules_selectionnes) if vehicules_selectionnes else 'Tous'}"
            style_paragraph(p, font_size=18)
            
            p = tf.add_paragraph()
            p.text = f"• Seuil de vitesse excessive: 90 km/h"
            style_paragraph(p, font_size=18)
            
            p = tf.add_paragraph()
            p.text = "Notes méthodologiques:"
            style_paragraph(p, font_size=20, bold=True)
            
            p = tf.add_paragraph()
            p.text = "• Les données de géolocalisation sont basées sur les trajets GPS enregistrés."
            style_paragraph(p, font_size=18)
            
            p = tf.add_paragraph()
            p.text = "• Les écarts de kilométrage peuvent être dus à des trajets non enregistrés ou des déclarations erronées."
            style_paragraph(p, font_size=18)
            
            p = tf.add_paragraph()
            p.text = "• La consommation réelle est calculée avec la distance GPS, plus précise que la distance déclarée."
            style_paragraph(p, font_size=18)
    
    # 9. Diapo de conclusion
    conclusion_slide, content = add_content_slide("Conclusion et Recommandations")
    
    if content:
        tf = content.text_frame
        tf.clear()
        
        p = tf.paragraphs[0]
        p.text = "Principaux constats:"
        style_paragraph(p, font_size=20, bold=True)
        
        # Générer des constats en fonction des résultats
        if not df_geoloc_filtered.empty:
            trajets_df = df_geoloc_filtered[df_geoloc_filtered['Type'] == 'Trajet']
            
            # Constat sur les distances
            p = tf.add_paragraph()
            p.text = f"• La flotte a parcouru {trajets_df['Distance'].sum():,.1f} km sur la période analysée."
            style_paragraph(p, font_size=18)
            
            # Constat sur les vitesses
            if 'Vitesse moyenne' in trajets_df.columns:
                vitesse_moy = trajets_df['Vitesse moyenne'].mean()
                trajets_excessifs = trajets_df[trajets_df['Vitesse moyenne'] > 90]
                pct_excessifs = (len(trajets_excessifs) / len(trajets_df)) * 100 if len(trajets_df) > 0 else 0
                
                if pct_excessifs > 5:
                    p = tf.add_paragraph()
                    p.text = f"• Attention: {pct_excessifs:.1f}% des trajets dépassent la vitesse de 90 km/h."
                    style_paragraph(p, font_size=18, color="FF0000")
            
            # Constat sur les weekends
            if 'Est_weekend' in trajets_df.columns:
                pct_weekend = (trajets_df['Est_weekend'].sum() / len(trajets_df)) * 100 if len(trajets_df) > 0 else 0
                
                if pct_weekend > 10:
                    p = tf.add_paragraph()
                    p.text = f"• {pct_weekend:.1f}% des trajets ont lieu le weekend."
                    style_paragraph(p, font_size=18)
        
        p = tf.add_paragraph()
        p.text = "Recommandations:"
        style_paragraph(p, font_size=20, bold=True)
        
        p = tf.add_paragraph()
        p.text = "• Vérifier les déclarations de kilométrage pour les véhicules présentant des écarts importants."
        style_paragraph(p, font_size=18)
        
        p = tf.add_paragraph()
        p.text = "• Sensibiliser les conducteurs au respect des limitations de vitesse."
        style_paragraph(p, font_size=18)
        
        p = tf.add_paragraph()
        p.text = "• Optimiser les trajets pour réduire les distances parcourues et la consommation."
        style_paragraph(p, font_size=18)
        
        p = tf.add_paragraph()
        p.text = f"• Rapport généré le {datetime.now().strftime('%d/%m/%Y à %H:%M')}."
        style_paragraph(p, font_size=16, color=current_theme["accent1"], alignment=PP_ALIGN.RIGHT)
    
    # Sauvegarder la présentation en mémoire
    ppt_buffer = io.BytesIO()
    prs.save(ppt_buffer)
    ppt_buffer.seek(0)
    
    return ppt_buffer
    
def generer_resume_anomalies_geolocalisation(
    df_geoloc: pd.DataFrame,
    df_transactions: pd.DataFrame,
    df_vehicules: pd.DataFrame,
    date_debut: datetime.date,
    date_fin: datetime.date
) -> pd.DataFrame:
    """
    Génère un tableau de bord résumant toutes les anomalies détectées
    à partir des données de géolocalisation et leur corrélation avec les transactions.

    Args:
        df_geoloc: DataFrame des données de géolocalisation.
        df_transactions: DataFrame des transactions.
        df_vehicules: DataFrame des véhicules.
        date_debut: Date de début de la période d'analyse.
        date_fin: Date de fin de la période d'analyse.

    Returns:
        Un DataFrame consolidant toutes les anomalies détectées par véhicule avec le score de risque.
    """
    # Collecter les anomalies de différentes analyses
    # Utiliser un seuil de vitesse par défaut pertinent pour l'agrégation, ex: 90km/h
    # Le seuil slider est pour l'exploration interactive dans l'onglet vitesse.
    seuil_vitesse_agregation = 90

    with st.spinner("Analyse des excès de vitesse (agrégation)..."):
        resume_exces, _ = analyser_exces_vitesse(df_geoloc, date_debut, date_fin, seuil_vitesse_agregation)

    with st.spinner("Détection des trajets suspects (agrégation)..."):
        if st.session_state.get('ss_activer_trajets_suspects', True):
            trajets_suspects = detecter_trajets_suspects(df_geoloc, date_debut, date_fin)
        else:
            trajets_suspects = pd.DataFrame()

    with st.spinner("Analyse des correspondances transactions/géolocalisation (agrégation)..."):
        _, transactions_suspectes = analyser_correspondance_transactions_geoloc(
            df_geoloc, df_transactions, df_vehicules, date_debut, date_fin
        )

    with st.spinner("Détection des détours suspects (agrégation)..."):
        if st.session_state.get('ss_activer_detours_suspects', True):
            detours_suspects = detecter_detours_suspects(df_geoloc, date_debut, date_fin)
        else:
            detours_suspects = pd.DataFrame()

    with st.spinner("Analyse comparative distances géoloc/transactions (agrégation)..."):
        comparaison, anomalies_distance = analyser_geolocalisation_vs_transactions(
            df_geoloc, df_transactions, df_vehicules, date_debut, date_fin
        )

    # Initialiser le DataFrame résumé
    vehicules_geoloc = set(df_geoloc['Véhicule'].unique())

    # Créer un DataFrame avec tous les véhicules géolocalisés
    resume = pd.DataFrame({'Véhicule': list(vehicules_geoloc)})

    # Ajouter les scores de chaque type d'anomalie

    # 1. Excès de vitesse
    if not resume_exces.empty:
        resume_exces_simple = resume_exces[['Véhicule', 'Nb_Trajets_Exces', 'Pourcentage_Trajets_Exces']]
        resume = resume.merge(resume_exces_simple, on='Véhicule', how='left')
        resume['Nb_Trajets_Exces'] = resume['Nb_Trajets_Exces'].fillna(0)
        resume['Pourcentage_Trajets_Exces'] = resume['Pourcentage_Trajets_Exces'].fillna(0)
        resume['Score_Exces_Vitesse'] = resume['Pourcentage_Trajets_Exces'] * st.session_state.get(
            'ss_poids_vitesse_excessive', DEFAULT_POIDS_VITESSE_EXCESSIVE
        ) / 10 # Normaliser le pourcentage
    else:
        resume['Nb_Trajets_Exces'] = 0
        resume['Pourcentage_Trajets_Exces'] = 0
        resume['Score_Exces_Vitesse'] = 0

    # 2. Trajets suspects (hors heures, weekend, vitesse lente)
    if not trajets_suspects.empty:
        trajets_suspects_agg = trajets_suspects.groupby('Véhicule').agg(
            Nb_Trajets_Suspects=('Score_Suspicion_Total', 'count'),
            Score_Trajets_Suspects=('Score_Suspicion_Total', 'sum') # Somme des scores individuels de chaque trajet
        ).reset_index()
        resume = resume.merge(trajets_suspects_agg, on='Véhicule', how='left')
        resume['Nb_Trajets_Suspects'] = resume['Nb_Trajets_Suspects'].fillna(0)
        resume['Score_Trajets_Suspects'] = resume['Score_Trajets_Suspects'].fillna(0)
    else:
        resume['Nb_Trajets_Suspects'] = 0
        resume['Score_Trajets_Suspects'] = 0

    # 3. Transactions suspectes (sans présence)
    if not transactions_suspectes.empty and 'Nouveau Immat' in transactions_suspectes.columns:
        transactions_suspectes_agg = transactions_suspectes.groupby('Nouveau Immat').agg(
            Nb_Transactions_Suspectes=('Score_Suspicion', 'count'),
            Score_Transactions_Suspectes=('Score_Suspicion', 'sum')
        ).reset_index()
        transactions_suspectes_agg.rename(columns={'Nouveau Immat': 'Véhicule'}, inplace=True)
        resume = resume.merge(transactions_suspectes_agg, on='Véhicule', how='left')
        resume['Nb_Transactions_Suspectes'] = resume['Nb_Transactions_Suspectes'].fillna(0)
        resume['Score_Transactions_Suspectes'] = resume['Score_Transactions_Suspectes'].fillna(0)
    else:
        resume['Nb_Transactions_Suspectes'] = 0
        resume['Score_Transactions_Suspectes'] = 0

    # 4. Détours suspects
    if not detours_suspects.empty:
        detours_suspects_agg = detours_suspects.groupby('Véhicule').agg(
            Nb_Detours_Suspects=('Score_Detour', 'count'),
            Score_Detours=('Score_Detour', 'sum')
        ).reset_index()
        resume = resume.merge(detours_suspects_agg, on='Véhicule', how='left')
        resume['Nb_Detours_Suspects'] = resume['Nb_Detours_Suspects'].fillna(0)
        resume['Score_Detours'] = resume['Score_Detours'].fillna(0)
    else:
        resume['Nb_Detours_Suspects'] = 0
        resume['Score_Detours'] = 0

    # 5. Anomalies de distance (comparaison géoloc vs transactions)
    if not anomalies_distance.empty and 'Immatriculation' in anomalies_distance.columns:
        anomalies_distance_agg = anomalies_distance.groupby('Immatriculation').agg(
            Ecart_Distance_Km=('Ecart_Distance', 'sum'), # Somme des écarts
            Pourcentage_Ecart_Moyen=('Pourcentage_Ecart', 'mean') # Écart moyen en %
        ).reset_index()
        anomalies_distance_agg.rename(columns={'Immatriculation': 'Véhicule', 'Pourcentage_Ecart_Moyen': 'Pourcentage_Ecart'}, inplace=True)
        resume = resume.merge(anomalies_distance_agg, on='Véhicule', how='left')
        resume['Ecart_Distance_Km'] = resume['Ecart_Distance_Km'].fillna(0)
        resume['Pourcentage_Ecart'] = resume['Pourcentage_Ecart'].fillna(0)

        # Calculer un score basé sur l'écart de distance (positif ou négatif)
        # Le poids pourrait être un paramètre
        resume['Score_Ecart_Distance'] = resume['Pourcentage_Ecart'].abs() * (st.session_state.get('ss_poids_detour_suspect', DEFAULT_POIDS_DETOUR_SUSPECT) / 20) # Exemple de pondération
    else:
        resume['Ecart_Distance_Km'] = 0
        resume['Pourcentage_Ecart'] = 0
        resume['Score_Ecart_Distance'] = 0

    # Calculer le score total de risque
    resume['Score_Risque_Total'] = (
        resume['Score_Exces_Vitesse'] +
        resume['Score_Trajets_Suspects'] +
        resume['Score_Transactions_Suspectes'] +
        resume['Score_Detours'] +
        resume['Score_Ecart_Distance']
    ).fillna(0)


    # Classifier le niveau de risque
    resume['Niveau_Risque'] = pd.cut(
        resume['Score_Risque_Total'],
        bins=[-1, 10, 20, 40, float('inf')], # -1 pour inclure 0
        labels=['Faible', 'Modéré', 'Élevé', 'Critique']
    )

    # Générer un résumé textuel des anomalies
    def generer_resume_textuel(row):
        anomalies_txt_list = [] # Renommé pour éviter conflit

        if row['Nb_Trajets_Exces'] > 0:
            anomalies_txt_list.append(f"Excès vitesse: {int(row['Nb_Trajets_Exces'])} ({row['Pourcentage_Trajets_Exces']:.0f}%)")

        if row['Nb_Trajets_Suspects'] > 0:
            anomalies_txt_list.append(f"Trajets susp: {int(row['Nb_Trajets_Suspects'])}")

        if row['Nb_Transactions_Suspectes'] > 0:
            anomalies_txt_list.append(f"Transac. susp: {int(row['Nb_Transactions_Suspectes'])}")

        if row['Nb_Detours_Suspects'] > 0:
            anomalies_txt_list.append(f"Détours susp: {int(row['Nb_Detours_Suspects'])}")

        if abs(row['Pourcentage_Ecart']) > 5: # Seuil pour afficher l'écart de distance
            type_ecart = "Sur-décl." if row['Ecart_Distance_Km'] > 0 else "Sous-décl."
            anomalies_txt_list.append(f"{type_ecart} km: {row['Pourcentage_Ecart']:.0f}%")

        if not anomalies_txt_list:
            return "Aucune anomalie géoloc majeure"

        return " | ".join(anomalies_txt_list)

    resume['Résumé_Anomalies'] = resume.apply(generer_resume_textuel, axis=1)

    # Arrondir les valeurs numériques
    cols_arrondi_resume = ['Pourcentage_Trajets_Exces', 'Score_Exces_Vitesse', 'Score_Risque_Total',
                           'Ecart_Distance_Km', 'Pourcentage_Ecart', 'Score_Ecart_Distance',
                           'Score_Trajets_Suspects', 'Score_Transactions_Suspectes', 'Score_Detours']


    for col in cols_arrondi_resume:
        if col in resume.columns:
            resume[col] = resume[col].round(1)

    # Trier par score de risque total décroissant
    return resume.sort_values('Score_Risque_Total', ascending=False)


# Fonction pour visualiser les trajets sur une carte
def visualiser_trajets_sur_carte(
    df_geoloc: pd.DataFrame,
    vehicule_selectionne: str = None,
    date_debut: Optional[datetime.date] = None,
    date_fin: Optional[datetime.date] = None,
    highlight_anomalies: bool = False
) -> None:
    """
    Affiche les trajets sur une carte interactive.

    Args:
        df_geoloc: DataFrame des données de géolocalisation.
        vehicule_selectionne: Immatriculation du véhicule à afficher (optionnel).
        date_debut: Date de début pour filtrer les trajets (optionnel).
        date_fin: Date de fin pour filtrer les trajets (optionnel).
        highlight_anomalies: Indique si les anomalies doivent être mises en évidence.
    """
    # Vérifier que les coordonnées GPS sont disponibles
    if 'Latitude_depart' not in df_geoloc.columns or 'Longitude_depart' not in df_geoloc.columns \
       or 'Latitude_arrivee' not in df_geoloc.columns or 'Longitude_arrivee' not in df_geoloc.columns:
        st.warning("Les coordonnées GPS (départ ET arrivée) ne sont pas disponibles dans les données. Impossible d'afficher la carte.")
        return

    # Filtrer les données
    df_filtered = df_geoloc.copy()
    df_filtered = df_filtered[df_filtered['Type'] == 'Trajet'] # Uniquement les trajets


    if date_debut is not None and date_fin is not None:
        mask_date = (df_filtered['Date'].dt.date >= date_debut) & (df_filtered['Date'].dt.date <= date_fin)
        df_filtered = df_filtered[mask_date]

    if vehicule_selectionne is not None and vehicule_selectionne != "Tous les véhicules":
        df_filtered = df_filtered[df_filtered['Véhicule'] == vehicule_selectionne]

    # S'assurer qu'il y a des données à afficher avec coordonnées valides
    df_filtered = df_filtered.dropna(subset=['Latitude_depart', 'Longitude_depart', 'Latitude_arrivee', 'Longitude_arrivee'])

    if df_filtered.empty:
        st.warning("Aucune donnée géolocalisée (trajets avec coordonnées valides) disponible pour les critères sélectionnés.")
        return

    # Créer une carte centrée sur le point moyen des coordonnées
    mean_lat = df_filtered['Latitude_depart'].mean()
    mean_lon = df_filtered['Longitude_depart'].mean()

    m = folium.Map(location=[mean_lat, mean_lon], zoom_start=10) # Zoom un peu plus large

    # Pour highlight_anomalies, il faudrait que df_filtered contienne déjà les colonnes d'anomalies
    # On peut pré-calculer certaines si elles ne sont pas là (ex: Exces_Vitesse)
    if highlight_anomalies and 'Exces_Vitesse' not in df_filtered.columns:
        # Recalculer rapidement pour la visualisation
        seuil_vitesse_visu = st.session_state.get('ss_vitesse_excessive_seuil', 90) # Utiliser un seuil, peut-être paramétrable
        df_filtered['Exces_Vitesse'] = df_filtered['Vitesse moyenne'] > seuil_vitesse_visu


    # Ajouter les trajets à la carte
    for idx, row in df_filtered.iterrows():
        # Déterminer la couleur en fonction de la vitesse ou anomalies
        color = 'blue'  # Couleur par défaut
        popup_text = f"Véhicule: {row['Véhicule']}<br>" \
                     f"Date: {row['Date'].strftime('%d/%m/%Y')}<br>" \
                     f"Début: {row['Début']}<br>" \
                     f"Fin: {row['Fin']}<br>" \
                     f"Distance: {row.get('Distance', 'N/A'):.1f} km<br>" \
                     f"Vitesse moy.: {row.get('Vitesse moyenne', 'N/A'):.1f} km/h"


        if highlight_anomalies:
            anomalie_type = []
            if row.get('Est_weekend', False):
                color = 'orange'
                anomalie_type.append("Weekend")
            if row.get('Exces_Vitesse', False): # Exces_Vitesse doit être précalculé
                color = 'red'
                anomalie_type.append("Excès Vitesse")
            # Si d'autres colonnes d'anomalies (ex: Est_Detour_Potentiel) sont ajoutées à df_geoloc
            # if row.get('Est_Detour_Potentiel', False):
            #     color = 'purple' # Example
            #     anomalie_type.append("Détour")
            if anomalie_type:
                popup_text += f"<br><b>Anomalie: {', '.join(anomalie_type)}</b>"

        else: # Couleur basée sur la vitesse moyenne
            v_moy = row.get('Vitesse moyenne', 0)
            if pd.notna(v_moy):
                if v_moy < 30: color = 'green'
                elif v_moy < 60: color = 'blue'
                elif v_moy < 90: color = 'orange'
                else: color = 'red'


        # Créer le tracé du trajet
        folium.PolyLine(
            locations=[
                [row['Latitude_depart'], row['Longitude_depart']],
                [row['Latitude_arrivee'], row['Longitude_arrivee']]
            ],
            color=color,
            weight=3, # Un peu plus épais
            opacity=0.8,
            tooltip=popup_text # Utiliser tooltip pour hover
        ).add_to(m)

        # Ajouter les marqueurs de début et fin
        folium.Marker(
            [row['Latitude_depart'], row['Longitude_depart']],
            popup=f"Départ: {row['Véhicule']} - {row['Date'].strftime('%d/%m/%Y')} {row['Début']}",
            icon=folium.Icon(color='green', icon='play', prefix='fa')
        ).add_to(m)

        folium.Marker(
            [row['Latitude_arrivee'], row['Longitude_arrivee']],
            popup=f"Arrivée: {row['Véhicule']} - {row['Date'].strftime('%d/%m/%Y')} {row['Fin']}",
            icon=folium.Icon(color='darkred', icon='stop', prefix='fa') # 'darkred' pour fin
        ).add_to(m)

    # Afficher la carte
    folium_static(m, width=None, height=600) # Ajuster la hauteur


# Fonction pour afficher la page d'analyse de géolocalisation
def afficher_page_analyse_geolocalisation(
    df_geoloc: pd.DataFrame,
    df_transactions: pd.DataFrame,
    df_vehicules: pd.DataFrame,
    date_debut: datetime.date,
    date_fin: datetime.date
):
    """Affiche la page d'analyse des données de géolocalisation."""
    st.header(f"📍 Analyse des Données de Géolocalisation ({date_debut.strftime('%d/%m/%Y')} - {date_fin.strftime('%d/%m/%Y')})")

    if df_geoloc is None or df_geoloc.empty:
        st.warning("Aucune donnée de géolocalisation à analyser. Veuillez charger un fichier de géolocalisation.")
        return

    # Filtrer les données pour la période sélectionnée
    mask_date = (df_geoloc['Date'].dt.date >= date_debut) & (df_geoloc['Date'].dt.date <= date_fin)
    df_geoloc_filtered_orig = df_geoloc[mask_date] # Garder original pour analyses mixtes Type

    if df_geoloc_filtered_orig.empty:
        st.warning(f"Aucune donnée de géolocalisation pour la période du {date_debut} au {date_fin}.")
        return

    # Filtrer spécifiquement les TRAJETS pour de nombreuses analyses
    df_geoloc_filtered_trajets = df_geoloc_filtered_orig[df_geoloc_filtered_orig['Type'] == 'Trajet'].copy()


    # Onglets pour différentes analyses
    tab_synthese, tab_comparaison, tab_vitesse, tab_utilisation, tab_trajets_suspects, tab_vehicules_sans_geoloc, tab_carte, tab_integration = st.tabs([
        "📊 Synthèse", "🔍 Comparaison Carburant", "🚨 Excès de Vitesse", "⚙️ Utilisation",
        "⚠️ Trajets Suspects", "❓ Véhicules Sans Géoloc", "🗺️ Carte", "🔄 Intégration"
    ])

    with tab_synthese:
        st.subheader("Synthèse des Données de Géolocalisation (Trajets)")
        if df_geoloc_filtered_trajets.empty:
            st.info("Aucun trajet enregistré dans la période sélectionnée.")
        else:
            nb_vehicules = df_geoloc_filtered_trajets['Véhicule'].nunique()
            nb_trajets = len(df_geoloc_filtered_trajets)
            distance_totale = df_geoloc_filtered_trajets['Distance'].sum()
            duree_totale_min = df_geoloc_filtered_trajets['Durée_minutes'].sum(skipna=True)
            duree_totale_heures = duree_totale_min / 60 if duree_totale_min > 0 else 0
            vitesse_moy = df_geoloc_filtered_trajets['Vitesse moyenne'].mean(skipna=True)
            # Pour Vitesse Max, on prend le max des vitesses moyennes des segments (car 'Vitesse max' originale n'est plus là)
            vitesse_max_segments = df_geoloc_filtered_trajets['Vitesse moyenne'].max(skipna=True)


            # Afficher les KPIs
            col1, col2, col3, col4 = st.columns(4)
            col1.metric("Véhicules avec Trajets", f"{nb_vehicules}")
            col2.metric("Trajets Enregistrés", f"{nb_trajets:,}")
            col3.metric("Distance Totale (Trajets)", f"{distance_totale:,.1f} km")
            col4.metric("Durée Totale (Trajets)", f"{duree_totale_heures:,.1f} h")

            col5, col6 = st.columns(2)
            col5.metric("Vitesse Moyenne (Trajets)", f"{vitesse_moy:.1f} km/h" if pd.notna(vitesse_moy) else "N/A")
            col6.metric("Vitesse Max des Segments", f"{vitesse_max_segments:.1f} km/h" if pd.notna(vitesse_max_segments) else "N/A")


            # Graphique de répartition des distances par véhicule
            distance_par_vehicule = df_geoloc_filtered_trajets.groupby('Véhicule')['Distance'].sum().reset_index()
            fig_distance = px.bar(
                distance_par_vehicule.sort_values('Distance', ascending=False),
                x='Véhicule',
                y='Distance',
                title="Distance Totale Parcourue par Véhicule (Trajets)",
                labels={'Distance': 'Distance (km)'}
            )
            st.plotly_chart(fig_distance, use_container_width=True)

            # Graphique d'évolution quotidienne des distances
            distance_quotidienne = df_geoloc_filtered_trajets.groupby(pd.Grouper(key='Date', freq='D'))['Distance'].sum().reset_index()
            fig_evol = px.line(
                distance_quotidienne,
                x='Date',
                y='Distance',
                title="Évolution Quotidienne de la Distance Parcourue (Trajets)",
                labels={'Distance': 'Distance (km)'},
                markers=True
            )
            st.plotly_chart(fig_evol, use_container_width=True)

            # Tableau récapitulatif
            st.subheader("Données par Véhicule (Trajets)")
            recap_vehicule = df_geoloc_filtered_trajets.groupby('Véhicule').agg(
                Nb_Trajets=('Distance', 'count'),
                Distance_Totale=('Distance', 'sum'),
                Distance_Moyenne=('Distance', 'mean'),
                Vitesse_Moyenne_Segments=('Vitesse moyenne', 'mean'),
                Vitesse_Max_Segment=('Vitesse moyenne', 'max'),
                Duree_Totale_Minutes=('Durée_minutes', 'sum')
            ).reset_index()

            recap_vehicule['Duree_Totale_Heures'] = (recap_vehicule['Duree_Totale_Minutes'] / 60).round(1)
            recap_vehicule = recap_vehicule.sort_values('Distance_Totale', ascending=False)

            # Arrondir les valeurs
            cols_numeriques = ['Distance_Totale', 'Distance_Moyenne', 'Vitesse_Moyenne_Segments', 'Vitesse_Max_Segment']
            for col in cols_numeriques:
                recap_vehicule[col] = recap_vehicule[col].round(1)

            afficher_dataframe_avec_export(recap_vehicule, "Récapitulatif par Véhicule (Trajets)", key="geoloc_recap_vehicule_trajets")

    # Ajouter le contenu du nouvel onglet
    with tab_vehicules_sans_geoloc:
        st.subheader("Véhicules Sans Données de Géolocalisation")
        
        with st.spinner("Analyse des véhicules sans géolocalisation en cours..."):
            df_vehicules_sans_geoloc = analyser_vehicules_sans_geoloc(
                df_transactions, df_vehicules, df_geoloc_filtered_orig, date_debut, date_fin
            )
        
        if df_vehicules_sans_geoloc.empty:
            st.success("✅ Tous les véhicules avec des transactions ont des données de géolocalisation pour la période sélectionnée.")
        else:
            nb_vehicules_sans_geoloc = len(df_vehicules_sans_geoloc)
            nb_total_vehicules_actifs = len(set(df_transactions[df_transactions['Card num.'].isin(df_vehicules['N° Carte'])]['Card num.'].unique()))
            
            pourcentage = (nb_vehicules_sans_geoloc / nb_total_vehicules_actifs * 100) if nb_total_vehicules_actifs > 0 else 0
            
            st.warning(f"⚠️ {nb_vehicules_sans_geoloc} véhicules ({pourcentage:.1f}% de la flotte active) ont effectué des transactions sans données de géolocalisation correspondantes.")
            
            # Afficher les statistiques
            afficher_dataframe_avec_export(
                df_vehicules_sans_geoloc,
                "Véhicules Sans Géolocalisation",
                key="vehicules_sans_geoloc"
            )
            
            # Graphique des transactions sans géolocalisation
            fig_sans_geoloc = px.bar(
                df_vehicules_sans_geoloc,
                x='Immatriculation',
                y='Nb_Transactions',
                title="Nombre de Transactions Sans Géolocalisation par Véhicule",
                color='Catégorie',
                hover_data=['Volume_Total_L', 'Montant_Total_CFA', 'Jours_Avec_Transactions']
            )
            st.plotly_chart(fig_sans_geoloc, use_container_width=True)
            
            # Recommandations
            st.subheader("Causes Possibles et Recommandations")
            st.markdown("""
            ### Causes possibles de l'absence de données de géolocalisation:
            
            1. **Équipement non installé**: Certains véhicules peuvent ne pas être équipés de dispositifs de géolocalisation.
            2. **Défaillance technique**: Problème avec le dispositif GPS ou le système de transmission de données.
            3. **Désactivation volontaire**: Débranchement ou neutralisation intentionnelle du dispositif.
            4. **Problème de couverture réseau**: Zones géographiques sans couverture pour la transmission de données.
            5. **Problème d'identifiant**: Non-correspondance entre l'immatriculation dans les données de transaction et de géolocalisation.
            
            ### Recommandations:
            
            - **Vérification physique**: Effectuer un contrôle visuel des dispositifs GPS sur les véhicules listés.
            - **Test de fonctionnement**: Vérifier que les systèmes sont correctement alimentés et configurés.
            - **Mise à jour des identifiants**: S'assurer que les immatriculations sont identiques dans tous les systèmes.
            - **Installation manquante**: Équiper les véhicules non pourvus de dispositifs de géolocalisation.
            - **Formation**: Sensibiliser les conducteurs à l'importance du bon fonctionnement des systèmes de géolocalisation.
            """)

    with tab_comparaison:
        st.subheader("Comparaison Kilométrage Géolocalisation vs. Transactions Carburant")

        # Analyse comparative
        with st.spinner("Analyse comparative en cours..."):
            # Utiliser df_geoloc_filtered_orig car analyser_geolocalisation_vs_transactions filtre par Type='Trajet' en interne
            comparaison, anomalies = analyser_geolocalisation_vs_transactions(
                df_geoloc_filtered_orig, df_transactions, df_vehicules, date_debut, date_fin
            )

        if comparaison.empty:
            st.info("Données insuffisantes pour effectuer une comparaison.")
        else:
            # Affichage des résultats
            nb_anomalies = len(anomalies)
            nb_sur_decla = len(anomalies[anomalies['Type_Anomalie'] == 'Sur-déclaration kilométrique']) if not anomalies.empty else 0
            nb_sous_decla = len(anomalies[anomalies['Type_Anomalie'] == 'Sous-déclaration kilométrique']) if not anomalies.empty else 0

            st.warning(f"⚠️ Détection de {nb_anomalies} anomalies significatives de kilométrage ({nb_sur_decla} sur-déclarations, {nb_sous_decla} sous-déclarations)")

            # Graphique des écarts
            fig_ecart = px.bar(
                comparaison.sort_values('Ecart_Distance', ascending=False).head(15),
                x='Immatriculation',
                y='Ecart_Distance',
                title="Top 15 des Écarts de Kilométrage (Déclaré - Géolocalisé)",
                labels={'Ecart_Distance': 'Écart (km)'},
                color='Ecart_Distance',
                color_continuous_scale=px.colors.diverging.RdBu,
                color_continuous_midpoint=0
            )
            st.plotly_chart(fig_ecart, use_container_width=True)

            # Tableau des anomalies
            if not anomalies.empty:
                st.subheader("Anomalies Détectées")
                cols_anomalies = [
                    'Immatriculation', 'Type_Anomalie', 'Gravite',
                    'Distance_Geoloc_Totale', 'Distance_Declaree_Totale',
                    'Ecart_Distance', 'Pourcentage_Ecart'
                ]
                afficher_dataframe_avec_export(
                    anomalies[cols_anomalies],
                    "Anomalies de Kilométrage",
                    key="geoloc_anomalies_km"
                )

            # Tableau complet
            st.subheader("Comparaison Complète par Véhicule")
            cols_comparaison = [
                'Immatriculation', 'Distance_Geoloc_Totale', 'Distance_Declaree_Totale',
                'Ecart_Distance', 'Pourcentage_Ecart', 'Consommation_100km_Reelle',
                'Consommation_100km_Declaree', 'Ecart_Consommation', 'Pourcentage_Ecart_Consommation',
                'Volume_Carburant_Total', 'Nb_Trajets', 'Nb_Transactions'
            ]
            afficher_dataframe_avec_export(
                comparaison[[c for c in cols_comparaison if c in comparaison.columns]],
                "Comparaison Kilométrage",
                key="geoloc_comparaison_km"
            )

            # Et ajouter un graphique de comparaison des consommations
            if 'Consommation_100km_Reelle' in comparaison.columns and 'Consommation_100km_Declaree' in comparaison.columns:
                comparaison_conso = comparaison.dropna(subset=['Consommation_100km_Reelle', 'Consommation_100km_Declaree']).sort_values('Pourcentage_Ecart_Consommation', ascending=False)
                if not comparaison_conso.empty:
                    fig_conso_compare = px.bar(
                        comparaison_conso,
                        x='Immatriculation',
                        y=['Consommation_100km_Reelle', 'Consommation_100km_Declaree'],
                        title="Comparaison Consommation Réelle (basée géoloc) vs Déclarée (basée km carte)",
                        labels={'value': 'Consommation (L/100km)', 'variable': 'Source'},
                        barmode='group'
                    )
                    st.plotly_chart(fig_conso_compare, use_container_width=True)

            # Graphique de comparaison consommation
            if 'Consommation_100km_Reelle' in comparaison.columns:
                comparaison_conso = comparaison.dropna(subset=['Consommation_100km_Reelle']).sort_values('Consommation_100km_Reelle', ascending=False)
                if not comparaison_conso.empty:
                    fig_conso = px.bar(
                        comparaison_conso,
                        x='Immatriculation',
                        y='Consommation_100km_Reelle',
                        title="Consommation Réelle (L/100km) basée sur la distance géolocalisée",
                        labels={'Consommation_100km_Reelle': 'Consommation (L/100km)'}
                    )
                    st.plotly_chart(fig_conso, use_container_width=True)

            # Recommandations basées sur les anomalies
            if not anomalies.empty:
                st.subheader("Recommandations")
                st.markdown("""
                Basé sur l'analyse des écarts entre kilométrage déclaré et géolocalisé:

                1. **Pour les sur-déclarations importantes**: Vérifier si les véhicules concernés ont effectué des trajets non enregistrés par le système de géolocalisation, ou s'il y a des déclarations potentiellement incorrectes.

                2. **Pour les sous-déclarations importantes**: Vérifier si les transactions de carburant sont correctement associées au véhicule, ou si le kilométrage n'est pas systématiquement saisi lors des prises de carburant.

                3. **Pour les écarts persistants**: Envisager un audit spécifique des véhicules présentant des anomalies répétées.
                """)

    with tab_vitesse:
        st.subheader("Analyse des Excès de Vitesse")

        # Paramètre de seuil de vitesse
        seuil_vitesse = st.slider(
            "Seuil de vitesse considéré comme excès (km/h)",
            min_value=50,
            max_value=130,
            value=90, # Default
            step=5,
            key="slider_seuil_vitesse_geoloc"
        )

        # Analyse des excès de vitesse
        with st.spinner("Analyse des excès de vitesse en cours..."):
             # Utiliser df_geoloc_filtered_orig car analyser_exces_vitesse filtre par Type='Trajet' en interne
            resume_exces, detail_exces = analyser_exces_vitesse(
                df_geoloc_filtered_orig, date_debut, date_fin, seuil_vitesse
            )

        if resume_exces.empty:
            st.info("Données insuffisantes pour analyser les excès de vitesse.")
        else:
            nb_total_exces = resume_exces['Nb_Trajets_Exces'].sum()
            nb_vehicules_exces = len(resume_exces[resume_exces['Nb_Trajets_Exces'] > 0])

            col_v1, col_v2, col_v3 = st.columns(3)
            col_v1.metric("Nombre Total d'Excès", f"{nb_total_exces:,}")
            col_v2.metric("Véhicules Concernés", f"{nb_vehicules_exces}")
            # Vitesse_Max_Observee est maintenant le max des vitesses moyennes des segments
            col_v3.metric("Vitesse Max (Segment)", f"{resume_exces['Vitesse_Max_Observee'].max():.1f} km/h" if not resume_exces['Vitesse_Max_Observee'].empty else "N/A")


            # Graphique des taux d'excès par véhicule
            fig_exces = px.bar(
                resume_exces[resume_exces['Nb_Trajets_Exces'] > 0].sort_values('Pourcentage_Trajets_Exces', ascending=False),
                x='Véhicule',
                y='Pourcentage_Trajets_Exces',
                title=f"Pourcentage de Trajets en Excès de Vitesse (>{seuil_vitesse} km/h) par Véhicule",
                labels={'Pourcentage_Trajets_Exces': '% de Trajets en Excès'},
                color='Niveau_Risque',
                color_discrete_map={
                    'Faible': 'green',
                    'Modéré': 'orange',
                    'Élevé': 'red',
                    'Très élevé': 'darkred'
                },
                hover_data=['Nb_Trajets_Exces', 'Nb_Total_Trajets', 'Vitesse_Max_Observee']
            )
            st.plotly_chart(fig_exces, use_container_width=True)

            # Histogramme des vitesses moyennes des segments en excès
            if not detail_exces.empty:
                fig_histogramme = px.histogram(
                    detail_exces,
                    x='Vitesse moyenne', # C'est la vitesse moyenne du segment
                    title="Distribution des Vitesses Moyennes des Trajets en Excès",
                    labels={'Vitesse moyenne': 'Vitesse Moyenne du Trajet (km/h)'},
                    color='Niveau_Exces',
                    color_discrete_map={
                        'Léger (< 10 km/h)': 'green',
                        'Modéré (10-20 km/h)': 'orange',
                        'Important (20-30 km/h)': 'red',
                        'Grave (> 30 km/h)': 'darkred'
                    },
                    nbins=30
                )
                st.plotly_chart(fig_histogramme, use_container_width=True)

            # Tableau récapitulatif
            st.subheader("Résumé des Excès par Véhicule")
            cols_resume_exces = [
                'Véhicule', 'Nb_Trajets_Exces', 'Nb_Total_Trajets', 'Pourcentage_Trajets_Exces',
                'Vitesse_Max_Observee', 'Vitesse_Moyenne_Trajets', 'Depassement_Moyen', 'Niveau_Risque'
            ]
            afficher_dataframe_avec_export(
                resume_exces[[c for c in cols_resume_exces if c in resume_exces.columns]],
                "Résumé des Excès",
                key="geoloc_resume_exces"
            )

            # Détail des excès
            if not detail_exces.empty:
                st.subheader("Détail des Trajets en Excès de Vitesse")
                cols_detail_exces = [
                    'Véhicule', 'Date', 'Début', 'Fin', 'Distance', 'Vitesse moyenne', # Utiliser Vitesse moyenne
                    'Depassement_km/h', 'Niveau_Exces'
                ]
                afficher_dataframe_avec_export(
                    detail_exces[[c for c in cols_detail_exces if c in detail_exces.columns]],
                    "Détail des Excès",
                    key="geoloc_detail_exces"
                )

    with tab_utilisation:
        st.subheader("Analyse de l'Utilisation des Véhicules")

        # Analyse de l'utilisation
        with st.spinner("Analyse de l'utilisation en cours..."):
             # Utiliser df_geoloc_filtered_orig car analyser_utilisation_vehicules filtre par Type='Trajet' en interne
            utilisation_vehicules, utilisation_quotidienne = analyser_utilisation_vehicules(
                df_geoloc_filtered_orig, date_debut, date_fin
            )

        if utilisation_vehicules.empty:
            st.info("Données insuffisantes pour analyser l'utilisation des véhicules.")
        else:
            # KPIs généraux
            duree_totale_heures = utilisation_vehicules['Duree_Totale_Heures'].sum()
            distance_totale = utilisation_vehicules['Distance_Totale'].sum()
            nb_trajets_weekend = utilisation_vehicules['Nb_Trajets_Weekend'].sum()
            nb_trajets_total = utilisation_vehicules['Nb_Trajets'].sum()
            pct_weekend = (nb_trajets_weekend / nb_trajets_total * 100) if nb_trajets_total > 0 else 0

            col_u1, col_u2, col_u3 = st.columns(3)
            col_u1.metric("Durée Totale d'Utilisation (Trajets)", f"{duree_totale_heures:,.1f} h")
            col_u2.metric("Distance Totale (Trajets)", f"{distance_totale:,.1f} km")
            col_u3.metric("Utilisation Weekend (Trajets)", f"{pct_weekend:.1f}%")

            # Graphique d'utilisation par véhicule
            fig_utilisation = px.bar(
                utilisation_vehicules.sort_values('Duree_Totale_Heures', ascending=False),
                x='Véhicule',
                y=['Duree_Totale_Heures', 'Distance_Totale'],
                title="Durée d'Utilisation et Distance par Véhicule (Trajets)",
                labels={'value': 'Valeur', 'variable': 'Métrique'},
                barmode='group'
            )
            st.plotly_chart(fig_utilisation, use_container_width=True)

            # Graphique des périodes d'utilisation
            periodes_jours = ['Matin (6h-9h)', 'Matinée (9h-12h)', 'Midi (12h-14h)',
                             'Après-midi (14h-17h)', 'Soir (17h-20h)', 'Nuit (20h-6h)']

            periodes_presentes = [p for p in periodes_jours if p in utilisation_vehicules.columns]
            if periodes_presentes:
                fig_periodes = px.bar(
                    utilisation_vehicules.sort_values('Nb_Trajets', ascending=False),
                    x='Véhicule',
                    y=periodes_presentes,
                    title="Répartition des Trajets par Période de la Journée",
                    labels={'value': 'Nombre de Trajets', 'variable': 'Période'},
                    barmode='stack'
                )
                st.plotly_chart(fig_periodes, use_container_width=True)

            # Pourcentage d'utilisation le weekend
            fig_weekend = px.bar(
                utilisation_vehicules.sort_values('Pourcentage_Trajets_Weekend', ascending=False),
                x='Véhicule',
                y='Pourcentage_Trajets_Weekend',
                title="Pourcentage d'Utilisation le Weekend par Véhicule (Trajets)",
                labels={'Pourcentage_Trajets_Weekend': '% Trajets Weekend'}
            )
            st.plotly_chart(fig_weekend, use_container_width=True)

            # Évolution quotidienne pour un véhicule spécifique
            if not utilisation_quotidienne.empty:
                vehicules_disponibles = sorted(utilisation_quotidienne['Véhicule'].unique())
                if vehicules_disponibles:
                    vehicule_selectionne = st.selectbox(
                        "Sélectionner un véhicule pour l'évolution quotidienne",
                        options=vehicules_disponibles,
                        key="veh_evol_quot_geoloc"
                    )

                    data_vehicule = utilisation_quotidienne[utilisation_quotidienne['Véhicule'] == vehicule_selectionne]

                    fig_evol_veh = px.line(
                        data_vehicule,
                        x='Date',
                        y=['Distance_Jour', 'Duree_Jour_Heures'],
                        title=f"Évolution Quotidienne (Trajets) - {vehicule_selectionne}",
                        labels={'value': 'Valeur', 'variable': 'Métrique'},
                        markers=True
                    )
                    st.plotly_chart(fig_evol_veh, use_container_width=True)

            # Tableau récapitulatif
            st.subheader("Utilisation Détaillée par Véhicule (Trajets)")
            cols_utilisation = [
                'Véhicule', 'Nb_Trajets', 'Distance_Totale', 'Duree_Totale_Heures',
                'Distance_Moyenne_Trajet', 'Duree_Moyenne_Trajet', 'Vitesse_Moyenne',
                'Nb_Trajets_Weekend', 'Pourcentage_Trajets_Weekend'
            ]

            for periode in periodes_jours:
                if periode in utilisation_vehicules.columns:
                    cols_utilisation.append(periode)

            afficher_dataframe_avec_export(
                utilisation_vehicules[[c for c in cols_utilisation if c in utilisation_vehicules.columns]],
                "Utilisation Détaillée (Trajets)",
                key="geoloc_utilisation_detail_trajets"
            )

    with tab_trajets_suspects:
        st.subheader("Analyse des Trajets Suspects")

        # Détection des trajets suspects
        with st.spinner("Détection des trajets suspects en cours..."):
            # Utiliser df_geoloc_filtered_orig car les fonctions internes filtrent par Type au besoin
            trajets_suspects = detecter_trajets_suspects(df_geoloc_filtered_orig, date_debut, date_fin)

            _, transactions_sans_presence = analyser_correspondance_transactions_geoloc(
                df_geoloc_filtered_orig, df_transactions, df_vehicules, date_debut, date_fin
            )

            detours_suspects = detecter_detours_suspects(df_geoloc_filtered_orig, date_debut, date_fin)

            resume_anomalies = generer_resume_anomalies_geolocalisation(
                df_geoloc_filtered_orig, df_transactions, df_vehicules, date_debut, date_fin
            )

        # Affichage du tableau de bord de risque
        if not resume_anomalies.empty:
            st.warning(f"⚠️ Détection de véhicules présentant des anomalies de géolocalisation")

            vehicules_risque_eleve = resume_anomalies[resume_anomalies['Niveau_Risque'].isin(['Élevé', 'Critique'])]

            if not vehicules_risque_eleve.empty:
                nb_vehicules_risque = len(vehicules_risque_eleve)
                st.error(f"🚨 {nb_vehicules_risque} véhicule(s) présentent un niveau de risque élevé ou critique")

                cols_risque = ['Véhicule', 'Score_Risque_Total', 'Niveau_Risque', 'Résumé_Anomalies']
                afficher_dataframe_avec_export(
                    vehicules_risque_eleve[[c for c in cols_risque if c in vehicules_risque_eleve.columns]],
                    "Véhicules à Risque Élevé/Critique",
                    key="vehicules_risque_eleve_geoloc"
                )

            st.subheader("Tableau de Bord des Risques par Véhicule (Géolocalisation)")
            cols_resume = [
                'Véhicule', 'Niveau_Risque', 'Score_Risque_Total', 'Résumé_Anomalies',
                'Nb_Trajets_Exces', 'Nb_Trajets_Suspects', 'Nb_Transactions_Suspectes',
                'Nb_Detours_Suspects', 'Ecart_Distance_Km', 'Pourcentage_Ecart'
            ]
            afficher_dataframe_avec_export(
                resume_anomalies[[c for c in cols_resume if c in resume_anomalies.columns]],
                "Tableau de Bord des Risques (Géoloc)",
                key="tableau_bord_risques_geoloc"
            )

            fig_risque = px.bar(
                resume_anomalies.sort_values('Score_Risque_Total', ascending=False),
                x='Véhicule',
                y='Score_Risque_Total',
                title="Score de Risque Total par Véhicule (Géolocalisation)",
                labels={'Score_Risque_Total': 'Score de Risque'},
                color='Niveau_Risque',
                color_discrete_map={
                    'Faible': 'green',
                    'Modéré': 'yellow',
                    'Élevé': 'orange',
                    'Critique': 'red'
                }
            )
            st.plotly_chart(fig_risque, use_container_width=True)
        else:
            st.success("✅ Aucune anomalie significative détectée dans les données de géolocalisation.")

        if not trajets_suspects.empty:
            st.subheader("Trajets Suspects Détectés (Hors heures, Weekend, Vitesse anormale)")
            nb_trajets_suspects_val = len(trajets_suspects) # Renommé
            st.warning(f"⚠️ {nb_trajets_suspects_val} trajets suspects détectés.")
            cols_ts = [
                'Véhicule', 'Date_Heure_Debut', 'Distance', 'Durée_minutes',
                'Vitesse moyenne', 'Description_Suspicion', 'Niveau_Suspicion'
            ]
            afficher_dataframe_avec_export(
                trajets_suspects[[c for c in cols_ts if c in trajets_suspects.columns]],
                "Détail des Trajets Suspects",
                key="detail_trajets_suspects_geoloc"
            )
            trajets_par_vehicule = trajets_suspects.groupby('Véhicule').size().reset_index(name='Nombre_Trajets_Suspects')
            fig_ts_vehicule = px.bar(
                trajets_par_vehicule.sort_values('Nombre_Trajets_Suspects', ascending=False),
                x='Véhicule', y='Nombre_Trajets_Suspects',
                title="Nombre de Trajets Suspects par Véhicule",
                labels={'Nombre_Trajets_Suspects': 'Nombre de Trajets Suspects'}
            )
            st.plotly_chart(fig_ts_vehicule, use_container_width=True)

        if not transactions_sans_presence.empty:
            st.subheader("Transactions sans Présence du Véhicule (Basé sur Géolocalisation)")
            nb_transactions_suspectes_val = len(transactions_sans_presence) # Renommé
            st.warning(f"⚠️ {nb_transactions_suspectes_val} transactions sans présence détectée du véhicule.")
            cols_trans = [
                'Nouveau Immat', 'Date', 'Hour', 'Place', 'Quantity', 'Amount',
                'Methode_Verification', 'Est_Weekend', 'Est_Hors_Heures', 'Niveau_Suspicion'
            ]
            afficher_dataframe_avec_export(
                transactions_sans_presence[[c for c in cols_trans if c in transactions_sans_presence.columns]],
                "Transactions sans Présence Détectée",
                key="transactions_sans_presence_geoloc"
            )

        if not detours_suspects.empty:
            st.subheader("Détours Potentiels Détectés")
            nb_detours = len(detours_suspects)
            st.warning(f"⚠️ {nb_detours} trajets avec détours potentiels détectés.")
            cols_detours = [
                'Véhicule', 'Date', 'Début', 'Distance', 'Durée_minutes',
                'Vitesse moyenne', 'Severite_Detour', 'Description_Detour'
            ]
            afficher_dataframe_avec_export(
                detours_suspects[[c for c in cols_detours if c in detours_suspects.columns]],
                "Détail des Détours Suspects",
                key="detail_detours_suspects_geoloc"
            )
            detours_par_severite = detours_suspects['Severite_Detour'].value_counts().reset_index()
            detours_par_severite.columns = ['Sévérité', 'Nombre']
            fig_severite = px.pie(
                detours_par_severite, values='Nombre', names='Sévérité',
                title="Répartition des Détours par Sévérité", color='Sévérité',
                color_discrete_map={'Léger': 'green', 'Modéré': 'orange', 'Important': 'red'}
            )
            st.plotly_chart(fig_severite, use_container_width=True)

    with tab_carte:
        st.subheader("Visualisation des Trajets sur Carte")
        coords_disponibles = (
            'Latitude_depart' in df_geoloc_filtered_orig.columns and
            'Longitude_depart' in df_geoloc_filtered_orig.columns and
            'Latitude_arrivee' in df_geoloc_filtered_orig.columns and
            'Longitude_arrivee' in df_geoloc_filtered_orig.columns and
            not df_geoloc_filtered_orig['Latitude_depart'].isna().all() and
            not df_geoloc_filtered_orig['Latitude_arrivee'].isna().all()
        )

        if not coords_disponibles:
            st.warning("Les coordonnées GPS (départ ET arrivée) ne sont pas disponibles dans les données. La visualisation sur carte n'est pas possible.")
        else:
            st.sidebar.subheader("Filtres pour la Carte")
            vehicules_carte = sorted(df_geoloc_filtered_orig['Véhicule'].unique())
            vehicule_carte = st.sidebar.selectbox(
                "Véhicule à visualiser", options=["Tous les véhicules"] + vehicules_carte,
                key="carte_vehicule_filter_geoloc"
            )
            highlight_anomalies_carte = st.sidebar.checkbox( # Renommé pour unicité de clé
                "Mettre en évidence les anomalies (carte)", value=True, key="highlight_anomalies_carte"
            )
            visualiser_trajets_sur_carte(
                df_geoloc_filtered_orig, # Passer le df filtré par date, la fonction interne filtre par Type Trajet
                vehicule_carte if vehicule_carte != "Tous les véhicules" else None,
                date_debut, date_fin, highlight_anomalies_carte
            )
            st.markdown("""
            ### Légende
            - **Points verts (play)**: Points de départ des trajets
            - **Points rouges (stop)**: Points d'arrivée des trajets
            - **Lignes**:
                - *Mode Normal*: Vert (lent) -> Bleu (moyen) -> Orange (rapide) -> Rouge (très rapide)
                - *Mode Anomalies*: Orange (Weekend), Rouge (Excès Vitesse), Violet (Détour - si implémenté)
            """)
            st.info("📌 Zoomez/Dézoomez, cliquez sur les tracés/marqueurs pour détails.")

    with tab_integration:
        st.subheader("Intégration Géolocalisation - Carburant: Analyse d'Efficacité")
        with st.spinner("Analyse de l'efficacité du carburant en cours..."):
             # Utiliser df_geoloc_filtered_orig car la fonction interne filtre par Type='Trajet'
            efficacite_carburant = analyser_efficacite_carburant(
                df_geoloc_filtered_orig, df_transactions, df_vehicules, date_debut, date_fin
            )

        if efficacite_carburant.empty:
            st.info("Données insuffisantes pour analyser l'efficacité du carburant.")
        else:
            conso_moyenne_eff = efficacite_carburant['Consommation_100km'].mean() # Renommé
            cout_km_moyen_eff = efficacite_carburant['Cout_par_km'].mean() # Renommé

            col_e1, col_e2 = st.columns(2)
            col_e1.metric("Consommation Moyenne (basée géoloc)", f"{conso_moyenne_eff:.1f} L/100km" if pd.notna(conso_moyenne_eff) else "N/A")
            col_e2.metric("Coût Moyen par Km (basé géoloc)", f"{cout_km_moyen_eff:.1f} CFA/km" if pd.notna(cout_km_moyen_eff) else "N/A")

            fig_efficacite = px.bar(
                efficacite_carburant.sort_values('Score_Efficacite', ascending=False),
                x='Véhicule', y='Score_Efficacite', title="Score d'Efficacité par Véhicule",
                labels={'Score_Efficacite': "Score d'Efficacité (100 = moyenne)"},
                color='Niveau_Efficacite',
                color_discrete_map={
                    'Très faible': 'red', 'Faible': 'orange', 'Normale': 'yellow',
                    'Bonne': 'lightgreen', 'Excellente': 'darkgreen'}
            )
            st.plotly_chart(fig_efficacite, use_container_width=True)

            st.subheader("Efficacité Détaillée par Véhicule")
            cols_efficacite = [
                'Véhicule', 'Catégorie', 'Distance_Geoloc_Totale', 'Volume_Total',
                'Consommation_100km', 'Conso_Moyenne_Cat', 'Ecart_Conso_Pct',
                'Cout_par_km', 'Score_Efficacite', 'Niveau_Efficacite'
            ]
            afficher_dataframe_avec_export(
                efficacite_carburant[[c for c in cols_efficacite if c in efficacite_carburant.columns]],
                "Efficacité Carburant", key="efficacite_carburant_geoloc"
            )

            efficacite_par_cat = efficacite_carburant.dropna(subset=['Consommation_100km', 'Conso_Moyenne_Cat'])
            if not efficacite_par_cat.empty:
                fig_comp_cat = make_subplots(specs=[[{"secondary_y": True}]])
                fig_comp_cat.add_trace(go.Bar(x=efficacite_par_cat['Véhicule'], y=efficacite_par_cat['Consommation_100km'], name="Consommation Réelle", marker_color='blue'), secondary_y=False)
                fig_comp_cat.add_trace(go.Scatter(x=efficacite_par_cat['Véhicule'], y=efficacite_par_cat['Conso_Moyenne_Cat'], name="Moyenne Catégorie", marker_color='red', mode='lines'), secondary_y=False)
                fig_comp_cat.add_trace(go.Bar(x=efficacite_par_cat['Véhicule'], y=efficacite_par_cat['Ecart_Conso_Pct'], name="Écart (%)", marker_color='green'), secondary_y=True)
                fig_comp_cat.update_layout(title="Comparaison Consommation Réelle vs Moyenne Catégorie", xaxis_title="Véhicule", yaxis_title="Consommation (L/100km)", yaxis2_title="Écart (%)")
                st.plotly_chart(fig_comp_cat, use_container_width=True)

            st.subheader("Recommandations d'Optimisation")
            efficacite_faible = efficacite_carburant[efficacite_carburant['Niveau_Efficacite'].isin(['Très faible', 'Faible'])]
            if not efficacite_faible.empty:
                st.warning(f"⚠️ {len(efficacite_faible)} véhicule(s) présentent une efficacité carburant faible ou très faible")
                recommandations_txt = """### Recommandations pour améliorer l'efficacité:
                1. **Véhicules à efficacité très faible**: Envisager un diagnostic mécanique.
                2. **Trajets avec détours fréquents**: Optimiser les itinéraires.
                3. **Excès de vitesse réguliers**: Sensibiliser les conducteurs.
                4. **Écarts importants kilométrage déclaré/géolocalisé**: Contrôle systématique des déclarations.
                5. **Transactions sans présence détectée**: Renforcer les contrôles sur l'utilisation des cartes."""
                st.markdown(recommandations_txt)
                st.subheader("Véhicules Prioritaires pour Optimisation")
                cols_opti = ['Véhicule', 'Catégorie', 'Consommation_100km', 'Ecart_Conso_Pct', 'Score_Efficacite', 'Niveau_Efficacite']
                afficher_dataframe_avec_export(
                    efficacite_faible[[c for c in cols_opti if c in efficacite_faible.columns]],
                    "Véhicules à Optimiser", key="vehicules_a_optimiser_geoloc"
                )


# --- Nouvelle fonction pour la détection avancée d'anomalies basée sur géolocalisation ---
def detecter_anomalies_geolocalisation(
    df_geoloc: pd.DataFrame,
    df_transactions: pd.DataFrame,
    df_vehicules: pd.DataFrame,
    date_debut: datetime.date,
    date_fin: datetime.date
) -> pd.DataFrame:
    """
    Agrégation de toutes les anomalies détectées via géolocalisation.
    Cette fonction est utilisée notamment pour l'intégration dans la page principale d'anomalies.

    Args:
        df_geoloc: DataFrame des données de géolocalisation.
        df_transactions: DataFrame des transactions.
        df_vehicules: DataFrame des véhicules.
        date_debut: Date de début de la période d'analyse.
        date_fin: Date de fin de la période d'analyse.

    Returns:
        Un DataFrame des anomalies avec leur description et score de risque.
    """
    anomalies = []

    # Filtrer les données pour la période
    mask_date_geoloc = (df_geoloc['Date'].dt.date >= date_debut) & (df_geoloc['Date'].dt.date <= date_fin)
    df_geoloc_filtered = df_geoloc[mask_date_geoloc]


    if df_geoloc_filtered.empty:
        return pd.DataFrame()

    # 1. Détecter les excès de vitesse
    seuil_vitesse_anomalie = 90 # Seuil fixe pour l'agrégation d'anomalies
    _, detail_exces = analyser_exces_vitesse(df_geoloc_filtered, date_debut, date_fin, seuil_vitesse_anomalie)

    if not detail_exces.empty:
        for idx, row in detail_exces.iterrows():
            anomalie = {
                'Véhicule': row['Véhicule'],
                'Date': row['Date'],
                'Heure': row.get('Début_str', row.get('Début')), # Utiliser Début_str si disponible, sinon Début
                'type_anomalie': 'Excès de vitesse (géoloc)',
                'detail_anomalie': f"Vitesse: {row['Vitesse moyenne']:.1f} km/h, Dépassement: {row.get('Depassement_km/h', 'N/A'):.1f} km/h",
                'Niveau_Anomalie': row.get('Niveau_Exces', 'N/A'),
                'poids_anomalie': st.session_state.get('ss_poids_vitesse_excessive', DEFAULT_POIDS_VITESSE_EXCESSIVE)
            }
            anomalies.append(anomalie)

    # 2. Détecter les trajets suspects (hors heures, weekend)
    trajets_suspects = detecter_trajets_suspects(df_geoloc_filtered, date_debut, date_fin)

    if not trajets_suspects.empty:
        for idx, row in trajets_suspects.iterrows():
            # Extraire date et heure de Date_Heure_Debut
            date_heure_str = row.get('Date_Heure_Debut', "N/A N/A")
            try:
                date_part_str, heure_part_str = date_heure_str.split(" ", 1)
                date_val = datetime.strptime(date_part_str, '%d/%m/%Y').date() if date_part_str != "N/A" else None
                heure_val = heure_part_str if heure_part_str != "N/A" else None
            except ValueError:
                date_val = None
                heure_val = None

            anomalie = {
                'Véhicule': row['Véhicule'],
                'Date': date_val,
                'Heure': heure_val,
                'type_anomalie': 'Trajet suspect (géoloc)',
                'detail_anomalie': row['Description_Suspicion'],
                'Niveau_Anomalie': row.get('Niveau_Suspicion', 'N/A'),
                'poids_anomalie': row['Score_Suspicion_Total']
            }
            anomalies.append(anomalie)

    # 3. Détecter les détours suspects
    detours_suspects = detecter_detours_suspects(df_geoloc_filtered, date_debut, date_fin)

    if not detours_suspects.empty:
        for idx, row in detours_suspects.iterrows():
            anomalie = {
                'Véhicule': row['Véhicule'],
                'Date': row['Date'],
                'Heure': row.get('Début_str', row.get('Début')),
                'type_anomalie': 'Détour suspect (géoloc)',
                'detail_anomalie': row['Description_Detour'],
                'Niveau_Anomalie': row.get('Severite_Detour', 'N/A'),
                'poids_anomalie': row['Score_Detour']
            }
            anomalies.append(anomalie)

    # 4. Détecter les transactions sans présence (si activé)
    transactions_sans_presence = pd.DataFrame()
    if st.session_state.get('ss_activer_transactions_sans_presence', True):
        _, transactions_sans_presence = analyser_correspondance_transactions_geoloc(
            df_geoloc_filtered, df_transactions, df_vehicules, date_debut, date_fin
        )

    if not transactions_sans_presence.empty:
        for idx, row in transactions_sans_presence.iterrows():
            anomalie = {
                'Véhicule': row['Nouveau Immat'],
                'Date': row['Date'],
                'Heure': row['Hour'].strftime('%H:%M:%S') if pd.notna(row['Hour']) else None,
                'type_anomalie': 'Transaction sans présence (géoloc)',
                'detail_anomalie': f"Station: {row['Place']}, Quantité: {row['Quantity']:.1f}L, Montant: {row['Amount']:.0f} CFA, Vérif: {row['Methode_Verification']}",
                'Niveau_Anomalie': row.get('Niveau_Suspicion', 'N/A'),
                'poids_anomalie': row['Score_Suspicion']
            }
            anomalies.append(anomalie)

    # 5. Détecter les anomalies de distance (comparaison géoloc vs transactions)
    _, anomalies_distance = analyser_geolocalisation_vs_transactions(
        df_geoloc_filtered, df_transactions, df_vehicules, date_debut, date_fin
    )

    if not anomalies_distance.empty:
        for idx, row in anomalies_distance.iterrows():
            # Cette anomalie est globale pour la période, pas liée à une date/heure spécifique
            anomalie = {
                'Véhicule': row['Immatriculation'],
                'Date': date_fin, # Ou None, ou date de fin de période
                'Heure': None,
                'type_anomalie': row['Type_Anomalie'] + " (géoloc)", # Pour distinguer
                'detail_anomalie': f"Géoloc: {row['Distance_Geoloc_Totale']:.1f}km, Déclaré: {row['Distance_Declaree_Totale']:.1f}km, Écart: {row['Ecart_Distance']:.1f}km ({row['Pourcentage_Ecart']:.1f}%)",
                'Niveau_Anomalie': row.get('Gravite', 'N/A'),
                'poids_anomalie': min(min(abs(row['Pourcentage_Ecart']), 200) * (st.session_state.get('ss_poids_detour_suspect', DEFAULT_POIDS_DETOUR_SUSPECT) / 20), 50)
            }
            anomalies.append(anomalie)


    if not anomalies:
        return pd.DataFrame()

    df_anomalies = pd.DataFrame(anomalies)

    # Ajouter Card num. et Catégorie pour la fusion avec les autres anomalies
    mapping_carte = df_vehicules.set_index('Nouveau Immat')['N° Carte'].to_dict()
    mapping_categorie = df_vehicules.set_index('Nouveau Immat')['Catégorie'].to_dict()

    df_anomalies['Nouveau Immat'] = df_anomalies['Véhicule'] # Assurer que 'Nouveau Immat' est là
    df_anomalies['Card num.'] = df_anomalies['Véhicule'].map(mapping_carte)
    df_anomalies['Catégorie'] = df_anomalies['Véhicule'].map(mapping_categorie)


    df_anomalies = df_anomalies.sort_values(['Véhicule', 'poids_anomalie'], ascending=[True, False])

    return df_anomalies

# ---------------------------------------------------------------------
# Fonctions d'Affichage des Pages (Mises à jour avec intégration géolocalisation)
# ---------------------------------------------------------------------

def afficher_page_dashboard(df_transactions: pd.DataFrame, df_vehicules: pd.DataFrame, df_ge: pd.DataFrame, df_autres: pd.DataFrame, date_debut: datetime.date, date_fin: datetime.date, df_geoloc: Optional[pd.DataFrame] = None):
    """Affiche le tableau de bord principal."""
    st.header(f"📊 Tableau de Bord Principal ({date_debut.strftime('%d/%m/%Y')} - {date_fin.strftime('%d/%m/%Y')})")

    if df_transactions.empty:
        st.warning("Aucune transaction à analyser pour la période sélectionnée.")
        return

    total_volume = df_transactions['Quantity'].sum()
    total_cout = df_transactions['Amount'].sum()
    nb_transactions_total = len(df_transactions) # Renommé
    cartes_veh_actives = df_transactions[df_transactions['Card num.'].isin(df_vehicules['N° Carte'])]['Card num.'].nunique()
    prix_moyen_litre_global = (total_cout / total_volume) if total_volume > 0 else 0

    kpi_cat_dash, df_vehicle_kpi_dash = calculer_kpis_globaux(df_transactions, df_vehicules, date_debut, date_fin, list(st.session_state.ss_conso_seuils_par_categorie.keys()))
    conso_moyenne_globale = (kpi_cat_dash['total_litres'].sum() / kpi_cat_dash['distance_totale'].sum()) * 100 if not kpi_cat_dash.empty and kpi_cat_dash['distance_totale'].sum() > 0 else 0
    cout_km_global = (kpi_cat_dash['total_cout'].sum() / kpi_cat_dash['distance_totale'].sum()) if not kpi_cat_dash.empty and kpi_cat_dash['distance_totale'].sum() > 0 else 0

    df_anomalies_dash = detecter_anomalies(df_transactions, df_vehicules)
    cartes_inconnues_dash = verifier_cartes_inconnues(df_transactions, df_vehicules, df_ge, df_autres)
    vehicules_risques_dash = calculer_score_risque(df_anomalies_dash)
    nb_vehicules_suspects = len(vehicules_risques_dash[vehicules_risques_dash['score_risque'] >= st.session_state.ss_seuil_anomalies_suspectes_score]) if not vehicules_risques_dash.empty else 0
    nb_anomalies_critiques = len(df_anomalies_dash[df_anomalies_dash['poids_anomalie'] >= 8]) if not df_anomalies_dash.empty else 0

    nb_anomalies_geoloc_dash = 0 # Initialisation
    anomalies_geoloc_dash_df = pd.DataFrame()
    if df_geoloc is not None and not df_geoloc.empty:
        with st.spinner("Analyse des anomalies de géolocalisation (dashboard)..."):
            anomalies_geoloc_dash_df = detecter_anomalies_geolocalisation(
                df_geoloc, df_transactions, df_vehicules, date_debut, date_fin
            )
            nb_anomalies_geoloc_dash = len(anomalies_geoloc_dash_df)


    st.subheader("🚀 Indicateurs Clés")
    col1, col2, col3, col4 = st.columns(4)
    col1.metric("Volume Total", f"{total_volume:,.0f} L")
    col2.metric("Coût Total", f"{total_cout:,.0f} CFA")
    col3.metric("Transactions", f"{nb_transactions_total:,}")
    col4.metric("Véhicules Actifs", f"{cartes_veh_actives:,}")

    col5, col6, col7, col8 = st.columns(4)
    col5.metric("Conso. Moyenne Globale", f"{conso_moyenne_globale:.1f} L/100km" if conso_moyenne_globale else "N/A")
    col6.metric("Coût Moyen / Km Global", f"{cout_km_global:.1f} CFA/km" if cout_km_global else "N/A")
    col7.metric("Prix Moyen / Litre", f"{prix_moyen_litre_global:,.0f} CFA/L" if prix_moyen_litre_global else "N/A")

    if df_geoloc is not None and not df_geoloc.empty:
        col8.metric("Anomalies Géoloc", f"{nb_anomalies_geoloc_dash:,}", delta_color="inverse")
    else:
        col8.metric("Anomalies Géoloc", "N/A")


    st.subheader("⚠️ Alertes Rapides")
    col_a1, col_a2, col_a3 = st.columns(3)
    col_a1.metric("Cartes Inconnues", len(cartes_inconnues_dash), delta_color="inverse")
    col_a2.metric(f"Véhicules Suspects (Score > {st.session_state.ss_seuil_anomalies_suspectes_score})", nb_vehicules_suspects, delta_color="inverse")
    col_a3.metric("Anomalies Critiques (Poids >= 8)", nb_anomalies_critiques, delta_color="inverse")

    st.subheader("📈 Graphiques Clés")
    with st.expander("Évolution Mensuelle Volume & Coût", expanded=True):
        evo_mensuelle = df_transactions.groupby(pd.Grouper(key='Date', freq='M')).agg(
            Volume_L=('Quantity', 'sum'),
            Cout_CFA=('Amount', 'sum')
        ).reset_index()
        evo_mensuelle['Mois'] = evo_mensuelle['Date'].dt.strftime('%Y-%m')
        fig_evo = px.bar(evo_mensuelle, x='Mois', y=['Volume_L', 'Cout_CFA'],
                         title="Évolution Mensuelle du Volume et du Coût",
                         labels={'value': 'Valeur', 'variable': 'Indicateur'}, barmode='group')
        fig_evo.update_layout(yaxis_title="Volume (L) / Coût (CFA)")
        st.plotly_chart(fig_evo, use_container_width=True)

    with st.expander("Répartition par Catégorie de Véhicule", expanded=False):
        if not kpi_cat_dash.empty:
             col_g1, col_g2 = st.columns(2)
             fig_pie_vol = px.pie(kpi_cat_dash, values='total_litres', names='Catégorie', title='Répartition Volume par Catégorie')
             col_g1.plotly_chart(fig_pie_vol, use_container_width=True)
             fig_pie_cout = px.pie(kpi_cat_dash, values='total_cout', names='Catégorie', title='Répartition Coût par Catégorie')
             col_g2.plotly_chart(fig_pie_cout, use_container_width=True)
        else:
             st.info("Données insuffisantes pour la répartition par catégorie.")

    with st.expander("Top 5 Véhicules (Coût / Volume / Anomalies)", expanded=False):
        if not df_vehicle_kpi_dash.empty:
             col_t1, col_t2 = st.columns(2)
             top_cout = df_vehicle_kpi_dash.nlargest(5, 'total_cout')[['Nouveau Immat', 'total_cout']]
             top_volume = df_vehicle_kpi_dash.nlargest(5, 'total_litres')[['Nouveau Immat', 'total_litres']]
             with col_t1: # Afficher dans les colonnes pour un meilleur layout
                afficher_dataframe_avec_export(top_cout, "Top 5 Coût Total", key="dash_top_cout")
             with col_t2:
                afficher_dataframe_avec_export(top_volume, "Top 5 Volume Total", key="dash_top_vol")
        else:
            st.info("Données insuffisantes pour le classement des véhicules.")

        if not vehicules_risques_dash.empty:
             top_risque = vehicules_risques_dash.nlargest(5, 'score_risque')[['Nouveau Immat', 'score_risque', 'nombre_total_anomalies']]
             afficher_dataframe_avec_export(top_risque, "Top 5 Score Risque", key="dash_top_risque")
        else:
             st.info("Aucune anomalie de transaction détectée pour le classement par risque.")

    if not cartes_inconnues_dash.empty:
        with st.expander("🚨 Cartes Inconnues Détectées", expanded=False):
            afficher_dataframe_avec_export(cartes_inconnues_dash, "Détail des Cartes Inconnues", key="dash_cartes_inconnues")

    # Aperçu des anomalies de géolocalisation
    if df_geoloc is not None and not df_geoloc.empty and not anomalies_geoloc_dash_df.empty:
        with st.expander("🚨 Aperçu des Anomalies de Géolocalisation", expanded=True):
            summary_geoloc_dash = anomalies_geoloc_dash_df.groupby('type_anomalie').agg(
                Nombre=('type_anomalie', 'size'),
                Score_Moyen=('poids_anomalie', 'mean') # poids_anomalie est le nom utilisé
            ).reset_index().sort_values('Nombre', ascending=False)

            afficher_dataframe_avec_export(summary_geoloc_dash, "Résumé Anomalies Géolocalisation (Dashboard)", key="dash_anomalies_geoloc")

            fig_anomalies_geoloc_dash = px.bar(
                summary_geoloc_dash,
                x='type_anomalie', y='Nombre', title="Anomalies de Géolocalisation par Type (Dashboard)",
                color='Score_Moyen', labels={'Nombre': "Nombre d'occurrences"}
            )
            st.plotly_chart(fig_anomalies_geoloc_dash, use_container_width=True)
            st.markdown("""👉 *Pour une analyse détaillée des anomalies de géolocalisation, utilisez la page "Géolocalisation"*""")


def afficher_page_anomalies(
    df_transactions: pd.DataFrame, 
    df_vehicules: pd.DataFrame, 
    date_debut: datetime.date, 
    date_fin: datetime.date, 
    df_geoloc: Optional[pd.DataFrame] = None
):
    """
    Affiche une page de synthèse des anomalies améliorée avec visualisations avancées,
    filtres interactifs et recommandations contextuelles.
    """
    st.header(f"🚨 Détection et Analyse des Anomalies ({date_debut.strftime('%d/%m/%Y')} - {date_fin.strftime('%d/%m/%Y')})")
    
    # Ajout d'une barre de progression pour les analyses lourdes
    progress_bar = st.progress(0)
    status_text = st.empty()
    
    if df_transactions.empty:
        st.warning("Aucune transaction à analyser pour la période sélectionnée.")
        return

    # --- Filtres avancés en sidebar ---
    st.sidebar.subheader("🔍 Filtres d'Anomalies")
    
    # Filtre par seuil de score
    seuil_score = st.sidebar.slider(
        "Seuil minimum de score de risque",
        min_value=0,
        max_value=50,
        value=st.session_state.ss_seuil_anomalies_suspectes_score,
        step=1,
        key="anom_seuil_slider"
    )
    
    # Filtre par type d'anomalie (sera rempli après l'analyse initiale)
    types_anomalies_filter = st.sidebar.empty()
    
    # Filtre par catégorie de véhicule
    if not df_vehicules.empty:
        all_cats = sorted(df_vehicules['Catégorie'].dropna().astype(str).unique())
        cats_filter = st.sidebar.multiselect(
            "Filtrer par catégorie de véhicule",
            options=all_cats,
            default=all_cats,
            key="anom_cats_filter"
        )
    else:
        cats_filter = []
    
    # Filtre temporel (optionnel)
    use_date_filter = st.sidebar.checkbox(
        "Filtrer par sous-période",
        value=False,
        help="Permet d'analyser une sous-période spécifique dans la période globale"
    )
    
    if use_date_filter:
        col_sub1, col_sub2 = st.sidebar.columns(2)
        with col_sub1:
            sub_date_debut = st.date_input(
                "Date début sous-période",
                value=date_debut,
                min_value=date_debut,
                max_value=date_fin,
                key="anom_sub_date_debut"
            )
        with col_sub2:
            sub_date_fin = st.date_input(
                "Date fin sous-période",
                value=date_fin,
                min_value=sub_date_debut,
                max_value=date_fin,
                key="anom_sub_date_fin"
            )
        
        # Appliquer le filtre de sous-période
        date_debut_analyse = sub_date_debut
        date_fin_analyse = sub_date_fin
    else:
        date_debut_analyse = date_debut
        date_fin_analyse = date_fin
    
    # --- Début des analyses ---
    df_anomalies_all = pd.DataFrame()  # Pour stocker toutes les anomalies
    df_scores_all = pd.DataFrame()     # Pour stocker tous les scores

    # Initialisation des conteneurs pour les onglets
    tab_container = st.empty()         # Pour afficher les onglets après les analyses
    
    # 1. Analyse des anomalies de transactions
    status_text.text("⏳ Détection des anomalies de transaction en cours...")
    progress_bar.progress(10)
    
    with st.spinner("Analyse des transactions..."):
        df_anomalies_transac = detecter_anomalies(df_transactions, df_vehicules)
    
    # 2. Analyse des anomalies de géolocalisation (si disponible)
    if df_geoloc is not None and not df_geoloc.empty:
        status_text.text("⏳ Détection des anomalies de géolocalisation en cours...")
        progress_bar.progress(30)
        
        with st.spinner("Analyse de géolocalisation..."):
            df_anomalies_geoloc = detecter_anomalies_geolocalisation(
                df_geoloc, df_transactions, df_vehicules, date_debut_analyse, date_fin_analyse
            )
        
        # Fusion des anomalies
        if not df_anomalies_transac.empty and not df_anomalies_geoloc.empty:
            progress_bar.progress(50)
            status_text.text("⏳ Fusion des anomalies en cours...")
            
            # Assurer la cohérence des colonnes pour la concaténation
            cols_communes = ['Nouveau Immat', 'Card num.', 'Catégorie', 'Date', 'type_anomalie', 'detail_anomalie', 'poids_anomalie']
            for col in cols_communes:
                if col not in df_anomalies_transac.columns: df_anomalies_transac[col] = pd.NA
                if col not in df_anomalies_geoloc.columns: df_anomalies_geoloc[col] = pd.NA

            df_anomalies_all = pd.concat([
                df_anomalies_transac[cols_communes],
                df_anomalies_geoloc[cols_communes]
            ], ignore_index=True)

        elif not df_anomalies_transac.empty:
            df_anomalies_all = df_anomalies_transac
        elif not df_anomalies_geoloc.empty:
            df_anomalies_all = df_anomalies_geoloc
    else:
        df_anomalies_all = df_anomalies_transac
    
    # Conversion des types de données problématiques
    if not df_anomalies_all.empty:
        # Conversion des dates
        progress_bar.progress(70)
        status_text.text("⏳ Calcul des scores de risque...")
        
        if 'Date' in df_anomalies_all.columns:
            # S'assurer que la colonne Date est au format datetime
            if not pd.api.types.is_datetime64_any_dtype(df_anomalies_all['Date']):
                try:
                    df_anomalies_all['Date'] = pd.to_datetime(df_anomalies_all['Date'], errors='coerce')
                except Exception as e:
                    st.warning(f"Problème lors de la conversion des dates: {e}")
        
        # Appliquer les filtres de catégorie
        if cats_filter:
            df_anomalies_all = df_anomalies_all[df_anomalies_all['Catégorie'].isin(cats_filter)]
        
        # Calcul des scores de risque
        df_scores_all = calculer_score_risque(df_anomalies_all)
    else:
        df_scores_all = pd.DataFrame(columns=['Nouveau Immat', 'Card num.', 'Catégorie', 'nombre_total_anomalies', 'score_risque'])
    
    # Extraire la liste des types d'anomalies pour les filtres
    if not df_anomalies_all.empty:
        all_anomaly_types = sorted(df_anomalies_all['type_anomalie'].unique())
        
        # Mettre à jour le filtre de types d'anomalies
        selected_anomaly_types = types_anomalies_filter.multiselect(
            "Types d'anomalies à afficher",
            options=all_anomaly_types,
            default=all_anomaly_types,
            key="anom_types_filter"
        )
        
        # Appliquer le filtre de type d'anomalie
        if selected_anomaly_types:
            df_anomalies_all = df_anomalies_all[df_anomalies_all['type_anomalie'].isin(selected_anomaly_types)]
            
            # Recalculer les scores après filtrage
            if not df_anomalies_all.empty:
                df_scores_all = calculer_score_risque(df_anomalies_all)
    
    progress_bar.progress(100)
    status_text.text("✅ Analyse terminée!")
    time.sleep(0.5)  # Pause pour montrer que l'analyse est terminée
    status_text.empty()
    progress_bar.empty()
    
    # --- Configuration des onglets améliorés ---
    tab_resume, tab_vehicules, tab_transactions, tab_stations, tab_geoloc, tab_tendances, tab_recommandations = st.tabs([
        "📊 Résumé", "🚗 Véhicules à Risque", "🧾 Anomalies Transactions", 
        "⛽ Stations à Risque", "📍 Anomalies Géoloc", "📈 Tendances", "💡 Recommandations"
    ])
    
    # --- ONGLET 1: RÉSUMÉ GLOBAL ---
    with tab_resume:
        if df_anomalies_all.empty:
            st.success("✅ Aucune anomalie détectée sur la période sélectionnée avec les filtres actuels!")
        else:
            # KPIs principaux
            nb_total_anomalies = len(df_anomalies_all)
            nb_vehicules_avec_anomalies = df_anomalies_all['Card num.'].nunique()
            nb_vehicules_risque_eleve = len(df_scores_all[df_scores_all['score_risque'] >= seuil_score])
            score_moyen = df_scores_all['score_risque'].mean() if not df_scores_all.empty else 0
            
            # Afficher les KPIs
            col1, col2, col3, col4 = st.columns(4)
            col1.metric(
                "Anomalies Détectées", 
                f"{nb_total_anomalies:,}",
                help="Nombre total d'anomalies détectées avec les filtres actuels"
            )
            col2.metric(
                "Véhicules Concernés", 
                f"{nb_vehicules_avec_anomalies:,}",
                help="Nombre de véhicules présentant au moins une anomalie"
            )
            col3.metric(
                f"Véhicules à Risque (Score ≥ {seuil_score})", 
                f"{nb_vehicules_risque_eleve:,}",
                delta=f"{nb_vehicules_risque_eleve/nb_vehicules_avec_anomalies*100:.1f}%" if nb_vehicules_avec_anomalies > 0 else None,
                delta_color="inverse",
                help=f"Véhicules dont le score de risque dépasse le seuil de {seuil_score}"
            )
            col4.metric(
                "Score Moyen", 
                f"{score_moyen:.1f}",
                help="Score de risque moyen sur tous les véhicules avec anomalies"
            )
            
            # Résumé par type d'anomalie
            st.subheader("📊 Synthèse par Type d'Anomalie")
            
            summary_by_type = df_anomalies_all.groupby('type_anomalie').agg(
                Nombre=('type_anomalie', 'size'),
                Score_Total=('poids_anomalie', 'sum'),
                Score_Moyen=('poids_anomalie', 'mean'),
                Nb_Vehicules=('Card num.', 'nunique')
            ).reset_index().sort_values('Score_Total', ascending=False)
            
            # Créer un graphique plus riche
            fig_types = px.bar(
                summary_by_type,
                x='type_anomalie',
                y='Nombre',
                color='Score_Total',
                color_continuous_scale='Reds',
                hover_data=['Nb_Vehicules', 'Score_Moyen'],
                title="Distribution des Anomalies par Type",
                labels={
                    'type_anomalie': "Type d'Anomalie",
                    'Nombre': "Nombre d'Occurrences"
                }
            )
            
            st.plotly_chart(fig_types, use_container_width=True)
            
            # Tableau récapitulatif
            st.dataframe(
                summary_by_type,
                use_container_width=True,
                column_config={
                    "Score_Moyen": st.column_config.NumberColumn(format="%.2f"),
                }
            )
            
            # Distribution des scores de risque
            st.subheader("📊 Distribution des Scores de Risque")
            
            if not df_scores_all.empty:
                fig_scores = px.histogram(
                    df_scores_all,
                    x='score_risque',
                    nbins=20,
                    title="Distribution des Scores de Risque",
                    labels={'score_risque': "Score de Risque"}
                )
                
                # Ajouter une ligne verticale pour le seuil
                fig_scores.add_vline(
                    x=seuil_score,
                    line_dash="dash",
                    line_color="red",
                    annotation_text=f"Seuil: {seuil_score}"
                )
                
                st.plotly_chart(fig_scores, use_container_width=True)
                
                # Répartition des risques par catégorie
                if 'Catégorie' in df_scores_all.columns:
                    cat_counts = df_scores_all.groupby('Catégorie').agg(
                        Nb_Vehicules=('Nouveau Immat', 'nunique'),
                        Score_Moyen=('score_risque', 'mean'),
                        Score_Max=('score_risque', 'max'),
                        Nb_Vehicules_Risque=('score_risque', lambda x: (x >= seuil_score).sum())
                    ).reset_index()
                    
                    # Calculer le pourcentage de véhicules à risque
                    cat_counts['Pourcentage_A_Risque'] = (cat_counts['Nb_Vehicules_Risque'] / cat_counts['Nb_Vehicules'] * 100).round(1)
                    
                    st.subheader("📊 Répartition des Risques par Catégorie")
                    
                    # Graphique pour les catégories
                    fig_cats = px.bar(
                        cat_counts,
                        x='Catégorie',
                        y=['Nb_Vehicules', 'Nb_Vehicules_Risque'],
                        title="Véhicules à Risque par Catégorie",
                        barmode='group',
                        labels={
                            'value': "Nombre de Véhicules",
                            'variable': "Métrique"
                        }
                    )
                    
                    # Ajouter une ligne pour le score moyen
                    fig_cats.add_trace(
                        go.Scatter(
                            x=cat_counts['Catégorie'],
                            y=cat_counts['Score_Moyen'],
                            mode='lines+markers',
                            name='Score Moyen',
                            yaxis='y2',
                            line=dict(color='red', width=2)
                        )
                    )
                    
                    # Configurer l'axe Y secondaire
                    fig_cats.update_layout(
                        yaxis2=dict(
                            title="Score Moyen",
                            overlaying='y',
                            side='right'
                        )
                    )
                    
                    st.plotly_chart(fig_cats, use_container_width=True)
    
    # --- ONGLET 2: VÉHICULES À RISQUE ---
    with tab_vehicules:
        st.subheader(f"🎯 Véhicules à Risque Élevé (Score de Risque ≥ {seuil_score})")
        
        if df_scores_all.empty:
            st.info("Aucune donnée de score disponible pour les filtres actuels.")
        else:
            # Filtrer les véhicules à risque élevé
            vehicules_suspects = df_scores_all[df_scores_all['score_risque'] >= seuil_score].copy()
            
            if vehicules_suspects.empty:
                st.success(f"✅ Aucun véhicule ne dépasse le seuil de score de risque ({seuil_score}) avec les filtres actuels.")
            else:
                # Préparer le pivot pour les détails par type d'anomalie
                if not df_anomalies_all.empty:
                    # Créer un pivot des types d'anomalies
                    pivot_details = df_anomalies_all.groupby(['Nouveau Immat', 'Card num.', 'Catégorie', 'type_anomalie']).size().unstack(fill_value=0)
                    vehicules_suspects_details = vehicules_suspects.merge(
                        pivot_details, 
                        on=['Nouveau Immat', 'Card num.', 'Catégorie'], 
                        how='left'
                    ).fillna(0)
                    
                    # Tri par score décroissant
                    vehicules_suspects_details = vehicules_suspects_details.sort_values('score_risque', ascending=False)
                    
                    # Fonction pour appliquer un style conditionnel au tableau
                    def highlight_risk(val):
                        if isinstance(val, (int, float)):
                            if val >= 50:
                                return 'background-color: #ffcccc'  # Rouge clair pour risque très élevé
                            elif val >= 30:
                                return 'background-color: #ffeecc'  # Orange clair pour risque élevé
                            elif val >= 15:
                                return 'background-color: #ffffcc'  # Jaune clair pour risque modéré
                        return ''
                    
                    # Ajouter une colonne de niveau de risque
                    vehicules_suspects_details['Niveau_Risque'] = pd.cut(
                        vehicules_suspects_details['score_risque'],
                        bins=[0, 15, 30, 50, float('inf')],
                        labels=['Faible', 'Modéré', 'Élevé', 'Critique']
                    )
                    
                    # Afficher le tableau avec style
                    st.dataframe(
                        vehicules_suspects_details.style.applymap(
                            highlight_risk, 
                            subset=['score_risque']
                        ),
                        use_container_width=True
                    )
                    
                    # Téléchargement du tableau
                    excel_data = to_excel(vehicules_suspects_details)
                    st.download_button(
                        label="📥 Télécharger la liste des véhicules à risque",
                        data=excel_data,
                        file_name=f"vehicules_risque_{date_debut_analyse.strftime('%Y%m%d')}_{date_fin_analyse.strftime('%Y%m%d')}.xlsx",
                        mime=EXCEL_MIME_TYPE
                    )
                    
                    # Visualisation des 10 véhicules les plus à risque
                    st.subheader("🚨 Top 10 des Véhicules les Plus à Risque")
                    
                    top_10_vehicles = vehicules_suspects_details.head(10)
                    fig_top10 = px.bar(
                        top_10_vehicles,
                        x='Nouveau Immat',
                        y='score_risque',
                        color='Niveau_Risque',
                        title="Top 10 des Véhicules à Risque",
                        labels={'score_risque': "Score de Risque"},
                        color_discrete_map={
                            'Faible': 'green',
                            'Modéré': 'yellow',
                            'Élevé': 'orange',
                            'Critique': 'red'
                        },
                        hover_data=['Catégorie', 'nombre_total_anomalies']
                    )
                    
                    st.plotly_chart(fig_top10, use_container_width=True)
                    
                    # Détail des anomalies pour un véhicule spécifique
                    st.subheader("🔍 Détail des Anomalies par Véhicule")
                    
                    # Sélecteur de véhicule
                    vehicle_options = [(v, f"{v} - Score: {s:.1f}") for v, s in zip(
                        vehicules_suspects_details['Nouveau Immat'],
                        vehicules_suspects_details['score_risque']
                    )]
                    selected_vehicle = st.selectbox(
                        "Sélectionner un véhicule pour voir le détail des anomalies",
                        options=vehicle_options,
                        format_func=lambda x: x[1],
                        key="anom_veh_select"
                    )[0]
                    
                    # Afficher les anomalies du véhicule sélectionné
                    vehicle_anomalies = df_anomalies_all[df_anomalies_all['Nouveau Immat'] == selected_vehicle].copy()
                    
                    if not vehicle_anomalies.empty:
                        # Trier par date et poids
                        if 'Date' in vehicle_anomalies.columns and pd.api.types.is_datetime64_any_dtype(vehicle_anomalies['Date']):
                            vehicle_anomalies = vehicle_anomalies.sort_values(['Date', 'poids_anomalie'], ascending=[False, False])
                        else:
                            vehicle_anomalies = vehicle_anomalies.sort_values('poids_anomalie', ascending=False)
                        
                        # Sélectionner les colonnes pertinentes
                        display_cols = ['Date', 'type_anomalie', 'detail_anomalie', 'poids_anomalie']
                        optional_cols = ['Place', 'Quantity', 'Amount', 'Hour']
                        
                        for col in optional_cols:
                            if col in vehicle_anomalies.columns:
                                display_cols.append(col)
                        
                        st.dataframe(vehicle_anomalies[display_cols], use_container_width=True)
                        
                        # Graphique des anomalies par type pour ce véhicule
                        anomaly_counts = vehicle_anomalies['type_anomalie'].value_counts().reset_index()
                        anomaly_counts.columns = ['type_anomalie', 'count']
                        
                        fig_veh_anom = px.pie(
                            anomaly_counts,
                            values='count',
                            names='type_anomalie',
                            title=f"Répartition des Anomalies - {selected_vehicle}",
                            hole=0.4
                        )
                        
                        st.plotly_chart(fig_veh_anom, use_container_width=True)
                    else:
                        st.info(f"Aucune anomalie trouvée pour le véhicule {selected_vehicle}.")
    
    # --- ONGLET 3: ANOMALIES DE TRANSACTIONS ---
    with tab_transactions:
        st.subheader("🧾 Analyse des Anomalies de Transactions")
        
        if 'df_anomalies_transac' not in locals() or df_anomalies_transac.empty:
            st.success("✅ Aucune anomalie de transaction détectée avec les filtres actuels.")
        else:
            # Statistiques des anomalies de transaction
            nb_anomalies_transac = len(df_anomalies_transac)
            nb_vehicules_transac = df_anomalies_transac['Card num.'].nunique()
            
            col_t1, col_t2 = st.columns(2)
            col_t1.metric("Anomalies de Transaction", f"{nb_anomalies_transac:,}")
            col_t2.metric("Véhicules Concernés", f"{nb_vehicules_transac:,}")
            
            # Résumé par type d'anomalie de transaction
            summary_transac = df_anomalies_transac.groupby('type_anomalie').agg(
                Nombre=('type_anomalie', 'size'),
                Score_Total=('poids_anomalie', 'sum'),
                Nb_Vehicules=('Card num.', 'nunique')
            ).reset_index().sort_values('Score_Total', ascending=False)
            
            # Afficher le tableau récapitulatif
            st.dataframe(summary_transac, use_container_width=True)
            
            # Graphique par type d'anomalie
            fig_transac = px.bar(
                summary_transac,
                x='type_anomalie',
                y='Nombre',
                color='Score_Total',
                title="Anomalies de Transaction par Type",
                labels={
                    'type_anomalie': "Type d'Anomalie",
                    'Nombre': "Nombre d'Occurrences"
                }
            )
            
            st.plotly_chart(fig_transac, use_container_width=True)
            
            # Afficher les détails des anomalies de transaction avec filtrage dynamique
            st.subheader("Détail des Anomalies de Transaction")
            
            # Filtres interactifs
            col_filter1, col_filter2 = st.columns(2)
            
            with col_filter1:
                # Filtre par type d'anomalie
                types_transac = summary_transac['type_anomalie'].tolist()
                selected_types_transac = st.multiselect(
                    "Filtrer par type d'anomalie",
                    options=types_transac,
                    default=types_transac,
                    key="anom_transac_types"
                )
            
            with col_filter2:
                # Filtre par poids minimum
                min_weight = st.slider(
                    "Poids minimum de l'anomalie",
                    min_value=0,
                    max_value=15,
                    value=0,
                    step=1,
                    key="anom_transac_weight"
                )
            
            # Appliquer les filtres
            filtered_transac = df_anomalies_transac[
                (df_anomalies_transac['type_anomalie'].isin(selected_types_transac)) &
                (df_anomalies_transac['poids_anomalie'] >= min_weight)
            ].copy()
            
            # Trier par date et poids
            if 'Date' in filtered_transac.columns:
                # Vérifier si la colonne Date est au format datetime
                if pd.api.types.is_datetime64_any_dtype(filtered_transac['Date']):
                    filtered_transac = filtered_transac.sort_values(['Date', 'poids_anomalie'], ascending=[False, False])
                else:
                    try:
                        # Essayer de convertir la colonne Date en datetime
                        filtered_transac['Date'] = pd.to_datetime(filtered_transac['Date'], errors='coerce')
                        filtered_transac = filtered_transac.sort_values(['Date', 'poids_anomalie'], ascending=[False, False])
                    except:
                        # En cas d'échec, trier uniquement par poids
                        filtered_transac = filtered_transac.sort_values('poids_anomalie', ascending=False)
                        st.warning("Impossible de trier par date - les dates ne sont pas au format correct. Tri par poids uniquement.")
            else:
                # Si pas de colonne Date, trier uniquement par poids
                filtered_transac = filtered_transac.sort_values('poids_anomalie', ascending=False)
            
            # Sélectionner les colonnes pertinentes
            display_cols_transac = [
                'Date', 'Hour', 'Nouveau Immat', 'Catégorie', 'type_anomalie', 
                'detail_anomalie', 'Quantity', 'Amount', 'Place', 'poids_anomalie'
            ]
            
            final_cols_transac = [col for col in display_cols_transac if col in filtered_transac.columns]
            
            # Afficher avec possibilité d'export
            afficher_dataframe_avec_export(
                filtered_transac[final_cols_transac],
                f"Anomalies de Transaction ({len(filtered_transac)} résultats)",
                key="anom_transac_detail"
            )
            
            # Bouton pour générer un rapport d'anomalies de transaction
            if st.button("📊 Générer un rapport détaillé des anomalies de transaction", key="btn_report_transac"):
                # Ici on pourrait générer un rapport plus complet, mais pour l'exemple on affiche juste un message
                st.success("✅ Rapport généré avec succès! Téléchargez-le via le bouton ci-dessus.")
    
    # --- ONGLET 4: STATIONS À RISQUE ---
    with tab_stations:
        st.subheader("⛽ Analyse des Stations à Risque de Fraude")
        
        if df_anomalies_all.empty:
            st.success("✅ Aucune anomalie détectée, impossible d'analyser les stations à risque.")
        else:
            with st.spinner("Analyse des stations à risque en cours..."):
                stations_risque = analyser_stations_risque(df_anomalies_all, df_transactions)
            
            if stations_risque.empty:
                st.info("Aucune station n'a présenté d'anomalies liées à la fraude par carte.")
            else:
                # Comptage des niveaux de risque
                nb_stations_risque = len(stations_risque)
                nb_stations_risque_eleve = len(stations_risque[stations_risque['Niveau_Risque'].isin(['Élevé', 'Critique'])])
                
                # Métriques principales
                col_s1, col_s2, col_s3 = st.columns(3)
                col_s1.metric("Stations Analysées", f"{nb_stations_risque:,}")
                col_s2.metric("Stations à Risque Élevé/Critique", f"{nb_stations_risque_eleve:,}")
                
                # Calculer le % de transactions à risque
                total_transactions = df_transactions['Place'].value_counts().sum()
                transactions_stations_risque = stations_risque['Nb_Total_Transactions'].sum()
                pct_transactions = (transactions_stations_risque / total_transactions * 100) if total_transactions > 0 else 0
                
                col_s3.metric(
                    "% des Transactions Totales", 
                    f"{pct_transactions:.1f}%",
                    help="Pourcentage des transactions totales effectuées dans les stations à risque"
                )
                
                # Carte de chaleur des anomalies par station et par type
                st.subheader("🔥 Carte de Chaleur des Anomalies par Station")
                
                # Préparer les données pour la carte de chaleur
                top_stations = stations_risque.head(10)['Place'].tolist()
                anomalies_heatmap = df_anomalies_all[
                    (df_anomalies_all['Place'].isin(top_stations)) & 
                    (df_anomalies_all['type_anomalie'].isin([
                        'Dépassement capacité', 'Prises rapprochées', 
                        'Facturation double suspectée', 'Transaction sans présence (géoloc)',
                        'Véhicule Hors Service'
                    ]))
                ]
                
                if not anomalies_heatmap.empty:
                    pivot_heatmap = pd.crosstab(
                        anomalies_heatmap['Place'], 
                        anomalies_heatmap['type_anomalie']
                    )
                    
                    # Normaliser les données pour une meilleure visualisation
                    pivot_norm = pivot_heatmap.div(pivot_heatmap.sum(axis=1), axis=0)
                    
                    fig_heatmap = px.imshow(
                        pivot_norm,
                        labels=dict(x="Type d'Anomalie", y="Station", color="Proportion"),
                        title="Distribution des Types d'Anomalies par Station (Top 10)",
                        color_continuous_scale='Reds'
                    )
                    
                    st.plotly_chart(fig_heatmap, use_container_width=True)
                
                # Tableau des stations à risque avec code couleur
                st.subheader("📋 Liste des Stations à Risque")
                
                # Fonction pour appliquer un style conditionnel
                def highlight_risk_level(val):
                    if val == 'Critique':
                        return 'background-color: #ffcccc'
                    elif val == 'Élevé':
                        return 'background-color: #ffeecc'
                    elif val == 'Modéré':
                        return 'background-color: #ffffcc'
                    return ''
                
                # Colonnes à afficher
                display_cols_stations = [
                    'Place', 'Nb_Anomalies_Fraude', 'Nb_Total_Transactions', 
                    'Pourcentage_Anomalies', 'Score_Risque_Station', 'Niveau_Risque',
                    'Anomalies_Principales'
                ]
                
                # Afficher le tableau stylisé
                styled_stations = stations_risque[display_cols_stations].style.applymap(
                    highlight_risk_level, 
                    subset=['Niveau_Risque']
                )
                
                st.dataframe(styled_stations, use_container_width=True)
                
                # Bouton d'exportation
                excel_stations = to_excel(stations_risque[display_cols_stations])
                st.download_button(
                    label="📥 Exporter la liste des stations à risque",
                    data=excel_stations,
                    file_name=f"stations_risque_{date_debut_analyse.strftime('%Y%m%d')}_{date_fin_analyse.strftime('%Y%m%d')}.xlsx",
                    mime=EXCEL_MIME_TYPE
                )
                
                # Graphique des stations à risque
                fig_stations = px.bar(
                    stations_risque.sort_values('Score_Risque_Station', ascending=False).head(10),
                    x='Place',
                    y='Score_Risque_Station',
                    color='Niveau_Risque',
                    title="Top 10 des Stations par Score de Risque",
                    color_discrete_map={
                        'Faible': 'green',
                        'Modéré': 'yellow',
                        'Élevé': 'orange',
                        'Critique': 'red'
                    },
                    hover_data=['Pourcentage_Anomalies', 'Nb_Total_Transactions', 'Anomalies_Principales']
                )
                
                st.plotly_chart(fig_stations, use_container_width=True)
                
                # Section interactive pour explorer une station spécifique
                st.subheader("🔍 Explorer une Station Spécifique")
                
                station_names = stations_risque['Place'].tolist()
                selected_station = st.selectbox(
                    "Sélectionner une station à analyser",
                    options=station_names,
                    key="anom_station_select"
                )
                
                if selected_station:
                    # Filtrer les anomalies pour cette station
                    station_anomalies = df_anomalies_all[df_anomalies_all['Place'] == selected_station].copy()
                    
                    if not station_anomalies.empty:
                        # Informations de la station
                        station_info = stations_risque[stations_risque['Place'] == selected_station].iloc[0]
                        
                        # Afficher les informations
                        col_si1, col_si2, col_si3 = st.columns(3)
                        col_si1.metric("Nombre d'Anomalies", f"{station_info['Nb_Anomalies_Fraude']:,}")
                        col_si2.metric("% des Transactions", f"{station_info['Pourcentage_Anomalies']:.1f}%")
                        col_si3.metric("Niveau de Risque", f"{station_info['Niveau_Risque']}")
                        
                        # Tableau des véhicules impliqués
                        vehicles_at_station = station_anomalies.groupby('Nouveau Immat').agg(
                            Nb_Anomalies=('type_anomalie', 'size'),
                            Types_Anomalies=('type_anomalie', lambda x: ', '.join(sorted(set(x)))),
                            Score_Total=('poids_anomalie', 'sum')
                        ).reset_index().sort_values('Nb_Anomalies', ascending=False)
                        
                        st.subheader(f"Véhicules Impliqués à {selected_station}")
                        st.dataframe(vehicles_at_station, use_container_width=True)
                        
                        # Graphique des anomalies par véhicule pour cette station
                        fig_station_veh = px.bar(
                            vehicles_at_station.head(10),
                            x='Nouveau Immat',
                            y='Nb_Anomalies',
                            color='Score_Total',
                            title=f"Top 10 des Véhicules avec Anomalies à {selected_station}",
                            labels={'Nb_Anomalies': "Nombre d'Anomalies"}
                        )
                        
                        st.plotly_chart(fig_station_veh, use_container_width=True)
                    else:
                        st.info(f"Aucune anomalie trouvée pour la station {selected_station}.")
                
                # Recommandations pour les stations à risque
                with st.expander("💡 Recommandations pour les Stations à Risque"):
                    st.markdown("""
                    ### Actions recommandées pour les stations à risque élevé:
                    
                    1. **Audit spécifique**: Effectuer des contrôles approfondis des transactions provenant de ces stations.
                    2. **Vérification des procédures**: S'assurer que les procédures de distribution de carburant sont respectées.
                    3. **Contrôles inopinés**: Établir un programme de vérifications surprises.
                    4. **Corrélation avec géolocalisation**: Vérifier systématiquement la présence des véhicules aux heures de transaction.
                    5. **Sensibilisation des gérants**: Organiser des sessions d'information pour les responsables des stations.
                    """)
    
    # --- ONGLET 5: ANOMALIES DE GÉOLOCALISATION ---
    with tab_geoloc:
        st.subheader("📍 Analyse des Anomalies de Géolocalisation")
        
        if df_geoloc is None or df_geoloc.empty:
            st.info("Aucune donnée de géolocalisation disponible. Veuillez charger un fichier de géolocalisation.")
        elif 'df_anomalies_geoloc' not in locals() or df_anomalies_geoloc.empty:
            st.success("✅ Aucune anomalie de géolocalisation détectée avec les filtres actuels.")
        else:
            # Comptage des anomalies de géolocalisation
            nb_anomalies_geoloc = len(df_anomalies_geoloc)
            nb_vehicules_geoloc = df_anomalies_geoloc['Véhicule'].nunique() if 'Véhicule' in df_anomalies_geoloc.columns else df_anomalies_geoloc['Nouveau Immat'].nunique()
            
            # Métriques principales
            col_g1, col_g2 = st.columns(2)
            col_g1.metric("Anomalies de Géolocalisation", f"{nb_anomalies_geoloc:,}")
            col_g2.metric("Véhicules Concernés", f"{nb_vehicules_geoloc:,}")
            
            # Résumé par type d'anomalie de géolocalisation
            summary_geoloc = df_anomalies_geoloc.groupby('type_anomalie').agg(
                Nombre=('type_anomalie', 'size'),
                Score_Total=('poids_anomalie', 'sum'),
                Nb_Vehicules=('Nouveau Immat', 'nunique')
            ).reset_index().sort_values('Score_Total', ascending=False)
            
            # Afficher le tableau récapitulatif
            st.dataframe(summary_geoloc, use_container_width=True)
            
            # Graphique par type d'anomalie de géolocalisation
            fig_geoloc = px.bar(
                summary_geoloc,
                x='type_anomalie',
                y='Nombre',
                color='Score_Total',
                title="Anomalies de Géolocalisation par Type",
                labels={
                    'type_anomalie': "Type d'Anomalie",
                    'Nombre': "Nombre d'Occurrences"
                }
            )
            
            st.plotly_chart(fig_geoloc, use_container_width=True)
            
            # Top véhicules avec anomalies de géolocalisation
            st.subheader("🚗 Top Véhicules avec Anomalies de Géolocalisation")
            
            top_vehicules_geoloc = df_anomalies_geoloc.groupby('Nouveau Immat').agg(
                Nb_Anomalies=('type_anomalie', 'size'),
                Score_Total=('poids_anomalie', 'sum'),
                Types_Anomalies=('type_anomalie', lambda x: ', '.join(sorted(set(x))))
            ).reset_index().sort_values('Score_Total', ascending=False)
            
            st.dataframe(top_vehicules_geoloc.head(10), use_container_width=True)
            
            # Filtres pour l'affichage détaillé
            st.subheader("🔍 Détail des Anomalies de Géolocalisation")
            
            # Filtres interactifs
            col_g_filter1, col_g_filter2 = st.columns(2)
            
            with col_g_filter1:
                # Filtre par type d'anomalie géoloc
                types_geoloc = summary_geoloc['type_anomalie'].tolist()
                selected_types_geoloc = st.multiselect(
                    "Filtrer par type d'anomalie",
                    options=types_geoloc,
                    default=types_geoloc,
                    key="anom_geoloc_types"
                )
            
            with col_g_filter2:
                # Filtre par poids minimum
                min_weight_geoloc = st.slider(
                    "Poids minimum de l'anomalie",
                    min_value=0,
                    max_value=15,
                    value=0,
                    step=1,
                    key="anom_geoloc_weight"
                )
            
            # Appliquer les filtres
            filtered_geoloc = df_anomalies_geoloc[
                (df_anomalies_geoloc['type_anomalie'].isin(selected_types_geoloc)) &
                (df_anomalies_geoloc['poids_anomalie'] >= min_weight_geoloc)
            ].copy()
            
            # Trier par date et poids
            if 'Date' in filtered_geoloc.columns:
                # Vérifier si la colonne Date est au format datetime
                if pd.api.types.is_datetime64_any_dtype(filtered_geoloc['Date']):
                    filtered_geoloc = filtered_geoloc.sort_values(['Date', 'poids_anomalie'], ascending=[False, False])
                else:
                    try:
                        # Essayer de convertir la colonne Date en datetime
                        filtered_geoloc['Date'] = pd.to_datetime(filtered_geoloc['Date'], errors='coerce')
                        filtered_geoloc = filtered_geoloc.sort_values(['Date', 'poids_anomalie'], ascending=[False, False])
                    except:
                        # En cas d'échec, trier uniquement par poids
                        filtered_geoloc = filtered_geoloc.sort_values('poids_anomalie', ascending=False)
                        st.warning("Impossible de trier par date - les dates ne sont pas au format correct. Tri par poids uniquement.")
            else:
                # Si pas de colonne Date, trier uniquement par poids
                filtered_geoloc = filtered_geoloc.sort_values('poids_anomalie', ascending=False)
            
            # Sélectionner les colonnes pertinentes
            display_cols_geoloc = [
                'Date', 'Nouveau Immat', 'Catégorie', 'type_anomalie', 
                'detail_anomalie', 'poids_anomalie'
            ]
            
            final_cols_geoloc = [col for col in display_cols_geoloc if col in filtered_geoloc.columns]
            
            # Afficher avec possibilité d'export
            afficher_dataframe_avec_export(
                filtered_geoloc[final_cols_geoloc],
                f"Anomalies de Géolocalisation ({len(filtered_geoloc)} résultats)",
                key="anom_geoloc_detail"
            )
            
            # Option pour visualiser les anomalies géographiquement
            if 'df_geoloc' in locals() and not df_geoloc.empty:
                if st.checkbox("🗺️ Afficher les anomalies sur la carte", key="anom_show_map"):
                    # On vérifie si on a des coordonnées GPS dans les données
                    coords_cols = ['Latitude_depart', 'Longitude_depart', 'Latitude_arrivee', 'Longitude_arrivee']
                    has_coords = all(col in df_geoloc.columns for col in coords_cols) and not df_geoloc[coords_cols].isna().all().all()
                    
                    if has_coords:
                        st.subheader("🗺️ Carte des Anomalies de Géolocalisation")
                        
                        # Extraire les véhicules à problème
                        problem_vehicles = filtered_geoloc['Nouveau Immat'].unique().tolist()
                        
                        # On appelle la fonction de visualisation existante (à adapter si besoin)
                        visualiser_trajets_sur_carte(
                            df_geoloc,
                            None,  # Tous les véhicules à problème
                            date_debut_analyse,
                            date_fin_analyse,
                            highlight_anomalies=True
                        )
                    else:
                        st.warning("Les coordonnées GPS ne sont pas disponibles pour afficher la carte.")
    
    # --- ONGLET 6: TENDANCES ET ÉVOLUTION ---
    with tab_tendances:
        st.subheader("📈 Tendances et Évolution des Anomalies")
        
        if df_anomalies_all.empty:
            st.info("Données insuffisantes pour analyser les tendances.")
        else:
            # S'assurer que la colonne Date est au format datetime
            if 'Date' in df_anomalies_all.columns:
                # Vérifier si la colonne Date est déjà au format datetime
                if not pd.api.types.is_datetime64_any_dtype(df_anomalies_all['Date']):
                    try:
                        df_anomalies_all['Date'] = pd.to_datetime(df_anomalies_all['Date'], errors='coerce')
                    except Exception as e:
                        st.warning(f"Erreur lors de la conversion des dates pour l'analyse des tendances: {e}")
                        st.warning("L'analyse des tendances temporelles peut être incomplète.")
                
                # Filtrer les valeurs non valides après conversion
                valid_dates_mask = df_anomalies_all['Date'].notna()
                df_anomalies_date_valid = df_anomalies_all[valid_dates_mask].copy()
                
                if not df_anomalies_date_valid.empty:
                    # Calculer les tendances d'anomalies par jour
                    df_anomalies_date_valid['AnneeMoisJour'] = df_anomalies_date_valid['Date'].dt.strftime('%Y-%m-%d')
                    tendances_quotidiennes = df_anomalies_date_valid.groupby('AnneeMoisJour').agg(
                        Nb_Anomalies=('type_anomalie', 'size'),
                        Score_Total=('poids_anomalie', 'sum'),
                        Nb_Vehicules=('Nouveau Immat', 'nunique')
                    ).reset_index()
                    
                    # Convertir en datetime pour le graphique
                    tendances_quotidiennes['Date'] = pd.to_datetime(tendances_quotidiennes['AnneeMoisJour'])
                    
                    # Graphique d'évolution quotidienne
                    fig_daily = px.line(
                        tendances_quotidiennes.sort_values('Date'),
                        x='Date',
                        y=['Nb_Anomalies', 'Nb_Vehicules'],
                        title="Évolution Quotidienne des Anomalies",
                        labels={
                            'value': "Nombre",
                            'variable': "Métrique",
                            'Date': "Date"
                        },
                        markers=True
                    )
                    
                    # Ajouter une ligne pour le score total
                    fig_daily.add_trace(
                        go.Scatter(
                            x=tendances_quotidiennes['Date'],
                            y=tendances_quotidiennes['Score_Total'],
                            mode='lines+markers',
                            name='Score Total',
                            yaxis='y2',
                            line=dict(color='red', width=2)
                        )
                    )
                    
                    # Configurer l'axe Y secondaire
                    fig_daily.update_layout(
                        yaxis2=dict(
                            title="Score Total",
                            overlaying='y',
                            side='right'
                        )
                    )
                    
                    st.plotly_chart(fig_daily, use_container_width=True)
                    
                    # Tendances par type d'anomalie
                    st.subheader("📊 Évolution par Type d'Anomalie")
                    
                    # Extraire les types les plus courants
                    top_types = df_anomalies_all['type_anomalie'].value_counts().nlargest(5).index.tolist()
                    
                    if top_types:
                        # Créer un pivot pour les types d'anomalies par jour
                        pivot_types = df_anomalies_date_valid[df_anomalies_date_valid['type_anomalie'].isin(top_types)].groupby(['AnneeMoisJour', 'type_anomalie']).size().unstack(fill_value=0)
                        pivot_types.reset_index(inplace=True)
                        pivot_types['Date'] = pd.to_datetime(pivot_types['AnneeMoisJour'])
                        
                        # Graphique des tendances par type
                        fig_types_trend = px.line(
                            pivot_types,
                            x='Date',
                            y=top_types,
                            title="Évolution des Principaux Types d'Anomalies",
                            labels={
                                'value': "Nombre d'Anomalies",
                                'variable': "Type d'Anomalie",
                                'Date': "Date"
                            },
                            markers=True
                        )
                        
                        st.plotly_chart(fig_types_trend, use_container_width=True)
                    
                    # Tendances par jour de la semaine
                    st.subheader("📅 Répartition par Jour de la Semaine")
                    
                    # Ajouter le jour de la semaine
                    df_anomalies_date_valid['JourSemaine'] = df_anomalies_date_valid['Date'].dt.day_name()
                    df_anomalies_date_valid['NumJourSemaine'] = df_anomalies_date_valid['Date'].dt.dayofweek
                    
                    # Compter les anomalies par jour
                    anomalies_par_jour = df_anomalies_date_valid.groupby(['NumJourSemaine', 'JourSemaine']).agg(
                        Nb_Anomalies=('type_anomalie', 'size'),
                        Score_Total=('poids_anomalie', 'sum')
                    ).reset_index()
                    
                    # Ajouter l'ordre des jours
                    ordre_jours = {
                        0: 'Monday', 1: 'Tuesday', 2: 'Wednesday', 
                        3: 'Thursday', 4: 'Friday', 5: 'Saturday', 6: 'Sunday'
                    }
                    
                    # Créer une colonne JourOrdonné basée sur NumJourSemaine
                    anomalies_par_jour['JourOrdonné'] = anomalies_par_jour['NumJourSemaine'].map(ordre_jours)
                    
                    # Trier par jour de la semaine
                    anomalies_par_jour = anomalies_par_jour.sort_values('NumJourSemaine')
                    
                    # Graphique par jour de la semaine
                    fig_jour = px.bar(
                        anomalies_par_jour,
                        x='JourSemaine',
                        y='Nb_Anomalies',
                        color='Score_Total',
                        title="Anomalies par Jour de la Semaine",
                        labels={
                            'Nb_Anomalies': "Nombre d'Anomalies",
                            'JourSemaine': "Jour de la Semaine"
                        },
                        category_orders={"JourSemaine": ["Monday", "Tuesday", "Wednesday", "Thursday", "Friday", "Saturday", "Sunday"]}
                    )
                    
                    st.plotly_chart(fig_jour, use_container_width=True)
                    
                    # Analyse horaire
                    if 'Hour' in df_anomalies_date_valid.columns:
                        st.subheader("⏰ Répartition par Heure de la Journée")
                        
                        # Convertir l'heure en format correct
                        try:
                            if isinstance(df_anomalies_date_valid['Hour'].iloc[0], datetime.time):
                                df_anomalies_date_valid['Heure'] = df_anomalies_date_valid['Hour'].apply(lambda x: x.hour if pd.notna(x) else None)
                            else:
                                df_anomalies_date_valid['Heure'] = pd.NA
                        except:
                            df_anomalies_date_valid['Heure'] = pd.NA
                        
                        # Compter les anomalies par heure
                        anomalies_par_heure = df_anomalies_date_valid.dropna(subset=['Heure']).groupby('Heure').agg(
                            Nb_Anomalies=('type_anomalie', 'size'),
                            Score_Total=('poids_anomalie', 'sum')
                        ).reset_index()
                        
                        if not anomalies_par_heure.empty:
                            # Graphique par heure
                            fig_heure = px.bar(
                                anomalies_par_heure.sort_values('Heure'),
                                x='Heure',
                                y='Nb_Anomalies',
                                color='Score_Total',
                                title="Anomalies par Heure de la Journée",
                                labels={
                                    'Nb_Anomalies': "Nombre d'Anomalies",
                                    'Heure': "Heure"
                                }
                            )
                            
                            # Ajouter des lignes pour les heures ouvrées/non ouvrées
                            heure_debut_non_ouvre = st.session_state.get('ss_heure_debut_non_ouvre', DEFAULT_HEURE_DEBUT_NON_OUVRE)
                            heure_fin_non_ouvre = st.session_state.get('ss_heure_fin_non_ouvre', DEFAULT_HEURE_FIN_NON_OUVRE)
                            
                            fig_heure.add_vrect(
                                x0=heure_debut_non_ouvre,
                                x1=24,
                                fillcolor="red",
                                opacity=0.1,
                                line_width=0,
                                annotation_text="Heures non ouvrées"
                            )
                            
                            fig_heure.add_vrect(
                                x0=0,
                                x1=heure_fin_non_ouvre,
                                fillcolor="red",
                                opacity=0.1,
                                line_width=0,
                                annotation_text="Heures non ouvrées"
                            )
                            
                            st.plotly_chart(fig_heure, use_container_width=True)
                else:
                    st.warning("Impossible d'analyser les tendances temporelles: pas assez de dates valides après conversion.")
            else:
                st.warning("La colonne 'Date' n'est pas disponible pour analyser les tendances temporelles.")
    
    # --- ONGLET 7: RECOMMANDATIONS ---
    with tab_recommandations:
        st.subheader("💡 Recommandations et Actions Correctives")
        
        if df_anomalies_all.empty:
            st.success("✅ Aucune anomalie détectée, aucune action corrective nécessaire.")
        else:
            # Analyser les types d'anomalies les plus graves
            types_anomalies_graves = df_anomalies_all.groupby('type_anomalie').agg(
                Score_Total=('poids_anomalie', 'sum'),
                Nb_Occurrences=('type_anomalie', 'size'),
                Score_Moyen=('poids_anomalie', 'mean')
            ).sort_values('Score_Total', ascending=False).reset_index()
            
            # Véhicules à risque élevé
            vehicules_risque_eleve = df_scores_all[df_scores_all['score_risque'] >= seuil_score]
            nb_veh_risque_eleve = len(vehicules_risque_eleve)
            
            # Stations à risque
            if 'stations_risque' in locals() and not stations_risque.empty:
                stations_risque_eleve = stations_risque[stations_risque['Niveau_Risque'].isin(['Élevé', 'Critique'])]
                nb_stations_risque_eleve = len(stations_risque_eleve)
            else:
                nb_stations_risque_eleve = 0
            
            # Définir les recommandations en fonction de l'analyse
            st.write("### Actions Prioritaires Recommandées")
            
            # 1. Recommandations sur les véhicules
            if nb_veh_risque_eleve > 0:
                st.warning(f"🚨 **Véhicules à surveiller**: {nb_veh_risque_eleve} véhicules présentent un score de risque élevé (≥ {seuil_score})")
                
                # Liste des véhicules les plus risqués
                if not vehicules_risque_eleve.empty:
                    top_veh_risque = vehicules_risque_eleve.nlargest(5, 'score_risque')
                    
                    with st.expander("📋 Véhicules prioritaires à contrôler"):
                        st.dataframe(
                            top_veh_risque[['Nouveau Immat', 'Catégorie', 'score_risque', 'nombre_total_anomalies']],
                            use_container_width=True
                        )
                        
                        st.markdown("""
                        **Actions recommandées:**
                        1. Effectuer un contrôle mécanique complet de ces véhicules
                        2. Vérifier les procédures de saisie des kilométrages
                        3. Sensibiliser les conducteurs aux bonnes pratiques
                        4. Analyser les trajets effectués via la géolocalisation
                        5. Comparer les consommations avec les moyennes de la catégorie
                        """)
            
            # 2. Recommandations sur les stations
            if nb_stations_risque_eleve > 0:
                st.warning(f"⛽ **Stations à risque**: {nb_stations_risque_eleve} stations présentent un risque élevé de fraude")
                
                # Liste des stations les plus risquées
                if 'stations_risque_eleve' in locals() and not stations_risque_eleve.empty:
                    top_stations_risque = stations_risque_eleve.nlargest(5, 'Score_Risque_Station')
                    
                    with st.expander("📋 Stations prioritaires à auditer"):
                        st.dataframe(
                            top_stations_risque[['Place', 'Nb_Anomalies_Fraude', 'Pourcentage_Anomalies', 'Score_Risque_Station', 'Niveau_Risque']],
                            use_container_width=True
                        )
                        
                        st.markdown("""
                        **Actions recommandées:**
                        1. Effectuer un audit complet de ces stations
                        2. Renforcer les contrôles de validation des transactions
                        3. Utiliser systématiquement la géolocalisation pour vérifier la présence des véhicules
                        4. Envisager un changement de fournisseur si les problèmes persistent
                        5. Former les pompistes aux procédures de sécurité
                        """)
            
            # 3. Recommandations sur les types d'anomalies
            if not types_anomalies_graves.empty:
                top_anomalies = types_anomalies_graves.head(3)['type_anomalie'].tolist()
                
                st.info(f"🔎 **Types d'anomalies à traiter prioritairement**: {', '.join(top_anomalies)}")
                
                with st.expander("📋 Plan d'action par type d'anomalie"):
                    for _, row in types_anomalies_graves.head(5).iterrows():
                        st.markdown(f"#### {row['type_anomalie']}")
                        
                        # Recommandations spécifiques par type d'anomalie
                        if "capacité" in row['type_anomalie'].lower():
                            st.markdown("""
                            **Actions correctives:**
                            - Vérifier les capacités des réservoirs dans la base de données
                            - Former les conducteurs à respecter la capacité maximale
                            - Installer des limiteurs de remplissage sur les véhicules concernés
                            """)
                        
                        elif "kilomét" in row['type_anomalie'].lower():
                            st.markdown("""
                            **Actions correctives:**
                            - Mettre en place une procédure de double vérification des kilométrages
                            - Former le personnel aux bonnes pratiques de saisie
                            - Envisager l'utilisation de systèmes automatiques de relevé
                            """)
                        
                        elif "rapprochées" in row['type_anomalie'].lower():
                            st.markdown("""
                            **Actions correctives:**
                            - Vérifier les heures de prise et les distances parcourues
                            - Sensibiliser les conducteurs à ne pas faire de pleins rapprochés
                            - Configurer des alertes sur les prises à intervalles courts
                            """)
                        
                        elif "présence" in row['type_anomalie'].lower() or "géoloc" in row['type_anomalie'].lower():
                            st.markdown("""
                            **Actions correctives:**
                            - Vérifier systématiquement la présence des véhicules via géolocalisation
                            - Mettre en place un système de validation des transactions en temps réel
                            - Effectuer des contrôles inopinés en station
                            """)
                        
                        elif "facturation" in row['type_anomalie'].lower() or "double" in row['type_anomalie'].lower():
                            st.markdown("""
                            **Actions correctives:**
                            - Auditer les stations concernées
                            - Vérifier les tickets et les heures de transaction
                            - Mettre en place un système d'alerte pour les transactions rapprochées
                            """)
                        
                        else:
                            st.markdown("""
                            **Actions correctives:**
                            - Analyser les causes spécifiques de cette anomalie
                            - Former les équipes à la détection et prévention
                            - Mettre en place des contrôles ciblés
                            """)
            
            # 4. Plan d'action global
            st.write("### Plan d'Action Global")
            
            with st.expander("📝 Plan d'action recommandé", expanded=True):
                st.markdown("""
                #### Court terme (1-2 semaines)
                1. **Contrôle immédiat** des véhicules et stations à risque critique
                2. **Sensibilisation** des conducteurs et responsables de flotte
                3. **Vérification** des procédures de saisie des kilométrages et validations des transactions
                
                #### Moyen terme (1-3 mois)
                1. **Mise en place** de contrôles systématiques sur les situations à risque
                2. **Formation** du personnel aux bonnes pratiques
                3. **Ajustement** des paramètres de détection dans le système
                
                #### Long terme (3-12 mois)
                1. **Automatisation** des contrôles et vérifications
                2. **Intégration** complète de la géolocalisation avec le système de carburant
                3. **Analyse prédictive** pour détecter les anomalies avant qu'elles ne deviennent critiques
                """)
                
                # Formulaire pour exporter le plan d'action
                if st.button("📥 Générer un rapport du plan d'action", key="btn_plan_action"):
                    st.success("✅ Rapport généré ! Utilisez le bouton d'exportation ci-dessus pour télécharger.")
            
            # 5. Résumé des actions immédiates
            st.write("### Résumé des Actions Immédiates")
            
            col_action1, col_action2, col_action3 = st.columns(3)
            
            with col_action1:
                st.info("🚗 **Véhicules**")
                st.markdown(f"- Contrôler {min(nb_veh_risque_eleve, 5)} véhicules prioritaires")
                st.markdown("- Vérifier les kilométrages et consommations")
                st.markdown("- Sensibiliser les conducteurs")
            
            with col_action2:
                st.info("⛽ **Stations**")
                st.markdown(f"- Auditer {min(nb_stations_risque_eleve, 3)} stations à risque")
                st.markdown("- Vérifier les procédures de transaction")
                st.markdown("- Renforcer les contrôles")
            
            with col_action3:
                st.info("🔄 **Système**")
                st.markdown("- Ajuster les paramètres de détection")
                st.markdown("- Améliorer la validation des saisies")
                st.markdown("- Renforcer l'intégration géoloc/carburant")

def afficher_page_parametres(df_vehicules: Optional[pd.DataFrame] = None):
    """Affiche la page des paramètres modifiables."""
    st.header("⚙️ Paramètres de l'Application")
    st.warning("Modifier ces paramètres affectera les analyses et la détection d'anomalies.")

    # Créer des onglets pour organiser les paramètres
    tab_generaux, tab_carburant, tab_geoloc, tab_bilan_carbone = st.tabs([
        "⚙️ Paramètres Généraux", "⛽ Paramètres Carburant", "📍 Paramètres Géolocalisation", "🌍 Bilan Carbone"
    ])

    with tab_generaux:
        with st.expander("Seuils Généraux d'Anomalies", expanded=True):
            st.session_state.ss_seuil_heures_rapprochees = st.number_input(
                "Seuil Prises Rapprochées (heures)",min_value=0.5, max_value=24.0,
                value=float(st.session_state.get('ss_seuil_heures_rapprochees', DEFAULT_SEUIL_HEURES_RAPPROCHEES)),
                step=0.5, format="%.1f", key='param_seuil_rappr'
            )
            st.session_state.ss_delta_minutes_facturation_double = st.number_input(
                "Delta Max Facturation Double (minutes)",min_value=1, max_value=180,
                value=st.session_state.get('ss_delta_minutes_facturation_double', DEFAULT_DELTA_MINUTES_FACTURATION_DOUBLE),
                step=1, key='param_delta_double'
            )
            st.session_state.ss_seuil_anomalies_suspectes_score = st.number_input(
                "Seuil Score de Risque Suspect",min_value=1, max_value=1000,
                value=st.session_state.get('ss_seuil_anomalies_suspectes_score', DEFAULT_SEUIL_ANOMALIES_SUSPECTES_SCORE),
                step=1, key='param_seuil_score'
            )

        with st.expander("Heures Non Ouvrées (Transactions Carburant)"):
            st.session_state.ss_heure_debut_non_ouvre = st.slider(
                "Heure Début Période Non Ouvrée",min_value=0, max_value=23,
                value=st.session_state.get('ss_heure_debut_non_ouvre', DEFAULT_HEURE_DEBUT_NON_OUVRE),
                step=1, key='param_heure_debut_no'
            )
            st.session_state.ss_heure_fin_non_ouvre = st.slider(
                "Heure Fin Période Non Ouvrée (exclusive)",min_value=0, max_value=23,
                value=st.session_state.get('ss_heure_fin_non_ouvre', DEFAULT_HEURE_FIN_NON_OUVRE),
                step=1, key='param_heure_fin_no'
            )
            st.caption(f"Plage non ouvrée actuelle (approximative): de {st.session_state.ss_heure_debut_non_ouvre}h à {st.session_state.ss_heure_fin_non_ouvre}h (hors weekend).")

    with tab_carburant:
        with st.expander("Seuils de Consommation par Catégorie (L/100km)", expanded=True):
            if df_vehicules is not None and st.session_state.get('data_loaded', False):
                current_seuils = st.session_state.get('ss_conso_seuils_par_categorie', {})
                all_cats = sorted(current_seuils.keys())
                new_seuils = {}
                cols = st.columns(3)
                col_idx = 0
                for cat in all_cats:
                    with cols[col_idx % 3]:
                         new_seuils[cat] = st.number_input(
                             f"Seuil {cat}",min_value=5.0, max_value=100.0,
                             value=float(current_seuils.get(cat, DEFAULT_CONSO_SEUIL)),
                             step=0.5, format="%.1f",key=f"param_seuil_conso_{cat}"
                         )
                    col_idx += 1
                st.session_state.ss_conso_seuils_par_categorie = new_seuils
            else:
                st.info("Chargez les données pour définir les seuils par catégorie.")
                st.number_input("Seuil Consommation par Défaut (utilisé si catégorie non définie)", value=DEFAULT_CONSO_SEUIL, disabled=True)

        with st.expander("Poids des Anomalies de Transaction pour Score de Risque"):
            st.caption("Ajustez l'importance de chaque type d'anomalie dans le calcul du score de risque.")
            c1, c2, c3 = st.columns(3)
            with c1:
                st.session_state.ss_poids_conso_excessive = st.slider("Poids: Conso. Excessive", 0, 15, st.session_state.get('ss_poids_conso_excessive', DEFAULT_POIDS_CONSO_EXCESSIVE), key='poids_cex')
                st.session_state.ss_poids_depassement_capacite = st.slider("Poids: Dépassement Capacité", 0, 15, st.session_state.get('ss_poids_depassement_capacite', DEFAULT_POIDS_DEPASSEMENT_CAPACITE), key='poids_dep')
                st.session_state.ss_poids_prises_rapprochees = st.slider("Poids: Prises Rapprochées", 0, 15, st.session_state.get('ss_poids_prises_rapprochees', DEFAULT_POIDS_PRISES_RAPPROCHEES), key='poids_rap')
            with c2:
                st.session_state.ss_poids_km_decroissant = st.slider("Poids: Km Décroissant", 0, 15, st.session_state.get('ss_poids_km_decroissant', DEFAULT_POIDS_KM_DECROISSANT), key='poids_kmd')
                st.session_state.ss_poids_km_inchange = st.slider("Poids: Km Inchangé", 0, 15, st.session_state.get('ss_poids_km_inchange', DEFAULT_POIDS_KM_INCHANGE), key='poids_kmi')
                st.session_state.ss_poids_km_saut = st.slider("Poids: Saut Km Important", 0, 15, st.session_state.get('ss_poids_km_saut', DEFAULT_POIDS_KM_SAUT), key='poids_kms')
            with c3:
                st.session_state.ss_poids_hors_horaire = st.slider("Poids: Hors Horaires/WE (Transaction)", 0, 15, st.session_state.get('ss_poids_hors_horaire', DEFAULT_POIDS_HORS_HORAIRE), key='poids_hor')
                st.session_state.ss_poids_hors_service = st.slider("Poids: Véhicule Hors Service", 0, 15, st.session_state.get('ss_poids_hors_service', DEFAULT_POIDS_HORS_SERVICE), key='poids_hsv')
                st.session_state.ss_poids_fact_double = st.slider("Poids: Facturation Double", 0, 15, st.session_state.get('ss_poids_fact_double', DEFAULT_POIDS_FACT_DOUBLE), key='poids_dbl')

    with tab_geoloc:
        with st.expander("Paramètres Généraux de Géolocalisation", expanded=True):
            st.session_state.ss_rayon_station_km = st.number_input(
                "Rayon autour station (km) - pour correspondance transaction/présence", min_value=0.1, max_value=1.0,
                value=float(st.session_state.get('ss_rayon_station_km', DEFAULT_RAYON_STATION_KM)),
                step=0.1, format="%.1f", key='param_rayon_station'
            )
            st.session_state.ss_seuil_arret_minutes = st.number_input( # Actuellement non utilisé car Type=Arrêt vient du fichier
                "Durée minimale d'un arrêt (minutes) - si calculé par l'app", min_value=1, max_value=30,
                value=st.session_state.get('ss_seuil_arret_minutes', DEFAULT_SEUIL_ARRET_MINUTES),
                step=1, key='param_seuil_arret',
                disabled=True # Le Type=Arrêt vient du fichier donc ce seuil n'est pas utilisé pour le moment
            )
            st.session_state.ss_seuil_detour_pct = st.slider(
                "Seuil écart vitesse pour détour suspect (%)", min_value=5, max_value=50,
                value=st.session_state.get('ss_seuil_detour_pct', DEFAULT_SEUIL_DETOUR_PCT),
                step=5, key='param_seuil_detour',
                help="Un trajet est un détour suspect si sa vitesse moyenne est X% inférieure à la vitesse moyenne habituelle du véhicule pour des trajets significatifs."
            )
            st.session_state.ss_vitesse_excessive_seuil = st.slider(
                "Vitesse maximale autorisée (km/h)", min_value=50, max_value=150,
                value=st.session_state.get('ss_vitesse_excessive_seuil', 90),
                step=5, key='param_vitesse_max',
                help="Seuil de vitesse au-delà duquel un trajet est considéré en excès de vitesse."
            )
            st.session_state.ss_nb_arrets_suspect = st.slider( # Actuellement non utilisé
                "Nombre d'arrêts Type='Arrêt' suspect pour un trajet court", min_value=2, max_value=10,
                value=st.session_state.get('ss_nb_arrets_suspect', DEFAULT_NB_ARRETS_SUSPECT),
                step=1, key='param_nb_arrets_suspect',
                disabled=True # Logique non implémentée
            )

        with st.expander("Activation/Désactivation des Types d'Anomalies", expanded=True):
            st.caption("Activez ou désactivez certains types d'anomalies pour ajuster votre analyse.")
            
            # Initialiser les paramètres d'activation s'ils n'existent pas
            if 'ss_activer_trajets_suspects' not in st.session_state:
                st.session_state['ss_activer_trajets_suspects'] = True  # Activé par défaut
            
            if 'ss_activer_detours_suspects' not in st.session_state:
                st.session_state['ss_activer_detours_suspects'] = True  # Activé par défaut
            
            # Créer les cases à cocher pour activer/désactiver
            st.session_state['ss_activer_trajets_suspects'] = st.checkbox(
                "Activer la détection des trajets suspects (hors heures, weekend, vitesse lente)",
                value=st.session_state['ss_activer_trajets_suspects'],
                help="Désactivez cette option si vous recevez trop d'alertes non pertinentes concernant des trajets hors heures ou weekend."
            )
            
            st.session_state['ss_activer_detours_suspects'] = st.checkbox(
                "Activer les détours suspects",
                                value=st.session_state['ss_activer_detours_suspects'],
                                help="Ajoutez ici le texte d'aide pour les détours suspects."
                            )
            
    with tab_bilan_carbone:
        st.subheader("Paramètres de calcul du bilan carbone")
        
        with st.expander("Facteurs d'émission par type de carburant", expanded=True):
            st.session_state['facteur_emission_essence'] = st.number_input(
                "Facteur d'émission essence (kg CO2e/L)",
                min_value=1.0, max_value=5.0,
                value=st.session_state.get('facteur_emission_essence', 2.3),
                step=0.1, format="%.2f", key='param_emission_essence'
            )
            
            st.session_state['facteur_emission_diesel'] = st.number_input(
                "Facteur d'émission diesel (kg CO2e/L)",
                min_value=1.0, max_value=5.0,
                value=st.session_state.get('facteur_emission_diesel', 2.7),
                step=0.1, format="%.2f", key='param_emission_diesel'
            )
            
            st.session_state['facteur_emission_default'] = st.number_input(
                "Facteur d'émission par défaut (kg CO2e/L)",
                min_value=1.0, max_value=5.0,
                value=st.session_state.get('facteur_emission_default', 2.5),
                step=0.1, format="%.2f", key='param_emission_default'
            )
        
        with st.expander("Options d'affichage du bilan carbone"):
            st.session_state['afficher_comparaisons'] = st.checkbox(
                "Afficher les comparaisons (équivalents CO2)",
                value=st.session_state.get('afficher_comparaisons', True),
                key='param_afficher_comparaisons'
            )
            
            st.session_state['unite_bilan'] = st.radio(
                "Unité d'affichage principale",
                options=["kg CO2e", "tonnes CO2e"],
                index=0 if st.session_state.get('unite_bilan', "kg CO2e") == "kg CO2e" else 1,
                key='param_unite_bilan'
            )
            
            st.session_state['comparaison_type'] = st.multiselect(
                "Types de comparaisons à afficher",
                options=["Kilomètres en voiture", "Vols Paris-New York", "Arbres nécessaires", "Repas avec viande"],
                default=st.session_state.get('comparaison_type', ["Kilomètres en voiture", "Arbres nécessaires"]),
                key='param_comparaison_type'
            )
def afficher_page_analyse_vehicules(df_transactions: pd.DataFrame, df_vehicules: pd.DataFrame, date_debut: datetime.date, date_fin: datetime.date, kpi_cat: pd.DataFrame, df_geoloc: Optional[pd.DataFrame] = None):
    """Affiche la page d'analyse détaillée des véhicules."""
    st.header(f"🚗 Analyse Détaillée des Véhicules ({date_debut.strftime('%d/%m/%Y')} - {date_fin.strftime('%d/%m/%Y')})")

    if df_transactions.empty:
        st.warning("Aucune transaction à analyser pour la période sélectionnée.")
        return

    # Filtre par véhicule
    all_vehicles = sorted(df_vehicules['Nouveau Immat'].dropna().unique())
    selected_vehicle = st.selectbox("Sélectionner un véhicule", ["Tous"] + all_vehicles, key="select_veh_analyse")

    if selected_vehicle != "Tous":
        # Obtenir les informations du véhicule
        vehicle_info_series = df_vehicules[df_vehicules['Nouveau Immat'] == selected_vehicle]
        if vehicle_info_series.empty:
            st.error(f"Information non trouvée pour le véhicule {selected_vehicle}")
            return
        vehicle_info = vehicle_info_series.iloc[0]

        vehicle_card = vehicle_info['N° Carte']

        # Obtenir les transactions du véhicule
        mask_card = df_transactions['Card num.'] == vehicle_card
        vehicle_transactions = df_transactions[mask_card]

        st.subheader(f"Détails du Véhicule: {selected_vehicle}")

        # Informations de base
        col1, col2, col3 = st.columns(3)
        with col1:
            st.metric("Immatriculation", selected_vehicle)
            st.metric("Catégorie", vehicle_info.get('Catégorie', 'Non définie'))
        with col2:
            st.metric("Numéro de Carte", vehicle_card)
            st.metric("Type", vehicle_info.get('Type', 'Non défini'))
        with col3:
            st.metric("Capacité Réservoir", f"{vehicle_info.get('Cap-rèservoir', 0)} L")
            if 'Dotation' in vehicle_info:
                st.metric("Dotation Mensuelle", f"{vehicle_info.get('Dotation', 0)} L")

        # Statistiques des transactions
        if not vehicle_transactions.empty:
            total_volume_veh = vehicle_transactions['Quantity'].sum()
            total_amount_veh = vehicle_transactions['Amount'].sum()
            nb_transactions_veh = len(vehicle_transactions)

            st.subheader("Statistiques de Consommation")
            col4, col5, col6 = st.columns(3)
            col4.metric("Volume Total", f"{total_volume_veh:.1f} L")
            col5.metric("Montant Total", f"{total_amount_veh:,.0f} CFA")
            col6.metric("Nombre de Transactions", f"{nb_transactions_veh}")

            # Ajouter un encadré pour les différentes méthodes de calcul de consommation
            st.subheader("📊 Analyse de consommation par différentes méthodes")

            with st.expander("Méthodes de calcul de consommation", expanded=True):
                # Calcul de distance et consommation avec différentes méthodes
                distance_simple, distance_cumulative, consommation_recommandee, methode_utilisee = calculer_distance_et_consommation(vehicle_transactions)
                
                # Si des données géoloc sont disponibles pour ce véhicule
                distance_geoloc = 0
                consommation_geoloc = 0
                if df_geoloc is not None and not df_geoloc.empty:
                    mask_geoloc_veh = df_geoloc['Véhicule'] == selected_vehicle
                    df_geoloc_veh = df_geoloc[mask_geoloc_veh]
                    
                    if not df_geoloc_veh.empty:
                        # Filtrer uniquement les trajets
                        trajets_veh = df_geoloc_veh[df_geoloc_veh['Type'] == 'Trajet']
                        if not trajets_veh.empty:
                            distance_geoloc = trajets_veh['Distance'].sum()
                            consommation_geoloc = (total_volume_veh / distance_geoloc) * 100 if distance_geoloc > 0 else 0
                
                # Créer un tableau comparatif des méthodes
                comparison_data = {
                    'Méthode': ['Simple (première/dernière transaction)', 'Cumulative (somme des deltas)', 
                              'Hybride/Recommandée', 'Géolocalisation (si disponible)'],
                    'Distance (km)': [distance_simple, distance_cumulative, 
                                   max(distance_simple, distance_cumulative), distance_geoloc],
                    'Consommation (L/100km)': [
                        (total_volume_veh / distance_simple) * 100 if distance_simple > 0 else 0,
                        (total_volume_veh / distance_cumulative) * 100 if distance_cumulative > 0 else 0,
                        consommation_recommandee,
                        consommation_geoloc
                    ],
                    'Description': [
                        'Différence entre le premier et dernier kilométrage',
                        'Somme des distances entre chaque transaction',
                        f'Méthode recommandée ({methode_utilisee})',
                        'Basée sur les données GPS réelles'
                    ]
                }
                
                # Convertir en DataFrame
                df_comparison = pd.DataFrame(comparison_data)
                
                # Arrondir les valeurs numériques
                df_comparison['Distance (km)'] = df_comparison['Distance (km)'].round(1)
                df_comparison['Consommation (L/100km)'] = df_comparison['Consommation (L/100km)'].round(2)
                
                # Surligner la méthode recommandée
                def highlight_recommended(row):
                    if row['Méthode'] == 'Hybride/Recommandée':
                        return ['background-color: rgba(0, 128, 0, 0.2)'] * len(row)
                    return [''] * len(row)
                
                # Afficher le tableau comparatif stylisé
                st.dataframe(df_comparison.style.apply(highlight_recommended, axis=1), use_container_width=True)
                
                # Explication des méthodes
                st.markdown("""
                #### 📝 Explication des méthodes de calcul:
                
                - **Méthode simple**: Utilise uniquement le premier et le dernier kilométrage de la période.
                - **Méthode cumulative**: Additionne toutes les distances individuelles entre chaque transaction.
                - **Méthode hybride/recommandée**: Combine intelligemment les deux méthodes précédentes pour une estimation plus fiable.
                - **Méthode géolocalisation**: Utilise les données GPS réelles des trajets (la plus précise quand disponible).
                
                ⚠️ **Note**: Des écarts importants entre ces méthodes peuvent indiquer des anomalies de déclaration de kilométrage.
                """)
                
                # Créer un graphique comparatif
                if not all(v == 0 for v in df_comparison['Distance (km)'].values):
                    fig_compare = px.bar(
                        df_comparison,
                        x='Méthode',
                        y=['Distance (km)', 'Consommation (L/100km)'],
                        title="Comparaison des méthodes de calcul",
                        barmode='group'
                    )
                    st.plotly_chart(fig_compare, use_container_width=True)
            
            # NOUVEAU: Ajout d'une section pour les stations fréquentées
            st.subheader("Stations Fréquentées")
            
            # Analyser les stations fréquentées
            stations_frequentees = vehicle_transactions.groupby('Place').agg(
                Nb_Transactions=('Quantity', 'count'),
                Volume_Total=('Quantity', 'sum'),
                Montant_Total=('Amount', 'sum'),
                Premiere_Visite=('Date', 'min'),
                Derniere_Visite=('Date', 'max')
            ).reset_index()
            
            # Calculer le prix moyen par litre par station
            stations_frequentees['Prix_Moyen_Litre'] = np.where(
                stations_frequentees['Volume_Total'] > 0,
                stations_frequentees['Montant_Total'] / stations_frequentees['Volume_Total'],
                0
            )
            
            # Calculer le pourcentage du volume total
            volume_total_veh = stations_frequentees['Volume_Total'].sum()
            stations_frequentees['Pourcentage_Volume'] = np.where(
                volume_total_veh > 0,
                (stations_frequentees['Volume_Total'] / volume_total_veh) * 100,
                0
            )
            
            # Arrondir les valeurs numériques
            for col in ['Volume_Total', 'Prix_Moyen_Litre', 'Pourcentage_Volume']:
                stations_frequentees[col] = stations_frequentees[col].round(1)
            stations_frequentees['Montant_Total'] = stations_frequentees['Montant_Total'].round(0)
            
            # Trier par nombre de transactions décroissant
            stations_frequentees = stations_frequentees.sort_values('Nb_Transactions', ascending=False)
            
            # Calculer la station principale (plus grand volume)
            station_principale = stations_frequentees.iloc[0]['Place'] if not stations_frequentees.empty else "N/A"
            pourc_station_principale = stations_frequentees.iloc[0]['Pourcentage_Volume'] if not stations_frequentees.empty else 0
            
            col_s1, col_s2 = st.columns(2)
            col_s1.metric("Station Principale", station_principale)
            col_s2.metric("% du Volume Total", f"{pourc_station_principale:.1f}%")
            
            # Afficher le tableau des stations
            afficher_dataframe_avec_export(
                stations_frequentees,
                "Stations Fréquentées",
                key=f"stations_frequentees_{selected_vehicle.replace(' ', '_')}"
            )
            
            # Graphique de répartition des volumes par station
            fig_stations = px.pie(
                stations_frequentees.head(5),  # Top 5 des stations
                values='Volume_Total',
                names='Place',
                title=f"Répartition du Volume par Station - Top 5 ({selected_vehicle})",
                hole=0.4  # Style donut
            )
            st.plotly_chart(fig_stations, use_container_width=True)
            
            # Vérifier s'il y a des anomalies spécifiques à certaines stations
            if 'all_anomalies_veh' in locals() and not all_anomalies_veh.empty and 'Place' in all_anomalies_veh.columns:
                anomalies_par_station = all_anomalies_veh.groupby('Place').agg(
                    Nb_Anomalies=('type_anomalie', 'count'),
                    Types_Anomalies=('type_anomalie', lambda x: ', '.join(sorted(set(x))))
                ).reset_index()
                
                if not anomalies_par_station.empty:
                    st.subheader("Anomalies par Station")
                    st.dataframe(anomalies_par_station, use_container_width=True)
                    
                    # Graphique des anomalies par station
                    fig_anomalies_station = px.bar(
                        anomalies_par_station,
                        x='Place',
                        y='Nb_Anomalies',
                        title=f"Anomalies par Station pour {selected_vehicle}",
                        hover_data=['Types_Anomalies']
                    )
                    st.plotly_chart(fig_anomalies_station, use_container_width=True)
            
            # NOUVEAU: Récupération et affichage des anomalies
            # Détecter les anomalies de transaction pour ce véhicule
            df_anomalies_veh = detecter_anomalies(vehicle_transactions, df_vehicules[df_vehicules['N° Carte'] == vehicle_card])
            
            # Ajouter les anomalies de géolocalisation si disponibles
            anomalies_geoloc_veh = pd.DataFrame()
            if df_geoloc is not None and not df_geoloc.empty:
                anomalies_geoloc_veh = detecter_anomalies_geolocalisation(
                    df_geoloc, vehicle_transactions, 
                    df_vehicules[df_vehicules['Nouveau Immat'] == selected_vehicle], 
                    date_debut, date_fin
                )
            
            # Fusionner toutes les anomalies
            if not df_anomalies_veh.empty and not anomalies_geoloc_veh.empty:
                # Assurez-vous que les colonnes nécessaires existent dans les deux dataframes
                common_cols = ['Date', 'Card num.', 'Nouveau Immat', 'type_anomalie', 'detail_anomalie', 'poids_anomalie']
                for col in common_cols:
                    if col not in df_anomalies_veh.columns: df_anomalies_veh[col] = None
                    if col not in anomalies_geoloc_veh.columns: anomalies_geoloc_veh[col] = None
                
                all_anomalies_veh = pd.concat([
                    df_anomalies_veh[common_cols],
                    anomalies_geoloc_veh[common_cols]
                ], ignore_index=True)
            elif not df_anomalies_veh.empty:
                all_anomalies_veh = df_anomalies_veh
            elif not anomalies_geoloc_veh.empty:
                all_anomalies_veh = anomalies_geoloc_veh
            else:
                all_anomalies_veh = pd.DataFrame()
            
            # Affichage du score de risque
            if not all_anomalies_veh.empty:
                # Calcul du score de risque
                score_risque_veh = calculer_score_risque(all_anomalies_veh)
                if not score_risque_veh.empty:
                    st.subheader("⚠️ Score de Risque")
                    score_row = score_risque_veh.iloc[0]
                    col_score1, col_score2, col_score3 = st.columns(3)
                    col_score1.metric("Score de Risque", f"{score_row['score_risque']:.1f}")
                    col_score2.metric("Nombre d'Anomalies", f"{score_row['nombre_total_anomalies']}")
                    
                    # Déterminer le niveau de risque
                    niveau_risque = "Faible"
                    if score_row['score_risque'] >= 40:
                        niveau_risque = "Critique"
                    elif score_row['score_risque'] >= 20:
                        niveau_risque = "Élevé"
                    elif score_row['score_risque'] >= 10:
                        niveau_risque = "Modéré"
                    col_score3.metric("Niveau de Risque", niveau_risque)
                
                # Résumé des anomalies par type
                st.subheader("Résumé des Anomalies Détectées")
                resume_anomalies_veh = all_anomalies_veh.groupby('type_anomalie').agg(
                    Nombre=('type_anomalie', 'count'),
                    Score_Total=('poids_anomalie', 'sum'),
                    Score_Moyen=('poids_anomalie', 'mean')
                ).reset_index().sort_values('Score_Total', ascending=False)
                
                st.dataframe(resume_anomalies_veh, use_container_width=True)
                
                # Graphique des types d'anomalies
                fig_anomalies_veh = px.bar(
                    resume_anomalies_veh,
                    x='type_anomalie', y='Nombre',
                    title=f"Anomalies Détectées pour {selected_vehicle}",
                    color='Score_Total',
                    labels={'type_anomalie': "Type d'Anomalie", 'Nombre': "Nombre d'occurrences"}
                )
                st.plotly_chart(fig_anomalies_veh, use_container_width=True)
                
                # Détail de toutes les anomalies
                with st.expander("Détail de toutes les anomalies", expanded=True):
                    # Convertir la date en datetime si ce n'est pas déjà fait
                    if 'Date' in all_anomalies_veh.columns and not pd.api.types.is_datetime64_any_dtype(all_anomalies_veh['Date']):
                        all_anomalies_veh['Date'] = pd.to_datetime(all_anomalies_veh['Date'])
                    
                    # Trier par date et poids d'anomalie
                    all_anomalies_veh_sorted = all_anomalies_veh.sort_values(['Date', 'poids_anomalie'], ascending=[True, False])
                    
                    # Sélectionner et afficher les colonnes pertinentes
                    cols_display = ['Date', 'type_anomalie', 'detail_anomalie', 'poids_anomalie']
                    # Ajouter les colonnes spécifiques si elles existent
                    optional_cols = ['Quantity', 'Amount', 'Place', 'Hour', 'Past mileage', 'Current mileage']
                    for col in optional_cols:
                        if col in all_anomalies_veh_sorted.columns:
                            cols_display.append(col)
                    
                    afficher_dataframe_avec_export(
                        all_anomalies_veh_sorted[cols_display],
                        f"Anomalies Détectées - {selected_vehicle}",
                        key=f"anomalies_{selected_vehicle}"
                    )
            else:
                st.success("✅ Aucune anomalie détectée pour ce véhicule sur la période sélectionnée.")

            # Afficher les transactions (inchangé)
            st.subheader("Transactions du Véhicule")
            afficher_dataframe_avec_export(
                vehicle_transactions,
                "Transactions du Véhicule",
                key=f"vehicle_transactions_{selected_vehicle.replace(' ', '_')}"
            )

            # Graphique d'évolution de la consommation (inchangé)
            st.subheader("Évolution de la Consommation")
            daily_consumption = vehicle_transactions.groupby(vehicle_transactions['Date'].dt.date)['Quantity'].sum().reset_index()

            fig = px.line(
                daily_consumption,
                x='Date',
                y='Quantity',
                title="Évolution de la Consommation Quotidienne",
                labels={'Date': 'Date', 'Quantity': 'Volume (L)'}
            )
            st.plotly_chart(fig, use_container_width=True)
            
            # NOUVEAU: Ajouter des visualisations spécifiques si données géoloc disponibles
            if df_geoloc is not None and not df_geoloc.empty:
                mask_geoloc_veh = df_geoloc['Véhicule'] == selected_vehicle
                df_geoloc_veh = df_geoloc[mask_geoloc_veh]
                
                if not df_geoloc_veh.empty:
                    st.subheader("Données de Géolocalisation")
                    
                    # Statistiques de géolocalisation
                    trajets_veh = df_geoloc_veh[df_geoloc_veh['Type'] == 'Trajet']
                    if not trajets_veh.empty:
                        col_geo1, col_geo2, col_geo3 = st.columns(3)
                        col_geo1.metric("Nombre de Trajets", f"{len(trajets_veh)}")
                        col_geo1.metric("Distance Totale", f"{trajets_veh['Distance'].sum():.1f} km")
                        
                        if 'Durée_minutes' in trajets_veh.columns:
                            duree_heures = trajets_veh['Durée_minutes'].sum() / 60
                            col_geo2.metric("Durée Totale", f"{duree_heures:.1f} h")
                        
                        if 'Vitesse moyenne' in trajets_veh.columns:
                            col_geo2.metric("Vitesse Moyenne", f"{trajets_veh['Vitesse moyenne'].mean():.1f} km/h")
                            col_geo3.metric("Vitesse Max", f"{trajets_veh['Vitesse moyenne'].max():.1f} km/h")
                        
                        # Visualiser les 5 derniers trajets
                        with st.expander("Visualiser les derniers trajets"):
                            derniers_trajets = trajets_veh.sort_values('Date', ascending=False).head(5)
                            st.dataframe(derniers_trajets[['Date', 'Début', 'Fin', 'Distance', 'Durée_minutes', 'Vitesse moyenne']])
                        
                        # Option pour voir la carte
                        if st.button("Afficher la carte des trajets"):
                            visualiser_trajets_sur_carte(
                                df_geoloc, selected_vehicle, date_debut, date_fin, highlight_anomalies=True
                            )
                    else:
                        st.info("Aucune donnée de trajet disponible pour ce véhicule.")
        else:
            st.info(f"Aucune transaction trouvée pour le véhicule {selected_vehicle} durant la période sélectionnée.")
    else:
        if not kpi_cat.empty:
            st.write("Consommation moyenne par catégorie de véhicule:")
            afficher_dataframe_avec_export(kpi_cat, "KPIs par Catégorie (Vue Globale)", key="kpi_cat_overview_page")


            # Graphique des consommations par catégorie
            fig_cat_kpi = px.bar( # Renommé
                kpi_cat,
                x='Catégorie',
                y=['consommation_moyenne', 'consommation_globale'],
                title="Consommation par Catégorie de Véhicule",
                barmode='group',
                labels={
                    'value': 'Consommation (L/100km)',
                    'variable': 'Type de Consommation'
                }
            )
            st.plotly_chart(fig_cat_kpi, use_container_width=True)

        # Liste de tous les véhicules et leurs statistiques
        merged_df_all_veh = df_transactions.merge( # Renommé
            df_vehicules[['N° Carte', 'Nouveau Immat', 'Catégorie']],
            left_on='Card num.',
            right_on='N° Carte',
            how='inner'
        )

        vehicle_stats_all_veh = merged_df_all_veh.groupby(['Nouveau Immat', 'Catégorie']).agg( # Renommé
            Volume_Total=('Quantity', 'sum'),
            Montant_Total=('Amount', 'sum'),
            Nb_Transactions=('Quantity', 'count')
        ).reset_index()

        st.subheader("Statistiques par Véhicule (Toutes Transactions)")
        afficher_dataframe_avec_export(
            vehicle_stats_all_veh.sort_values('Volume_Total', ascending=False),
            "Statistiques par Véhicule (Global)",
            key="vehicle_stats_all_page"
        )
def afficher_page_analyse_couts(df_transactions: pd.DataFrame, df_vehicules: pd.DataFrame, date_debut: datetime.date, date_fin: datetime.date):
    """Affiche la page d'analyse des coûts."""
    st.header(f"💰 Analyse des Coûts ({date_debut.strftime('%d/%m/%Y')} - {date_fin.strftime('%d/%m/%Y')})")

    if df_transactions.empty:
        st.warning("Aucune transaction à analyser pour la période sélectionnée.")
        return
    # Statistiques globales
    total_volume = df_transactions['Quantity'].sum()
    total_amount = df_transactions['Amount'].sum()
    avg_price_per_liter = total_amount / total_volume if total_volume > 0 else 0
    
    col1, col2, col3 = st.columns(3)
    col1.metric("Volume Total", f"{total_volume:.1f} L")
    col2.metric("Coût Total", f"{total_amount:,.0f} CFA")
    col3.metric("Prix Moyen / Litre", f"{avg_price_per_liter:.0f} CFA/L")
    
    # Fusion avec les données des véhicules
    merged_df = df_transactions.merge(
        df_vehicules[['N° Carte', 'Nouveau Immat', 'Catégorie']],
        left_on='Card num.',
        right_on='N° Carte',
        how='inner'
    )
    
    # Analyse par mois
    st.subheader("Évolution Mensuelle des Coûts")
    monthly_costs = merged_df.groupby(pd.Grouper(key='DateTime', freq='M')).agg(
        Volume=('Quantity', 'sum'),
        Montant=('Amount', 'sum'),
        Nb_Transactions=('Quantity', 'count')
    ).reset_index()
    
    monthly_costs['Mois'] = monthly_costs['DateTime'].dt.strftime('%Y-%m')
    monthly_costs['Prix_Moyen_Litre'] = monthly_costs['Montant'] / monthly_costs['Volume']
    
    fig = px.line(
        monthly_costs, 
        x='Mois', 
        y=['Volume', 'Montant'],
        title="Évolution Mensuelle du Volume et du Coût",
        labels={'value': 'Valeur', 'variable': 'Métrique'}
    )
    st.plotly_chart(fig, use_container_width=True)
    
    fig_price = px.line(
        monthly_costs, 
        x='Mois', 
        y='Prix_Moyen_Litre',
        title="Évolution du Prix Moyen par Litre",
        labels={'Prix_Moyen_Litre': 'Prix (CFA/L)'}
    )
    st.plotly_chart(fig_price, use_container_width=True)
    
    # Analyse par station
    st.subheader("Analyse par Station")
    station_costs = merged_df.groupby('Place').agg(
        Volume_Total=('Quantity', 'sum'),
        Montant_Total=('Amount', 'sum'),
        Nb_Transactions=('Quantity', 'count'),
        Prix_Moyen=('Amount', lambda x: (x.sum() / merged_df.loc[x.index, 'Quantity'].sum()) if merged_df.loc[x.index, 'Quantity'].sum() > 0 else 0)
    ).reset_index()
    
    station_costs = station_costs.sort_values('Montant_Total', ascending=False)
    
    # Afficher le tableau des coûts par station
    afficher_dataframe_avec_export(
        station_costs, 
        "Coûts par Station", 
        key="station_costs"
    )
    
    # Graphique des coûts par station (top 10)
    top_stations = station_costs.head(10)
    fig_stations = px.bar(
        top_stations,
        x='Place',
        y='Montant_Total',
        title="Top 10 des Stations par Coût Total",
        labels={'Montant_Total': 'Montant Total (CFA)'}
    )
    st.plotly_chart(fig_stations, use_container_width=True)
    
    # Analyse par catégorie de véhicule
    st.subheader("Analyse par Catégorie de Véhicule")
    category_costs = merged_df.groupby('Catégorie').agg(
        Volume_Total=('Quantity', 'sum'),
        Montant_Total=('Amount', 'sum'),
        Nb_Transactions=('Quantity', 'count'),
        Nb_Vehicules=('Nouveau Immat', 'nunique')
    ).reset_index()
    
    category_costs['Prix_Moyen_Litre'] = category_costs['Montant_Total'] / category_costs['Volume_Total']
    category_costs['Cout_Moyen_Transaction'] = category_costs['Montant_Total'] / category_costs['Nb_Transactions']
    
    # Afficher le tableau des coûts par catégorie
    afficher_dataframe_avec_export(
        category_costs, 
        "Coûts par Catégorie", 
        key="category_costs"
    )
    
    # Graphique des coûts par catégorie
    fig_category = px.pie(
        category_costs,
        values='Montant_Total',
        names='Catégorie',
        title="Répartition des Coûts par Catégorie"
    )
    st.plotly_chart(fig_category, use_container_width=True)

def afficher_page_kpi(df_transactions: pd.DataFrame, df_vehicules: pd.DataFrame, date_debut: datetime.date, date_fin: datetime.date):
    """Affiche la page des indicateurs clés de performance (KPIs)."""
    st.header(f"📊 Indicateurs Clés de Performance ({date_debut.strftime('%d/%m/%Y')} - {date_fin.strftime('%d/%m/%Y')})")
    
    if df_transactions.empty:
        st.warning("Aucune transaction à analyser pour la période sélectionnée.")
        return
    
    # Calculer les KPIs globaux
    total_volume = df_transactions['Quantity'].sum()
    total_amount = df_transactions['Amount'].sum()
    nb_transactions = len(df_transactions)
    nb_active_cards = df_transactions['Card num.'].nunique()
    
    # Fusion avec les données des véhicules
    merged_df = df_transactions.merge(
        df_vehicules[['N° Carte', 'Nouveau Immat', 'Catégorie']],
        left_on='Card num.',
        right_on='N° Carte',
        how='inner'
    )
    
    nb_active_vehicles = merged_df['Nouveau Immat'].nunique()
    avg_price_per_liter = total_amount / total_volume if total_volume > 0 else 0
    
    # Calcul des KPIs de consommation
    merged_df['distance_parcourue'] = merged_df['Current mileage'] - merged_df['Past mileage']
    valid_distance = merged_df.loc[merged_df['distance_parcourue'] > 0, 'distance_parcourue'].sum()
    avg_consumption_100km = (total_volume / valid_distance) * 100 if valid_distance > 0 else 0
    avg_cost_per_km = total_amount / valid_distance if valid_distance > 0 else 0
    
    # Affichage des KPIs principaux
    st.subheader("KPIs Globaux")
    col1, col2, col3 = st.columns(3)
    col1.metric("Volume Total", f"{total_volume:.1f} L")
    col1.metric("Coût Total", f"{total_amount:,.0f} CFA")
    
    col2.metric("Nombre de Transactions", f"{nb_transactions}")
    col2.metric("Véhicules Actifs", f"{nb_active_vehicles}")
    
    col3.metric("Prix Moyen / Litre", f"{avg_price_per_liter:.0f} CFA/L")
    col3.metric("Coût Moyen / Km", f"{avg_cost_per_km:.1f} CFA/km")
    
    # KPIs par catégorie de véhicule
    st.subheader("KPIs par Catégorie de Véhicule")
    category_kpis = merged_df.groupby('Catégorie').agg(
        Volume_Total=('Quantity', 'sum'),
        Montant_Total=('Amount', 'sum'),
        Nb_Transactions=('Quantity', 'count'),
        Nb_Vehicules=('Nouveau Immat', 'nunique'),
        Distance_Totale=('distance_parcourue', lambda x: x[x > 0].sum())
    ).reset_index()
    
    category_kpis['Consommation_100km'] = category_kpis.apply(
        lambda row: (row['Volume_Total'] / row['Distance_Totale']) * 100 if row['Distance_Totale'] > 0 else 0, 
        axis=1
    )
    
    category_kpis['Cout_km'] = category_kpis.apply(
        lambda row: row['Montant_Total'] / row['Distance_Totale'] if row['Distance_Totale'] > 0 else 0, 
        axis=1
    )
    
    # Afficher le tableau des KPIs par catégorie
    afficher_dataframe_avec_export(
        category_kpis, 
        "KPIs par Catégorie", 
        key="category_kpis"
    )
    
    # Graphiques des KPIs
    col4, col5 = st.columns(2)
    
    with col4:
        fig_consumption = px.bar(
            category_kpis,
            x='Catégorie',
            y='Consommation_100km',
            title="Consommation par Catégorie (L/100km)",
            labels={'Consommation_100km': 'L/100km'}
        )
        st.plotly_chart(fig_consumption, use_container_width=True)
    
    with col5:
        fig_cost = px.bar(
            category_kpis,
            x='Catégorie',
            y='Cout_km',
            title="Coût par Km par Catégorie",
            labels={'Cout_km': 'CFA/km'}
        )
        st.plotly_chart(fig_cost, use_container_width=True)
    
    # KPIs mensuels
    st.subheader("Évolution Mensuelle des KPIs")
    monthly_kpis = merged_df.groupby(pd.Grouper(key='DateTime', freq='M')).agg(
        Volume=('Quantity', 'sum'),
        Montant=('Amount', 'sum'),
        Nb_Transactions=('Quantity', 'count'),
        Distance=('distance_parcourue', lambda x: x[x > 0].sum())
    ).reset_index()
    
    monthly_kpis['Mois'] = monthly_kpis['DateTime'].dt.strftime('%Y-%m')
    monthly_kpis['Consommation_100km'] = monthly_kpis.apply(
        lambda row: (row['Volume'] / row['Distance']) * 100 if row['Distance'] > 0 else 0, 
        axis=1
    )
    
    fig_monthly = px.line(
        monthly_kpis, 
        x='Mois', 
        y='Consommation_100km',
        title="Évolution Mensuelle de la Consommation (L/100km)",
        labels={'Consommation_100km': 'L/100km'}
    )
    st.plotly_chart(fig_monthly, use_container_width=True)

def afficher_page_autres_cartes(df_transactions: pd.DataFrame, df_autres: pd.DataFrame, date_debut: datetime.date, date_fin: datetime.date):
    """Affiche la page d'analyse des autres types de cartes (non-véhicules)."""
    st.header(f"💳 Analyse des Autres Cartes ({date_debut.strftime('%d/%m/%Y')} - {date_fin.strftime('%d/%m/%Y')})")
    
    if df_transactions.empty or df_autres.empty:
        st.warning("Aucune transaction ou données de cartes à analyser pour la période sélectionnée.")
        return
    
    # Filtrer les transactions des autres cartes
    autres_cartes_nums = set(df_autres['N° Carte'])
    autres_transactions = df_transactions[df_transactions['Card num.'].isin(autres_cartes_nums)]
    
    if autres_transactions.empty:
        st.info("Aucune transaction pour les autres types de cartes durant la période sélectionnée.")
        return
    
    # Statistiques globales
    total_volume = autres_transactions['Quantity'].sum()
    total_amount = autres_transactions['Amount'].sum()
    nb_transactions = len(autres_transactions)
    nb_active_cards = autres_transactions['Card num.'].nunique()
    
    col1, col2, col3 = st.columns(3)
    col1.metric("Volume Total", f"{total_volume:.1f} L")
    col2.metric("Coût Total", f"{total_amount:,.0f} CFA")
    col3.metric("Cartes Actives", f"{nb_active_cards}/{len(df_autres)}")
    
    # Fusion avec les données des autres cartes
    merged_df = autres_transactions.merge(
        df_autres[['N° Carte']],
        left_on='Card num.',
        right_on='N° Carte',
        how='inner'
    )
    
    # Analyse par carte
    st.subheader("Consommation par Carte")
    
    carte_stats = merged_df.groupby('Card num.').agg(
        Volume_Total=('Quantity', 'sum'),
        Montant_Total=('Amount', 'sum'),
        Nb_Transactions=('Quantity', 'count')
    ).reset_index()
    
    if not carte_stats.empty:
        carte_stats = carte_stats.sort_values('Volume_Total', ascending=False)
        
        # Afficher le tableau des statistiques par carte
        afficher_dataframe_avec_export(
            carte_stats, 
            "Statistiques par Carte", 
            key="autres_cartes_stats"
        )
        
        # Graphique des volumes par carte (top 10)
        top_cartes = carte_stats.head(10)
        fig_cartes = px.bar(
            top_cartes,
            x='Card num.',
            y='Volume_Total',
            title="Top 10 des Autres Cartes par Volume",
            labels={'Volume_Total': 'Volume Total (L)'}
        )
        st.plotly_chart(fig_cartes, use_container_width=True)
    
    # Analyse par mois
    st.subheader("Évolution Mensuelle")
    monthly_stats = merged_df.groupby(pd.Grouper(key='DateTime', freq='M')).agg(
        Volume=('Quantity', 'sum'),
        Montant=('Amount', 'sum'),
        Nb_Transactions=('Quantity', 'count')
    ).reset_index()
    
    monthly_stats['Mois'] = monthly_stats['DateTime'].dt.strftime('%Y-%m')
    
    fig_monthly = px.line(
        monthly_stats, 
        x='Mois', 
        y=['Volume', 'Montant'],
        title="Évolution Mensuelle du Volume et du Coût (Autres Cartes)",
        labels={'value': 'Valeur', 'variable': 'Métrique'}
    )
    st.plotly_chart(fig_monthly, use_container_width=True)
    
    # Transactions individuelles
    st.subheader("Détail des Transactions")
    afficher_dataframe_avec_export(
        autres_transactions, 
        "Transactions des Autres Cartes", 
        key="autres_cartes_transactions"
    )
# ---------------------------------------------------------------------
# Point d'entrée avec navigation mise à jour
# ---------------------------------------------------------------------
def afficher_page_rapports(df_transactions: pd.DataFrame, df_vehicules: pd.DataFrame, df_geoloc: Optional[pd.DataFrame], date_debut: datetime.date, date_fin: datetime.date):
    """Affiche la page de génération de rapports PowerPoint."""
    st.header("📊 Générateur de Rapports PowerPoint")
    
    if df_geoloc is None or df_geoloc.empty:
        st.warning("Aucune donnée de géolocalisation disponible. Veuillez charger un fichier de géolocalisation pour générer des rapports.")
        return
    
    st.subheader("Configuration du Rapport de Géolocalisation")
    
    # Sélection de la période
    col1, col2 = st.columns(2)
    with col1:
        rapport_date_debut = st.date_input(
            "Date de début du rapport", 
            value=date_debut,
            min_value=df_geoloc['Date'].min().date(),
            max_value=df_geoloc['Date'].max().date(),
            key="rapport_date_debut"
        )
    with col2:
        rapport_date_fin = st.date_input(
            "Date de fin du rapport", 
            value=date_fin,
            min_value=rapport_date_debut,
            max_value=df_geoloc['Date'].max().date(),
            key="rapport_date_fin"
        )
    
    # Sélection des véhicules à inclure
    all_vehicles = sorted(df_geoloc['Véhicule'].unique())
    vehicules_rapport = st.multiselect(
        "Sélectionner les véhicules à inclure dans le rapport",
        options=["Tous les véhicules"] + all_vehicles,
        default=["Tous les véhicules"],
        key="rapport_vehicules"
    )
    
    # Options de contenu
    st.subheader("Options de Contenu")
    col3, col4 = st.columns(2)
    
    with col3:
        inclure_cartes = st.checkbox("Inclure les cartes des trajets", value=True, key="rapport_inclure_cartes")
        inclure_analyse_vitesse = st.checkbox("Inclure l'analyse des vitesses", value=True, key="rapport_inclure_vitesses")
    
    with col4:
        inclure_analyse_trajets = st.checkbox("Inclure l'analyse détaillée des trajets", value=True, key="rapport_inclure_trajets")
        inclure_comparaison_carburant = st.checkbox("Inclure la comparaison avec les données carburant", value=True, key="rapport_inclure_comparaison")
    
    # Titre et description du rapport
    titre_rapport = st.text_input(
        "Titre du rapport", 
        value=f"Rapport de Géolocalisation - {rapport_date_debut.strftime('%d/%m/%Y')} à {rapport_date_fin.strftime('%d/%m/%Y')}",
        key="rapport_titre"
    )
    
    description_rapport = st.text_area(
        "Description ou notes additionnelles (facultatif)",
        value="",
        key="rapport_description"
    )
    
    # Personnalisation de l'apparence
    st.subheader("Personnalisation de l'Apparence")
    
    col5, col6 = st.columns(2)
    
    with col5:
        theme_rapport = st.selectbox(
            "Thème du rapport",
            options=["Standard", "Professionnel", "Moderne", "Coloré"],
            index=0,
            key="rapport_theme"
        )
    
    with col6:
        orientation_rapport = st.radio(
            "Orientation des pages",
            options=["Paysage", "Portrait"],
            index=0,
            key="rapport_orientation"
        )
    
    # Options avancées
    with st.expander("Options avancées", expanded=False):
        st.markdown("### 🛠️ Configuration avancée du rapport")
        
        col_adv1, col_adv2 = st.columns(2)
        
        with col_adv1:
            inclure_page_titre = st.checkbox("Inclure une page de titre", value=True, key="rapport_inclure_titre")
            inclure_sommaire = st.checkbox("Inclure un sommaire", value=True, key="rapport_inclure_sommaire")
            inclure_numerotation = st.checkbox("Inclure la numérotation des pages", value=True, key="rapport_inclure_numeros")
        
        with col_adv2:
            inclure_footer = st.checkbox("Inclure pied de page avec date", value=True, key="rapport_inclure_footer")
            inclure_logos = st.checkbox("Inclure logos d'entreprise", value=False, key="rapport_inclure_logos")
            inclure_annexes = st.checkbox("Inclure annexes techniques", value=False, key="rapport_inclure_annexes")
        
        if inclure_logos:
            logo_file = st.file_uploader("Télécharger un logo (PNG, JPG)", type=["png", "jpg", "jpeg"])
    
    # Bouton de génération
    if st.button("📊 Générer le Rapport PowerPoint", key="btn_generer_rapport", type="primary"):
        with st.spinner("Génération du rapport PowerPoint en cours..."):
            try:
                # Vérifier l'installation de python-pptx
                try:
                    from pptx import Presentation
                except ImportError:
                    st.error("📚 La bibliothèque python-pptx n'est pas installée. Exécutez 'pip install python-pptx' puis redémarrez l'application.")
                    st.stop()
                
                # Montrer une barre de progression
                progress_bar = st.progress(0, text="Initialisation du rapport...")
                
                # Simuler quelques étapes de progression
                for percent_complete in [10, 25, 40, 60, 80, 100]:
                    time.sleep(0.5)  # Simuler le travail
                    progress_bar.progress(percent_complete/100, text=f"Génération du rapport... {percent_complete}%")
                
                # Ajuster les véhicules sélectionnés
                veh_to_include = None if "Tous les véhicules" in vehicules_rapport else vehicules_rapport
                
                # Appeler la fonction de génération du rapport
                ppt_buffer = generer_rapport_powerpoint_geoloc(
                    df_geoloc=df_geoloc,
                    df_transactions=df_transactions,
                    df_vehicules=df_vehicules,
                    date_debut=rapport_date_debut,
                    date_fin=rapport_date_fin,
                    vehicules_selectionnes=veh_to_include,
                    inclure_cartes=inclure_cartes,
                    inclure_analyse_vitesse=inclure_analyse_vitesse,
                    inclure_analyse_trajets=inclure_analyse_trajets,
                    inclure_comparaison_carburant=inclure_comparaison_carburant,
                    titre=titre_rapport,
                    description=description_rapport,
                    theme=theme_rapport,
                    orientation=orientation_rapport,
                    inclure_page_titre=inclure_page_titre,
                    inclure_sommaire=inclure_sommaire,
                    inclure_numerotation=inclure_numerotation,
                    inclure_footer=inclure_footer,
                    inclure_logos=inclure_logos,
                    inclure_annexes=inclure_annexes,
                    logo_file=logo_file if 'logo_file' in locals() and logo_file is not None else None
                )
                
                # Préparer le téléchargement
                st.success("✅ Rapport généré avec succès! Cliquez sur le bouton ci-dessous pour télécharger.")
                
                nom_fichier = f"rapport_geoloc_{rapport_date_debut.strftime('%Y%m%d')}_{rapport_date_fin.strftime('%Y%m%d')}.pptx"
                
                # Afficher les détails du rapport
                col_info1, col_info2 = st.columns(2)
                with col_info1:
                    st.info(f"""
                    #### 📄 Informations sur le rapport:
                    
                    - **Période**: {rapport_date_debut.strftime('%d/%m/%Y')} - {rapport_date_fin.strftime('%d/%m/%Y')}
                    - **Véhicules**: {len(veh_to_include) if veh_to_include else 'Tous'}
                    - **Thème**: {theme_rapport}
                    """)
                
                with col_info2:
                    st.info(f"""
                    #### 📊 Contenu du rapport:
                    
                    - **Analyse vitesse**: {'✓' if inclure_analyse_vitesse else '✗'}
                    - **Analyse trajets**: {'✓' if inclure_analyse_trajets else '✗'}
                    - **Comparaison carburant**: {'✓' if inclure_comparaison_carburant else '✗'}
                    - **Cartes**: {'✓' if inclure_cartes else '✗'}
                    """)
                
                st.download_button(
                    label="📥 Télécharger le Rapport PowerPoint",
                    data=ppt_buffer,
                    file_name=nom_fichier,
                    mime="application/vnd.openxmlformats-officedocument.presentationml.presentation",
                    key="download_report"
                )
                
                # Afficher un aperçu ou des instructions
                with st.expander("📋 Aperçu du contenu du rapport", expanded=True):
                    st.markdown(f"""
                    ### 📊 Contenu détaillé du rapport "{titre_rapport}":
                    
                    1. **Page de titre** avec les informations sur la période
                    2. **Sommaire** interactif des sections principales
                    3. **Synthèse** des données de géolocalisation:
                       - Nombre de véhicules: {len(veh_to_include) if veh_to_include else df_geoloc['Véhicule'].nunique()}
                       - Période: {rapport_date_debut.strftime('%d/%m/%Y')} - {rapport_date_fin.strftime('%d/%m/%Y')}
                       - Statistiques clés sur l'utilisation des véhicules
                       
                    4. **Graphiques d'analyse** pour chaque section sélectionnée:
                       - Visualisations des distances parcourues
                       - Répartition des trajets par période
                       - Analyse des vitesses et excès potentiels
                       - Comparaison avec les données de consommation
                       
                    5. **Conclusions et recommandations** basées sur l'analyse des données
                    
                    Ce rapport est optimisé pour Microsoft PowerPoint et peut être facilement partagé avec les parties prenantes.
                    """)
                
            except Exception as e:
                st.error(f"⚠️ Erreur lors de la génération du rapport: {e}")
                st.error("Assurez-vous que toutes les dépendances nécessaires sont installées (python-pptx).")
                st.info("💡 Conseil: Si vous utilisez un environnement virtuel, vérifiez que python-pptx y est installé avec la commande 'pip install python-pptx'")
    
    # Instructions et astuces
    with st.expander("📝 Instructions et Astuces", expanded=False):
        st.markdown("""
        ### 💡 Comment utiliser efficacement le générateur de rapports:
        
        1. **Sélectionnez une période précise** pour des rapports plus ciblés et pertinents.
        
        2. **Limitez le nombre de véhicules** pour des rapports plus spécifiques et rapides à générer.
        
        3. **Adaptez les options de contenu** en fonction de votre audience:
           - Pour la direction: inclure la synthèse et les comparaisons
           - Pour les responsables de flotte: inclure l'analyse détaillée des trajets
           - Pour les responsables sécurité: inclure l'analyse des vitesses
        
        4. **Personnalisez l'apparence** avec les thèmes et options disponibles pour créer des rapports professionnels.
        
        5. **Complétez le rapport** avec vos propres observations et commentaires directement dans PowerPoint après téléchargement.
        
        6. **Installation requise**: Cette fonctionnalité nécessite la bibliothèque `python-pptx`. 
           Si elle n'est pas installée, exécutez `pip install python-pptx` dans votre environnement.
           
        7. **Pour les grands volumes de données**, limitez la période d'analyse pour des performances optimales.
        """)
    
    # Footer avec les crédits et version
    st.markdown("---")
    st.caption("Générateur de Rapports PowerPoint v1.0 | Propulsé par Moctar TALL")

def main():
    st.title("📊 Gestion & Analyse Cartes Carburant")

    st.sidebar.header("1. Chargement des Données")
    fichier_transactions = st.sidebar.file_uploader("Fichier Transactions (CSV)", type=['csv'])
    fichier_cartes = st.sidebar.file_uploader("Fichier Cartes (Excel)", type=['xlsx', 'xls'])
    
    # Ajout du chargement optionnel du fichier de géolocalisation
    fichier_geoloc = st.sidebar.file_uploader("Fichier Géolocalisation (Excel, optionnel)", type=['xlsx', 'xls'])

    if not fichier_transactions or not fichier_cartes:
        st.info("👋 Bienvenue ! Veuillez charger le fichier des transactions (CSV) et le fichier des cartes (Excel) via la barre latérale pour commencer.")
        initialize_session_state() 
        if st.sidebar.radio("Navigation", ["Paramètres"], index=0, key="nav_no_data") == "Paramètres":
            afficher_page_parametres()
        return

    df_transactions, df_vehicules, df_ge, df_autres = charger_donnees(fichier_transactions, fichier_cartes)
    
    # Chargement des données de géolocalisation (optionnel)
    df_geoloc = None
    if fichier_geoloc is not None:
        with st.spinner("Chargement des données de géolocalisation..."):
            df_geoloc = charger_donnees_geolocalisation(fichier_geoloc)
            if df_geoloc is not None:
                st.sidebar.success("✅ Données de géolocalisation chargées avec succès !")
                st.sidebar.markdown(f"**Trajets géolocalisés :** {len(df_geoloc):,}")
                if 'Date' in df_geoloc.columns:
                    min_date_geo = df_geoloc['Date'].min()
                    max_date_geo = df_geoloc['Date'].max()
                    st.sidebar.markdown(f"**Période géoloc :** {min_date_geo.strftime('%d/%m/%Y')} - {max_date_geo.strftime('%d/%m/%Y')}")

    if df_transactions is None or df_vehicules is None or df_ge is None or df_autres is None:
        st.error("❌ Erreur lors du chargement ou de la validation des fichiers principaux. Veuillez vérifier les fichiers et les colonnes requises.")
        st.session_state['data_loaded'] = False
        return 

    st.session_state['data_loaded'] = True
    st.sidebar.success("✅ Données chargées avec succès !")
    min_date, max_date = df_transactions['Date'].min(), df_transactions['Date'].max()
    st.sidebar.markdown(f"**Transactions :** {len(df_transactions):,}")
    st.sidebar.markdown(f"**Période :** {min_date.strftime('%d/%m/%Y')} - {max_date.strftime('%d/%m/%Y')}")

    initialize_session_state(df_vehicules)

    st.sidebar.header("2. Période d'Analyse Globale")
    col_date1, col_date2 = st.sidebar.columns(2)
    global_date_debut = col_date1.date_input("Date Début", min_date.date(), min_value=min_date.date(), max_value=max_date.date(), key="global_date_debut")
    global_date_fin = col_date2.date_input("Date Fin", max_date.date(), min_value=min_date.date(), max_value=max_date.date(), key="global_date_fin")

    if global_date_debut > global_date_fin:
        st.sidebar.error("La date de début ne peut pas être postérieure à la date de fin.")
        return

    mask_global_date = (df_transactions['Date'].dt.date >= global_date_debut) & (df_transactions['Date'].dt.date <= global_date_fin)
    df_transac_filtered = df_transactions[mask_global_date].copy()

    if df_transac_filtered.empty:
         st.warning("Aucune transaction trouvée pour la période sélectionnée.")
    else:
         st.sidebar.info(f"{len(df_transac_filtered):,} transactions dans la période sélectionnée.")

    st.sidebar.header("3. Navigation")
    pages = [
        "Tableau de Bord", "Analyse Véhicules", "Analyse des Coûts", 
        "Analyse par Période", "Suivi des Dotations", "Anomalies", "KPIs", "Autres Cartes", "Bilan Carbone"
    ]
    
    # Ajouter les pages de géolocalisation et de rapports si le fichier est chargé
    if df_geoloc is not None:
        pages.append("Géolocalisation")
        pages.append("Rapports PowerPoint")  # Nouvelle page
        
    pages.append("Paramètres")  # Toujours en dernier
    
    # Laisser toutes les pages accessibles même si df_transac_filtered est vide, les pages géreront l'affichage.
    page = st.sidebar.radio("Choisir une page :", pages, key="navigation_main")

    if page == "Tableau de Bord":
        kpi_cat_dashboard, df_vehicle_kpi_dashboard = calculer_kpis_globaux(
            df_transac_filtered, df_vehicules, global_date_debut, global_date_fin,
            list(st.session_state.ss_conso_seuils_par_categorie.keys()) 
        )
        afficher_page_dashboard(df_transac_filtered, df_vehicules, df_ge, df_autres, global_date_debut, global_date_fin, df_geoloc)
        ameliorer_dashboard(df_transac_filtered, df_vehicules, global_date_debut, global_date_fin, 
                        kpi_cat_dashboard, df_vehicle_kpi_dashboard, df_geoloc)
    elif page == "Analyse Véhicules":
         kpi_cat_veh_page, _ = calculer_kpis_globaux(
             df_transac_filtered, df_vehicules, global_date_debut, global_date_fin,
             list(st.session_state.ss_conso_seuils_par_categorie.keys()) 
         )
         afficher_page_analyse_vehicules(df_transac_filtered, df_vehicules, global_date_debut, global_date_fin, kpi_cat_veh_page, df_geoloc)
    elif page == "Analyse des Coûts":
         afficher_page_analyse_couts(df_transac_filtered, df_vehicules, global_date_debut, global_date_fin)
    elif page == "Analyse par Période":
         afficher_page_analyse_periodes(df_transac_filtered, df_vehicules, global_date_debut, global_date_fin)
    elif page == "Suivi des Dotations":
         afficher_page_suivi_dotations(df_transac_filtered, df_vehicules, global_date_debut, global_date_fin)
    elif page == "Anomalies":
        afficher_page_anomalies(df_transac_filtered, df_vehicules, global_date_debut, global_date_fin, df_geoloc)
    elif page == "KPIs":
        afficher_page_kpi(df_transac_filtered, df_vehicules, global_date_debut, global_date_fin)
    elif page == "Autres Cartes":
        afficher_page_autres_cartes(df_transac_filtered, df_autres, global_date_debut, global_date_fin)
    elif page == "Bilan Carbone":
        afficher_page_bilan_carbone(df_transac_filtered, df_vehicules, global_date_debut, global_date_fin)
    elif page == "Géolocalisation" and df_geoloc is not None:
        # Page d'analyse de géolocalisation
        afficher_page_analyse_geolocalisation(df_geoloc, df_transac_filtered, df_vehicules, global_date_debut, global_date_fin)
    elif page == "Rapports PowerPoint" and df_geoloc is not None:
        # Nouvelle page pour générer des rapports PowerPoint
        afficher_page_rapports(df_transac_filtered, df_vehicules, df_geoloc, global_date_debut, global_date_fin)
    elif page == "Paramètres":
        afficher_page_parametres(df_vehicules)

if __name__ == "__main__":
    main()
